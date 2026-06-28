#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build-template.py — PPTX builder (ppt-lab clean-room rebuild).

Generates 1920x1080 (20 x 11.25 inch) slides from a spec.json using python-pptx.
- Loads design-tokens.json + layouts.json (no hardcoded hex / fonts).
- 좌표는 전부 layouts.json (grid.json SSOT 파생). 매직넘버 금지.
- ARCHETYPE_TO_BUILDER dict routes an archetype slug -> builder function.
- 23 아키타입(7 family) + media() 미디어 슬롯 + chart() ~20 차트형.
- --palette N (valid ids: 1,2,3,4,5,6,8 — no 7) 팔레트 스왑 · --style / --look 직교(레이아웃 불변).
- Korean font fallback set on latin + east-asian + complex-script runs.

Usage:
    python3 build-template.py OUT.pptx [--palette N] [--style S] [--look L] [--spec spec.json]
    python3 build-template.py /tmp/demo.pptx              # 내장 데모(전 아키타입)

spec.json format:
    {
      "title": "...",
      "palette": 2,
      "slides": [
        { "archetype": "cover", "data": {...} },
        { "archetype": "bars",  "data": {...} },
        ...
      ]
    }
"""
import json
import math
import os
import sys
import argparse
import hashlib
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

REF_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # tools/ -> references/

# ---------------------------------------------------------------------------
# Token loading + palette
# ---------------------------------------------------------------------------
TOKENS = json.load(open(os.path.join(REF_DIR, "design-tokens.json"), encoding="utf-8"))
LAYOUTS = json.load(open(os.path.join(REF_DIR, "layouts.json"), encoding="utf-8"))

COLORS = dict(TOKENS["colors"])          # mutated by palette swap
PALETTES = TOKENS["palettes"]
TYPO = TOKENS["typography"]
SLIDE = TOKENS["slide"]
SP = TOKENS["spacing"]
RAD = TOKENS["radius"]
FONT_PRIMARY = TOKENS["fonts"]["primary"]
FONTS = TOKENS["fonts"]
STYLES = TOKENS.get("styles", {})
LOOKS = TOKENS.get("looks", {})          # design-pick 보존 룩(자기완결 style+palette)

# spec.json 의 디렉토리 — 미디어 상대경로 해소 기준(media() 헬퍼가 참조).
# main()/build()에서 spec 로드 시 설정한다.
SPEC_DIR = None

# Active style state — mutated by apply_style() like COLORS is by apply_palette().
# A "style" = the absorbed identity of a reference deck: its fonts + component look.
# Default split: Latin/numerals -> english face (Inter), Hangul -> primary (Pretendard).
# A second optional tier — *_display — is used for headline tiers (DISPLAY_TIERS) so a
# look can pair a display face (e.g. high-contrast 송명) for covers/headings with a
# readable body face. Falls back to the base latin/ea when a look omits `display`.
ACTIVE_FONT = {"latin": FONTS.get("english", FONT_PRIMARY), "ea": FONT_PRIMARY,
               "latin_display": FONTS.get("english", FONT_PRIMARY),
               "ea_display": FONT_PRIMARY,
               # Monospace face for the meta-label LAYER (eyebrow/label/caption +
               # list/flow index numbers + chart value/axis labels). Used ONLY when
               # a look declares grammar.mono_meta; otherwise these tiers keep the
               # proportional latin face (no-op for the other 109 looks). Default
               # is the global tokens.fonts.mono so a look need only flip the switch.
               "mono": FONTS.get("mono", FONTS.get("english", FONT_PRIMARY))}
# Tiers that render in the display face. Mirrors the headline pre-wrap set exactly:
# covers/section/slide-titles + big stat numbers. Sub-headings (h3/h4) stay on body.
DISPLAY_TIERS = ("display", "h1", "h2")
# Meta-label tiers — the eyebrow/kicker, small labels, and captions. grammar.mono_meta
# renders THESE (and these only) in the monospace face + optional `//`/`[ ]` glyph
# decoration. Headlines, body, and DATA VALUES are never mono-routed (their meaning
# must stay intact); the mono treatment is a label SKIN, exactly like brutalism's
# upper_tiers, not a content rewrite.
_META_TIERS = ("eyebrow", "label", "caption")
# Body-copy tiers — a look's `grammar.body_bold` forces these to render bold
# (brutalism keeps body at 700+). Headings/labels carry their own weight.
_BODY_TIERS = ("body", "body-lg", "body-md", "caption")
COMPONENTS = {}                          # active component presets (card/callout/...)
GRAMMAR = {}                             # active look's COMPOSITION grammar (deterministic
                                         # rules a look declares beyond skin: uppercase tiers,
                                         # bold body, cover title-card, square badges, primary
                                         # block fills, brutal charts...). Empty for the 110
                                         # skin-only looks → no behavior change. The "grammar
                                         # layer" that generalizes a deck's composition language
                                         # to every archetype (not just the ones it shipped).
ACCENTS = ["blue", "orange"]             # active palette's accent ROLES (primary, secondary, ...); set by apply_palette
DARK = False                             # True when the active look declares a dark canvas; set by apply_look
_LOOK_EXPLICIT_ROLES = frozenset()       # role tokens a look pinned by hand (surface/divider/...). build()'s
                                         # final resolve_dark() must preserve these so a look that locks its
                                         # surface ladder (dark-tech, prismatic) isn't clobbered by the
                                         # canvas-derived defaults on the re-settle pass. Set by apply_look.
_LOOK_COVER_BG = None                    # a look's cover/title-slide background, when it differs from the
                                         # body canvas (dark cover + light body: navy IR decks, full-bleed
                                         # keynote covers). None → cover follows `canvas` (derived). A spec's
                                         # per-slide `bg` still wins over this. Set by apply_look.

PX_TO_PT = SLIDE["px_to_pt"]             # 0.75


def apply_palette(pal_id):
    """Overlay palette overrides onto base COLORS."""
    if not pal_id:
        return None
    for p in PALETTES:
        if p["id"] == int(pal_id):
            COLORS.update(p["overrides"])
            ACCENTS[:] = p.get("accents", ["blue", "orange"])
            return p
    return None


def accent_cycle(n):
    """n colors from the active palette's accent ROLES (primary→secondary→...),
    padding with primary shades when the palette defines fewer accents than n.

    Mono palette → mono; multi palette → multi — but ALWAYS inside the chosen
    deck's palette. No out-of-palette color can leak (the cause of the old
    stats-green bug). 'One deck, one DEFINED palette' (mono or multi)."""
    # Swiss ("무채색 + 스팟 1색") — 거의 잉크, 첫 항목만 스팟. 슬라이드당 1색 원칙.
    # 브루탈의 accent_repeat(원색 반복)과 정반대 축. mono_accent 선언 룩만.
    # 비강조 채움은 캔버스-인지 무채색: 라이트=잉크(navy), 다크=중간 그레이(muted).
    # navy를 다크 캔버스에 쓰면(예: dark-luxury navy=#382E1E) near-black 배경에
    # 묻혀 사라진다 → 다크 mono 룩이 깨지므로 DARK일 때 muted로 플립.
    if GRAMMAR.get("mono_accent") and ACCENTS:
        return [ACCENTS[0]] + (["muted"] if DARK else ["navy"]) * (n - 1)
    # Brutalism/Memphis ("원색 충돌") repeats the DEFINED primaries verbatim — the
    # raw clash IS the identity — so it opts OUT of tone variation via accent_repeat.
    if GRAMMAR.get("accent_repeat") and ACCENTS:
        return [ACCENTS[i % len(ACCENTS)] for i in range(n)]
    # GLOBAL default: the declared accents, then in-palette TONE VARIANTS (lower
    # saturation + lighter) of the SAME hues when more colors are needed than the
    # palette declares. Never an identical solid duplicate, never an off-palette hue
    # (the old blue-2/blue-light/blue-pale pad leaked blue into non-blue palettes).
    # n ≤ len(ACCENTS) → returns the base accents unchanged (byte-identical to before).
    if not ACCENTS:
        return ["blue"] * n
    na = len(ACCENTS)
    return [ACCENTS[i] if i < na else _palette_variant(ACCENTS[i % na], i // na)
            for i in range(n)]


def apply_style(style_id):
    """Resolve a style preset (fonts + components). Deterministic per chosen deck.

    Falls back to the 'house' style when style_id is None or unknown, so every
    build has a coherent default. Latin/EA fonts are resolved separately: Latin
    follows the deck's typeface, Hangul (ea/cs) falls back to a Korean-capable
    font so absorbing a Western reference (e.g. Calibri) never tofu-boxes Korean.
    """
    st = STYLES.get(str(style_id)) if style_id else None
    if st is None:
        st = STYLES.get("house")
    if st is None:
        return None
    _apply_fonts(st.get("fonts", {}))
    COMPONENTS.clear()
    COMPONENTS.update(st.get("components", {}))
    # A plain style carries no composition grammar; an explicit --style override
    # therefore strips a look's grammar (the absorbed identity is being replaced).
    GRAMMAR.clear()
    GRAMMAR.update(st.get("grammar", {}))
    return st


def _apply_fonts(f):
    """Set ACTIVE_FONT from a look/style `fonts` block. The optional nested
    `display: {latin, ea}` overrides the headline-tier faces; absent → same as
    body faces (so existing single-font looks are unchanged)."""
    base_latin = f.get("latin", FONTS.get("english", FONT_PRIMARY))
    base_ea = f.get("ea", FONT_PRIMARY)
    disp = f.get("display", {})
    ACTIVE_FONT["latin"] = base_latin
    ACTIVE_FONT["ea"] = base_ea
    ACTIVE_FONT["latin_display"] = disp.get("latin", base_latin)
    ACTIVE_FONT["ea_display"] = disp.get("ea", base_ea)
    # Monospace meta face: a look may name its own (`fonts.mono`), else the global
    # tokens default stands. Resets every look so a prior dark-tech build can't leak
    # a mono face into the next look (the mono_meta GATE still controls whether it's
    # ever USED — this just keeps the slot correct).
    ACTIVE_FONT["mono"] = f.get("mono", FONTS.get("mono", base_latin))


def apply_look(look_id):
    """Apply a self-contained design-pick look in one shot: fonts + components +
    palette (colors + accent roles). A look bundles a style and a palette so that
    `--look <slug>` reproduces that look's identity, while the LAYOUT axis (variant)
    stays orthogonal. Explicit --style / --palette can still override afterwards.
    Returns the look dict, or None if unknown."""
    lk = LOOKS.get(str(look_id))
    if lk is None:
        return None
    _apply_fonts(lk.get("fonts", {}))
    COMPONENTS.clear()
    COMPONENTS.update(lk.get("components", {}))
    GRAMMAR.clear()
    GRAMMAR.update(lk.get("grammar", {}))   # composition rules (empty for skin-only looks)
    pal = dict(lk.get("palette", {}))
    accents = pal.pop("accents", ["blue", "orange"])
    explicit_roles = set(pal.keys())   # roles the look set by hand (respected)
    COLORS.update(pal)
    ACCENTS[:] = accents
    global _LOOK_EXPLICIT_ROLES, _LOOK_COVER_BG
    _LOOK_EXPLICIT_ROLES = frozenset(explicit_roles)
    _LOOK_COVER_BG = lk.get("cover_bg")   # dark cover on a light body, when the look declares one
    resolve_dark(explicit_roles, force=lk.get("dark"))
    return lk


# ---------------------------------------------------------------------------
# Size axis (4th orthogonal axis: layout × style × palette × SIZE)
# ---------------------------------------------------------------------------
# Type SIZE is its own axis, independent of the look. A look defines the visual
# identity (fonts/colors/components) and its OWN natural type scale; the size
# axis re-targets the on-screen type to a reading distance. Body-led: body/labels
# grow most, headings little (a deck shouldn't get title-heavy just because the
# room is big). Absolute per-tier scale that OVERRIDES the look's own `type.scale`
# (line-height/tracking from the look are preserved). "소" = no override → the
# look's natural size stands. Reused by build_matrix's row-height auto-fit.
SIZE_TIERS = {
    "중": {"body": 1.5,   "body-md": 1.5,   "body-lg": 1.35,
           "h4": 1.2,  "h3": 1.18, "h2": 1.06, "label": 1.5, "eyebrow": 1.2, "caption": 1.3},
    "대": {"body": 1.667, "body-md": 1.667, "body-lg": 1.45,
           "h4": 1.25, "h3": 1.22, "h2": 1.1,  "label": 1.6, "eyebrow": 1.25, "caption": 1.35},
}
_SIZE_ALIASES = {
    "소": "소", "s": "소", "small": "소", "sm": "소",
    "중": "중", "m": "중", "medium": "중", "md": "중",
    "대": "대", "l": "대", "large": "대", "lg": "대",
}


def apply_size(tier):
    """Apply a size tier (소/중/대) on top of the active look. Replaces only the
    per-tier size `scale` inside GRAMMAR['type'] with a FRESH dict (look line/
    tracking preserved; LOOKS source never mutated). '소'/unknown → no-op so the
    look's own type metrics stand. Returns the canonical tier or None."""
    key = _SIZE_ALIASES.get(str(tier).strip().lower())
    prof = SIZE_TIERS.get(key)   # '소' has no profile → look stands
    if not prof:
        return key  # canonical (e.g. '소') or None — caller may log
    t = dict(GRAMMAR.get("type") or {})
    t["scale"] = dict(prof)
    GRAMMAR["type"] = t
    return key


def resolve_dark(explicit_roles=frozenset(), force=None):
    """Decide DARK and fill the role palette to match the active canvas.
    A look is dark when its `canvas` token resolves to a dark color (or force=True).
    For a dark canvas, every role NOT explicitly set by the look is auto-derived
    from the canvas hue, so any look becomes readable dark by declaring one color.
    Also fixes on-accent in BOTH modes from accent luminance (light accents get
    dark text). Idempotent — safe to call again after a palette/style override."""
    global DARK
    canvas_hex = _resolve("canvas")
    accent_hex = _resolve(ACCENTS[0] if ACCENTS else "blue")
    DARK = bool(force) if force is not None else is_dark(canvas_hex)
    if DARK:
        for k, v in derive_dark_roles(canvas_hex, accent_hex).items():
            if k not in explicit_roles:
                COLORS[k] = v
    else:
        # Light mode: roles keep their base aliases; only correct on-accent so a
        # light accent (e.g. mint/lime) doesn't get invisible white text.
        if "on-accent" not in explicit_roles:
            COLORS["on-accent"] = "#0F172A" if _rel_lum(accent_hex) > 0.6 else "white"


# ---------------------------------------------------------------------------
# Unit + color helpers
# ---------------------------------------------------------------------------
def IN(px):
    return Inches(px / 96.0)


def PT_PX(px):
    """px (Figma) -> pt for python-pptx."""
    return Pt(px * PX_TO_PT)


def _resolve(token, _depth=0):
    """Resolve a token through one or more alias hops to a raw #hex.
    Role tokens (ink/muted/surface/canvas…) alias other tokens, so a dark look
    can repoint a role in its palette and every builder that reads the role flips
    at once. Stops at a value that is a #hex literal (or any non-key string)."""
    v = COLORS.get(token, token)
    if isinstance(v, str) and not v.startswith("#") and v in COLORS and _depth < 8:
        return _resolve(v, _depth + 1)
    return v


def C(token):
    """token name, role alias, or raw #hex -> RGBColor."""
    if token is None:
        token = "ink"
    hexv = _resolve(token)
    return RGBColor.from_string(hexv.lstrip("#").upper())


# ---------------------------------------------------------------------------
# Dark-canvas auto-derivation — generalizes dark mode to ANY look.
# A look declares ONE color: `canvas` (its content background). Everything else
# (ink ladder, surfaces, dividers, on-accent text) is derived deterministically
# from that canvas hue so the surface palette always MATCHES the deck's
# background — no per-look hand tuning. Light looks keep the base aliases.
# ---------------------------------------------------------------------------
def _hex2rgb(h):
    h = str(h).lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _rgb2hex(r, g, b):
    return "#%02X%02X%02X" % tuple(
        max(0, min(255, int(round(c)))) for c in (r, g, b))


def _rel_lum(hexv):
    """WCAG relative luminance 0..1 from a #hex."""
    def lin(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = _hex2rgb(hexv)
    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def _mix(a, b, t):
    """Linear blend a->b by t in [0,1], in sRGB space (good enough for UI tints)."""
    ar, ag, ab = _hex2rgb(a)
    br, bg, bb = _hex2rgb(b)
    return _rgb2hex(ar + (br - ar) * t, ag + (bg - ag) * t, ab + (bb - ab) * t)


def is_dark(hexv):
    """A canvas is dark when its luminance sits below mid — text must invert."""
    try:
        return _rel_lum(hexv) < 0.30
    except Exception:
        return False


def _ink_on(fill_token):
    """Readable text color for a flat block of `fill_token`: dark ink on light
    fills (e.g. brutalist yellow), white on dark/saturated fills (red/blue).
    Used where a builder paints a card/cell in a PALETTE color (block_fill
    grammar) so per-block text contrast is correct regardless of the hue."""
    try:
        return "navy" if _rel_lum(_resolve(fill_token)) > 0.6 else "white"
    except Exception:
        return "white"


def _palette_variant(hexv, tier):
    """A deterministic IN-PALETTE variant of a base accent for when a builder needs
    more colors than the palette declares. Instead of repeating an identical solid
    (visual duplicate) or padding with an off-palette blue, we vary the SAME hue by
    cutting saturation and lifting lightness per `tier` (1,2,…) → a recessive sibling
    tone of the same color family. Hue is preserved, so the deck never leaves its
    palette; near-neutral hues (charcoal) lean on the lightness lift since they have
    little saturation to give. Used by accent_cycle's overflow (global rule)."""
    try:
        r, g, b = (c / 255.0 for c in _hex2rgb(_resolve(hexv)))
    except Exception:
        return hexv
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    t = max(1, int(tier))
    s2 = s * (0.6 ** t)                              # desaturate ~40% per tier
    l2 = l + (0.90 - l) * min(0.86, 0.36 * t)        # lift toward a pale tone, ramped
    r2, g2, b2 = colorsys.hls_to_rgb(h, min(0.94, max(0.0, l2)), max(0.0, s2))
    return _rgb2hex(round(r2 * 255), round(g2 * 255), round(b2 * 255))


# ---------------------------------------------------------------------------
# 일반 가독성 레이어 — 배경 atmosphere 결정론 도출 (resolve_dark 의 가독성 짝)
# ---------------------------------------------------------------------------
# resolve_dark() 가 "캔버스 휘도 → ink/surface 역할토큰"을 자동 도출하듯,
# _bg_atmosphere() 는 "배경 속성(그라디언트/글로우/이미지) → 카드 반투명·차트 백드롭
# 정책"을 자동 도출한다. 룩이 배경을 *선언*하기만 하면 가독성이 따라온다 — 흡수 때
# 카드 alpha 를 수동 지정할 필요가 없어진다(룩별 수공예 → 엔진 일반 기능).
#
# 두 비트:
#   depth — 배경이 "공간감/광원"을 가지나? (그라디언트·라디얼글로우·이미지)
#           → 솔리드 카드를 프로스티드/다크글래스 패널로 띄울 근거. grid_bg/glow 만으론
#             depth 아님(dark-tech 의 솔리드 코드블록 카드 유지 위해 — 무회귀 1차 가드).
#   busy  — 그 배경이 "다채(다색·고채도)"라 데이터 마크 가독성을 위협하나?
#           → 메시(vivid) = busy(프로스티드 카드 + fragile 차트 백드롭),
#             단색·동색조 어두운 그라디언트(hyundai 네이비·engineered 바이올렛) = calm.
# 결정론·렌더 불필요. 임계값은 전부 명명 상수(매직넘버 금지).
import colorsys  # noqa: E402  (HSV 변환 — busy 판정용, 표준 라이브러리)

# depth/busy 임계 상수 (이름으로 의미 고정)
_BG_BUSY_HUE_BUCKET_DEG = 30      # 색상환을 30° 버킷으로 양자화해 "서로 다른 hue" 카운트
_BG_BUSY_DISTINCT_HUES = 3        # 서로 다른 hue 버킷 ≥3 → 다채 (vivid 4정점 메시)
_BG_BUSY_HUE_RANGE_DEG = 120      # 또는 hue 범위 ≥120° AND
_BG_BUSY_MEAN_SAT = 0.40          #        평균 채도 ≥0.40 → 다채
# 카드 알파 레짐 2개 (busy=frost 흰틴트 저알파, calm-depth=옅은 다크글래스)
_FROST_ALPHA = 16    # busy: 흰 프로스트 틴트 16% 불투명 — 메시가 비치되 텍스트 대비 확보
_FROST_TINT = "#FFFFFF"   # frost 레짐은 흰 틴트로 패널을 칠해 어떤 메시색 위에서도 균일
_DARKGLASS_ALPHA = 84     # calm-depth: 채움 84% 불투명 — 배경 글로우가 가장자리로 은은히


def _hue_sat_of(hexv):
    """#hex → (hue 0..360, saturation 0..1). busy 판정용 HSV."""
    r, g, b = _hex2rgb(hexv)
    h, s, v = colorsys.rgb_to_hsv(r / 255.0, g / 255.0, b / 255.0)
    return h * 360.0, s


def _bg_atmosphere():
    """활성 룩의 배경 속성에서 (depth, busy) 를 결정론적으로 도출. 룩 무관 — 어떤 룩이든
    배경을 선언한 *방식*만 본다. 반환: dict(depth=bool, busy=bool).

      depth = grammar.bg_gradient(스톱≥2) OR grammar.radial_glow OR 이미지 배경
              중 하나라도 존재 (grid_bg/glow 는 제외 — 솔리드 카드 룩 보존).
      busy  = bg_gradient 스톱들을 HSV 변환해, 서로 다른 hue 버킷 ≥3
              (다색 메시), 또는 hue 범위 ≥120° AND 평균 채도 ≥0.40.
              그라디언트가 없으면(라디얼글로우·이미지만) busy=False (단광원 = calm).

    DARK 무관하게 계산하되, 자동 카드/백드롭 적용부에서 DARK 게이트를 별도로 건다
    (라이트 캔버스는 반투명 패널이 의미 없음)."""
    g = GRAMMAR.get("bg_gradient")
    mesh = GRAMMAR.get("bg_mesh")
    # 메시 룩은 정점색(verts) 또는 파생원 bg_gradient 를 busy 판정 스톱으로 쓴다.
    mesh_verts = (mesh.get("verts") if isinstance(mesh, dict) else None) if mesh else None
    stops = [s for s in (mesh_verts or g or [])]
    has_grad = bool(stops) and len(stops) >= 2
    has_glow = bool(GRAMMAR.get("radial_glow"))
    has_image = bool(GRAMMAR.get("bg_image"))   # 미래 이미지 배경 씨임(현재 룩엔 없음)
    depth = has_grad or has_glow or has_image or bool(mesh)
    busy = False
    if has_grad:
        try:
            hs = [_hue_sat_of(_resolve(c)) for c in stops]
            hues = [h for h, s in hs]
            sats = [s for h, s in hs]
            buckets = {int(h // _BG_BUSY_HUE_BUCKET_DEG) for h in hues}
            hue_range = max(hues) - min(hues)
            mean_sat = sum(sats) / len(sats) if sats else 0.0
            busy = (len(buckets) >= _BG_BUSY_DISTINCT_HUES or
                    (hue_range >= _BG_BUSY_HUE_RANGE_DEG and mean_sat >= _BG_BUSY_MEAN_SAT))
        except Exception:
            busy = False
    return {"depth": depth, "busy": busy}


def derive_dark_roles(canvas_hex, accent_hex):
    """Deterministically build the full role palette for a dark canvas.
    Ink lifts toward white at decreasing strength (primary→tertiary); surfaces
    lift the canvas itself by small amounts (so cards read as elevated panels of
    the SAME hue); dividers/hairlines lift a touch more. on-* text inverts to the
    canvas so chips/accents stay legible. Accent-on-light is handled by luminance:
    when the accent is bright, text on it goes dark."""
    W = "#FFFFFF"
    roles = {
        "ink": W,
        "ink-2": _mix(canvas_hex, W, 0.85),
        "body": _mix(canvas_hex, W, 0.70),
        "muted": _mix(canvas_hex, W, 0.56),
        "subtle": _mix(canvas_hex, W, 0.44),
        "surface": _mix(canvas_hex, W, 0.07),
        "surface-2": _mix(canvas_hex, W, 0.11),
        "surface-3": _mix(canvas_hex, W, 0.17),
        "band": _mix(canvas_hex, W, 0.11),
        "divider": _mix(canvas_hex, W, 0.17),
        "hairline": _mix(canvas_hex, W, 0.30),
        # Emphasis slab (versus card / table header): in light mode a near-black
        # navy block; in dark mode a *subtly elevated dark* panel (NOT inverted to
        # white, which would scream on a dark deck). Text on it stays light.
        "slab": _mix(canvas_hex, W, 0.05),
        "on-slab": W,
        "on-ink": canvas_hex,
    }
    roles["on-accent"] = canvas_hex if _rel_lum(accent_hex) > 0.55 else W
    return roles


def _type_override():
    """The active look's per-tier typography override (size scale / line / tracking),
    layered on top of the GLOBAL TYPO scale. Empty for skin-only looks → global
    metrics unchanged. This is the seam that makes a look's TYPE METRICS absorbable
    (not just its font face): brutalism's compressed/chunky scale + tight leading +
    tight tracking live here, scoped to the look."""
    return GRAMMAR.get("type") or {}


def _size_px(key):
    """Tier size in px after the active look's size-scale override (×multiplier)."""
    base = TYPO[key]["size_px"]
    return base * (_type_override().get("scale") or {}).get(key, 1.0)


def style(key):
    """typography key -> (pt, bold, line, upper). Size and line-height reflect the
    active look's `type` override (global TYPO is the default)."""
    s = TYPO[key]
    bold = bool(s.get("bold") or s.get("medium"))
    line = (_type_override().get("line") or {}).get(key, s.get("line", 1.2))
    return PT_PX(_size_px(key)), bold, line, bool(s.get("upper"))


# ---------------------------------------------------------------------------
# Korean-safe pre-wrap (headline tiers only)
# ---------------------------------------------------------------------------
# PowerPoint wraps CJK text at ANY character boundary, so a Korean 어절 like
# "한다" can split as "한"/"다" at a line end. The OOXML eaLnBrk attribute that
# should control this is ignored by PowerPoint's render engine (confirmed
# empirically) and U+2060 word-joiners are ignored too — so the ONLY reliable
# fix is to measure the rendered width with the real font and insert explicit
# line breaks at whitespace (어절) boundaries BEFORE the text reaches PowerPoint.
# Applied to headline tiers (display/h1/h2) only, where mid-word breaks are most
# visible; body/caption keep PowerPoint's native wrap. Pure measurement — the
# rendered runs are unchanged. Units are Figma px: PX_TO_PT=0.75 makes the
# 96-DPI export pixel size equal the token size_px, and layout boxes are Figma
# px too, so box width and font px share one space (no DPI juggling).
import os as _os
_FONT_DIR = "/mnt/c/dev/ppt-fonts"
_PIL_FONT_CACHE = {}
_FONT_PATH_CACHE = {}


# Weight words that may appear IN a face name (e.g. "Pretendard Black", "Archivo
# Black"). When present, the prewrap measurement picks the matching weight file so
# the measured width matches the rendered heavy face (Black is wider than Bold).
_WEIGHT_WORDS = ("black", "heavy", "extrabold", "semibold", "bold", "medium",
                 "light", "thin", "regular")


def _font_path(face, bold):
    key = (face, bold)
    if key in _FONT_PATH_CACHE:
        return _FONT_PATH_CACHE[key]
    path = None
    face_l = face.lower()
    # A weight baked into the face name (Pretendard "Black") wins over the bold flag.
    named_weight = next((wd for wd in _WEIGHT_WORDS if wd in face_l), None)
    for folder in (face.replace(" ", "_"), face.split(" ")[0]):
        d = _os.path.join(_FONT_DIR, folder)
        if _os.path.isdir(d):
            ttfs = [f for f in _os.listdir(d) if f.lower().endswith((".ttf", ".otf"))]
            pool = [f for f in ttfs if "italic" not in f.lower()] or ttfs
            pick = None
            if named_weight:
                pick = next((f for f in pool if named_weight in f.lower()), None)
            if pick is None:
                want = "bold" if bold else "regular"
                pick = next((f for f in pool if want in f.lower()
                             and "extrabold" not in f.lower()
                             and "semibold" not in f.lower()), None)
                pick = pick or next((f for f in pool if want in f.lower()), None)
            pick = pick or (pool[0] if pool else None)
            if pick:
                path = _os.path.join(d, pick)
                break
    if path is None and face != "Inter":  # proprietary/missing → proportional fallback
        path = _font_path("Inter", bold)
    _FONT_PATH_CACHE[key] = path
    return path


def _pil_font(face, bold, px):
    key = (face, bold, px)
    if key not in _PIL_FONT_CACHE:
        from PIL import ImageFont
        p = _font_path(face, bold)
        _PIL_FONT_CACHE[key] = ImageFont.truetype(p, px) if p else ImageFont.load_default()
    return _PIL_FONT_CACHE[key]


def _is_cjk(ch):
    o = ord(ch)
    return (0xAC00 <= o <= 0xD7A3 or 0x3130 <= o <= 0x318F
            or 0x3000 <= o <= 0x303F or 0x4E00 <= o <= 0x9FFF
            or 0xFF00 <= o <= 0xFFEF)


def _measure(s, latin_face, ea_face, bold, px):
    """Width of s in px, routing CJK glyphs to the ea font, others to latin."""
    total, seg, seg_cjk = 0.0, "", None
    for ch in s:
        c = _is_cjk(ch)
        if seg and c != seg_cjk:
            total += _pil_font(ea_face if seg_cjk else latin_face, bold, px).getlength(seg)
            seg = ""
        seg += ch
        seg_cjk = c
    if seg:
        total += _pil_font(ea_face if seg_cjk else latin_face, bold, px).getlength(seg)
    return total


def _prewrap(content, box_w, size_px, bold, latin_face, ea_face, upper=False):
    """Greedy word-pack at whitespace so no 어절 splits mid-line. Returns a list
    of lines. On any failure (missing font, PIL error) returns [content] so the
    build never breaks; the worst case is PowerPoint's original native wrap."""
    try:
        px = max(1, int(round(size_px)))
        limit = box_w - 2  # tiny safety margin vs render rounding
        out = []
        for src in str(content).split("\n"):
            cur = ""
            for w in src.split(" "):
                trial = w if not cur else cur + " " + w
                m = _measure(trial.upper() if upper else trial,
                             latin_face, ea_face, bold, px)
                if not cur or m <= limit:
                    cur = trial
                else:
                    out.append(cur)
                    cur = w
            out.append(cur)
        return out or [content]
    except Exception:
        return [content]


# ---------------------------------------------------------------------------
# Low-level drawing primitives
# ---------------------------------------------------------------------------
def _set_run_font(run, latin=None, ea=None):
    """Resolve typeface from the active style: Latin glyphs use the deck's latin
    face, Hangul (east-asian + complex-script) uses the ea face (Korean fallback).
    Both default to the active style set by apply_style()."""
    latin = latin or ACTIVE_FONT["latin"]
    ea = ea or ACTIVE_FONT["ea"]
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", ea)


def _apply_tracking(run, style_key, size_px):
    """Letter-spacing (자간) via OOXML a:rPr@spc — python-pptx exposes no API for it.
    The active look's `type.tracking[tier]` is a fraction of em (negative = tight,
    e.g. brutalism headline -0.02); converted to centipoints (1/100 pt) at the
    rendered size. No override → attribute omitted → unchanged."""
    tr = (_type_override().get("tracking") or {}).get(style_key)
    if not tr:
        return
    spc = int(round(tr * size_px * PX_TO_PT * 100))   # em-fraction → centipoints
    run._r.get_or_add_rPr().set("spc", str(spc))


def text(slide, box, content, style_key, color="navy", align="left",
         anchor="top", wrap=True, line=None, upper=None, shrink=False):
    """Add a textbox. content may be a string or list of (string,) paragraphs."""
    x, y, w, h = box
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if shrink:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.NONE
    pt, bold, lh, up = style(style_key)
    # Composition grammar (empty for skin-only looks → no change):
    #   body_bold   — a look whose body weight is heavy (brutalism: "본문도 700")
    #   upper_tiers — tiers a look renders ALL CAPS (brutalism headlines/labels)
    if GRAMMAR.get("body_bold") and not bold and style_key in _BODY_TIERS:
        bold = True
    if upper is not None:
        up = upper
    elif GRAMMAR.get("upper_tiers") and style_key in GRAMMAR["upper_tiers"]:
        up = True
    # Headline tiers render in the display face (falls back to body face when the
    # look declares no `display` font), so covers/headings can use a display serif.
    is_disp = style_key in DISPLAY_TIERS
    lat = ACTIVE_FONT["latin_display"] if is_disp else ACTIVE_FONT["latin"]
    ea = ACTIVE_FONT["ea_display"] if is_disp else ACTIVE_FONT["ea"]
    # mono-meta grammar: route the META tiers (eyebrow/label/caption) to the
    # monospace face — the dark-tech terminal label layer. No-op for every look
    # that doesn't declare mono_meta, and never touches headline/body/data tiers.
    if _is_mono_tier(style_key):
        lat, ea = _mono_face(is_disp)
    # Korean-safe pre-wrap for headline tiers: break at 어절 boundaries so a word
    # never splits mid-line (PowerPoint's CJK char-level wrap, eaLnBrk ignored).
    sz_px = _size_px(style_key)   # tier size after the look's size-scale override
    if wrap and is_disp and isinstance(content, str):
        content = _prewrap(content, w, sz_px, bold, lat, ea, upper=up)
    al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT}[align]
    lines = content if isinstance(content, list) else [content]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        p.line_spacing = line or lh
        s = (str(ln).upper() if up else str(ln))
        r = p.add_run()
        r.text = s
        r.font.size = pt
        r.font.bold = bold
        r.font.color.rgb = C(color)
        _set_run_font(r, lat, ea)
        _apply_tracking(r, style_key, sz_px)
    return tb


def bullets(slide, box, items, style_key="body", color="ink-2",
            marker=True, marker_color="blue", gap=1.5):
    x, y, w, h = box
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    pt, bold, lh, _ = style(style_key)
    if GRAMMAR.get("body_bold") and not bold and style_key in _BODY_TIERS:
        bold = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap
        p.space_after = Pt(6)
        if marker:
            rb = p.add_run()
            rb.text = "•  "
            rb.font.size = pt
            rb.font.bold = True
            rb.font.color.rgb = C(marker_color)
            _set_run_font(rb)
        r = p.add_run()
        r.text = str(it)
        r.font.size = pt
        r.font.bold = bold
        r.font.color.rgb = C(color)
        _set_run_font(r)
    return tb


def rect(slide, box, fill=None, line_color=None, line_w=0, radius=0):
    x, y, w, h = box
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp_type, IN(x), IN(y), IN(w), IN(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line_color:
        sp.line.color.rgb = C(line_color)
        sp.line.width = Pt(line_w or 1)
    else:
        sp.line.fill.background()
    if radius:
        try:
            sp.adjustments[0] = min(0.5, (radius / 96.0) / (min(w, h) / 96.0))
        except Exception:
            pass
    return sp


_A_NS = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'


def _apply_alpha(sp, pct):
    """솔리드 채움 srgbClr에 a:alpha 삽입 → 반투명(글래스). python-pptx가 fill 알파를
    직접 노출하지 않아 XML로 주입한다. pct=55 → 55% 불투명. 채움이 솔리드일 때만 동작
    (그 외 무시 = 무회귀)."""
    try:
        sf = sp._element.spPr.find(qn("a:solidFill"))
        srgb = sf.find(qn("a:srgbClr")) if sf is not None else None
        if srgb is None:
            return
        for old in srgb.findall(qn("a:alpha")):
            srgb.remove(old)
        srgb.append(parse_xml('<a:alpha %s val="%d"/>' % (_A_NS, _alpha_val(pct))))
    except Exception:
        pass


def _apply_frost(sp, hexv, pct):
    """프로스트 레짐: 솔리드 채움 색을 *흰 틴트*(hexv)로 교체하고 저알파를 건다 →
    busy(다색 메시) 배경 위에서 어떤 메시색이 비쳐도 패널이 균일한 프로스트로 읽힌다.
    surface 색을 그대로 저알파화하면 메시색이 패널 안에서 일렁여 텍스트 대비가 깨지므로
    흰 틴트로 통일한다. 채움이 솔리드일 때만 동작(무회귀)."""
    try:
        sf = sp._element.spPr.find(qn("a:solidFill"))
        srgb = sf.find(qn("a:srgbClr")) if sf is not None else None
        if srgb is None:
            return
        srgb.set("val", str(hexv).lstrip("#").upper())
        for old in srgb.findall(qn("a:alpha")):
            srgb.remove(old)
        srgb.append(parse_xml('<a:alpha %s val="%d"/>' % (_A_NS, _alpha_val(pct))))
    except Exception:
        pass


# --- radial corner glow primitive (grammar.radial_glow) ------------------------
# The #1 identity cue of the ENGINEERED-DARK pack: "거의 검정 배경에 슬라이드당 단 1개
# 의 추상 라디얼 글로우(보라-청록)가 코너 한 곳에서 은은히 퍼진다". This is the opposite
# texture from dark-tech's cyber GRID — engineered-dark FORBIDS grids/patterns and instead
# floats one soft radial bloom in a corner. Implemented as a big rounded shape filled with
# an OOXML radial <a:gradFill path="circle"> whose centre stop is the violet→teal accent at
# low alpha and whose outer stop fades to 0% alpha (transparent → the charcoal canvas shows
# through). The focal rect (<a:fillToRect>) is pinned to the declared corner so the bloom
# blooms FROM that corner. Gated HARD on grammar.radial_glow + dark canvas; every other look
# adds zero shapes (byte-identical). Corner index alternates per-slide so the deck doesn't
# put the bloom in the same corner every page (the pack: "코너 또는 가장자리에만", 1/slide).
_RG_CORNERS = ("tr", "bl", "tl", "br")          # cycle order: top-right first (per 적용 예)


def _radial_glow(slide, idx=0):
    """Paint one soft corner radial bloom for an engineered-dark slide. No-op unless
    grammar.radial_glow is declared AND the canvas is dark.

    grammar.radial_glow (dict, all optional):
      from     — centre hue token/hex (default 'blue' = soft violet #8B7BF0)
      to       — outer-flow hue token/hex (default 'teal' #3FB8C4 — the violet→teal drift)
      alpha    — centre opacity 0..100 (default 18, per spec '중심 18% → 외곽 0%')
      size     — bloom diameter in px (default 760 ≈ a generous corner wash)
      corner   — fixed corner 'tr'|'tl'|'br'|'bl'; default = cycle by slide idx"""
    cfg = GRAMMAR.get("radial_glow")
    if not cfg or not DARK:
        return
    if cfg is True:
        cfg = {}
    c_from = _resolve(cfg.get("from", "blue")).lstrip("#")
    c_to = _resolve(cfg.get("to", "teal")).lstrip("#")
    a_ctr = int(max(0, min(100, cfg.get("alpha", 18))))
    size = cfg.get("size", 760)
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    corner = cfg.get("corner") or _RG_CORNERS[idx % len(_RG_CORNERS)]
    # Place the bloom so its centre sits in/near the chosen corner (bleed off-canvas).
    half = size / 2.0
    cx = {"tr": W, "br": W, "tl": 0, "bl": 0}[corner]
    cy = {"tr": 0, "tl": 0, "br": H, "bl": H}[corner]
    x, y = cx - half, cy - half
    # focal point (0..100000) inside the shape that the radial blooms FROM = the corner.
    fx = 100000 if corner in ("tr", "br") else 0
    fy = 0 if corner in ("tr", "tl") else 100000
    sp = oval(slide, [x, y, size, size], fill=("#" + c_from))
    # 3-stop radial: centre = from@alpha, mid = to@~half-alpha (the violet→teal drift),
    # outer = to@0% (transparent). fillToRect pins the bright focus to the corner.
    grad = parse_xml(
        '<a:gradFill %s rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '<a:gs pos="55000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="%s"><a:alpha val="0"/></a:srgbClr></a:gs>'
        '</a:gsLst><a:path path="circle"><a:fillToRect l="%d" t="%d" r="%d" b="%d"/>'
        '</a:path></a:gradFill>'
        % (_A_NS, c_from, a_ctr * 1000, c_to, int(a_ctr * 0.6) * 1000, c_to,
           fx, fy, 100000 - fx, 100000 - fy))
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious if ln is not None else spPr.append)(grad)
    # no outline on the bloom — it's pure light, not a disc.
    if ln is None:
        spPr.append(parse_xml('<a:ln %s><a:noFill/></a:ln>' % _A_NS))
    return sp


# --- soft 3D sphere decor primitive (grammar.orb_decor) ------------------------
# The msit/government 'soft sphere' motif: translucent balls with an offset specular
# highlight float as decor (cover/section), reading as frosted 3D orbs. Unlike radial_glow
# (a DARK-only corner WASH that bleeds off-canvas), these are PLACEABLE discrete spheres
# that work on a LIGHT canvas too — a radial <a:gradFill path="circle"> whose focal point
# is pinned UPPER-LEFT (the light source) so the highlight sits off-centre like a real ball,
# fading to the base hue at the rim and a soft (low-alpha) edge so there is no hard outline.
# Gated HARD on grammar.orb_decor; drawn once per FRAME slide (cover/section/closing) right
# after the bg, below content. Every other look adds zero shapes (byte-identical).
def _orb_decor(slide):
    """Draw the look's decorative soft spheres (grammar.orb_decor). No-op unless declared.

    grammar.orb_decor (dict):
      spheres — [{at:[fx,fy] 0..1 centre, r: radius px, color: base hue token/hex,
                  highlight: token/hex (default '#FFFFFF'), alpha: 0..100 overall opacity
                  (default 100), edge: 0..1 rim-alpha fraction for the soft edge (default
                  0.55), focal: [l,t] 0..100 specular position (default [22,22] upper-left)}]"""
    cfg = GRAMMAR.get("orb_decor")
    if not cfg:
        return
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    for s in (cfg.get("spheres") or []):
        fx, fy = (s.get("at") or [0.5, 0.5])[:2]
        r = float(s.get("r", 180))
        base = _resolve(s.get("color", "blue-light")).lstrip("#")
        hi = _resolve(s.get("highlight", "#FFFFFF")).lstrip("#")
        a = int(max(0, min(100, s.get("alpha", 100))))
        edge = float(s.get("edge", 0.55))
        fl, ft = (s.get("focal") or [22, 22])[:2]
        x, y = fx * W - r, fy * H - r
        sp = oval(slide, [x, y, 2 * r, 2 * r], fill="#" + base)
        grad = parse_xml(
            '<a:gradFill %s rotWithShape="1"><a:gsLst>'
            '<a:gs pos="0"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
            '<a:gs pos="55000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
            '<a:gs pos="100000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
            '</a:gsLst><a:path path="circle">'
            '<a:fillToRect l="%d" t="%d" r="%d" b="%d"/></a:path></a:gradFill>'
            % (_A_NS, hi, _alpha_val(a * 1000), base, _alpha_val(a * 1000),
               base, _alpha_val(int(a * edge) * 1000),
               int(fl * 1000), int(ft * 1000),
               int((100 - fl) * 1000), int((100 - ft) * 1000)))
        spPr = sp._element.spPr
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
            e = spPr.find(qn(tag))
            if e is not None:
                spPr.remove(e)
        ln = spPr.find(qn("a:ln"))
        (ln.addprevious if ln is not None else spPr.append)(grad)
        if ln is None:
            spPr.append(parse_xml('<a:ln %s><a:noFill/></a:ln>' % _A_NS))
    return None


def _gradient_bg(slide, colors):
    """풀블리드 다색 대각 선형 그라디언트 배경. 4정점 메시는 PPTX가 직접 지원하지 않아
    대각 선형(45°) 다중 스톱으로 근사한다(글래스모피즘 배경의 정직한 한계). colors=hex 목록."""
    sp = rect(slide, [0, 0, SLIDE["width_px"], SLIDE["height_px"]], fill=colors[0])
    n = max(2, len(colors))
    gs = "".join(
        '<a:gs pos="%d"><a:srgbClr val="%s"/></a:gs>'
        % (int(i * 100000 / (n - 1)), str(c).lstrip("#"))
        for i, c in enumerate(colors))
    grad = parse_xml(
        '<a:gradFill %s><a:gsLst>%s</a:gsLst>'
        '<a:lin ang="2700000" scaled="1"/></a:gradFill>' % (_A_NS, gs))
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    return sp


# --- baked mesh background primitive (grammar.bg_mesh) -------------------------
# A TRUE multi-vertex mesh gradient is the defining identity of the glassmorphism /
# vivid-gradient / iridescent looks: distinct hues anchored at the four CORNERS blend
# into a 2D color FIELD (violet TL, blue TR, pink BL, teal BR) — optionally with soft
# additive radial BLOOMS floating inside it. OOXML cannot express a 4-corner mesh; the
# old `_gradient_bg` flattened the same colors into a single diagonal LINEAR sweep,
# collapsing the field to one axis (the teal corner vanished, all looked like a ramp).
# So we BAKE the mesh deterministically (numpy bilinear corners + additive blooms, the
# same math as gen-texture.py) to a content-addressed PNG and drop it full-bleed behind
# the content via add_picture. Determinism: identical (verts, blooms, size, blur) → same
# bytes (no PNG timestamp) → python-pptx dedups the embedded media part by blob hash, so
# every slide of a deck shares ONE image part. Gated HARD on grammar.bg_mesh; every other
# look adds zero shapes (byte-identical). numpy/PIL missing → graceful linear fallback.
_MESH_CACHE = {}                          # digest -> baked file path (per build process)
_MESH_DIR = os.path.join(tempfile.gettempdir(), "pptlab-mesh")


def _bake_mesh(verts, blooms, w, h, blur=0.0):
    """Deterministically bake a 4-corner bilinear mesh (+ additive radial blooms) to a
    PNG and return its path. verts = [TL, TR, BL, BR] hex (short lists are padded by
    repeating the last). blooms = [{at:[fx,fy] 0..1, color:hex, intensity:0..~1,
    radius:0..~1 of min(w,h)}]. Content-addressed → byte-stable, reused across slides."""
    import numpy as np                    # lazy: only mesh looks pay the import
    from PIL import Image, ImageFilter
    key = json.dumps([verts, blooms, w, h, round(float(blur), 3)], sort_keys=True)
    digest = hashlib.sha1(key.encode("utf-8")).hexdigest()[:16]
    if digest in _MESH_CACHE:
        return _MESH_CACHE[digest]
    os.makedirs(_MESH_DIR, exist_ok=True)
    path = os.path.join(_MESH_DIR, "mesh-%s.png" % digest)
    if not os.path.exists(path):
        cols = [np.array(_hex2rgb(c), float) for c in (verts or ["#FFFFFF"])]
        while len(cols) < 4:
            cols.append(cols[-1])
        tl, tr, bl, br = cols[:4]
        x = np.linspace(0, 1, w)[None, :, None]
        y = np.linspace(0, 1, h)[:, None, None]
        top = tl * (1 - x) + tr * x
        bot = bl * (1 - x) + br * x
        img = top * (1 - y) + bot * y                 # bilinear corner field
        if blooms:
            yy, xx = np.mgrid[0:h, 0:w]
            for b in blooms:
                fx, fy = (b.get("at") or [0.5, 0.5])[:2]
                col = np.array(_hex2rgb(b.get("color", "#FFFFFF")), float)
                inten = float(b.get("intensity", 0.4))
                rad = float(b.get("radius", 0.55)) * min(w, h)
                dist = np.sqrt((xx - fx * w) ** 2 + (yy - fy * h) ** 2)
                fall = np.clip(1 - dist / max(rad, 1.0), 0, 1) ** 2.0
                img = img + col[None, None, :] * fall[..., None] * inten   # additive light
        im = Image.fromarray(np.clip(img, 0, 255).astype("uint8"), "RGB")
        if blur and blur > 0:
            im = im.filter(ImageFilter.GaussianBlur(blur))
        im.save(path)
    _MESH_CACHE[digest] = path
    return path


def _mesh_bg(slide, cfg):
    """Full-bleed baked mesh background for grammar.bg_mesh. cfg=True → derive the four
    corner colors from grammar.bg_gradient (so a look opts in with one extra key). cfg=
    dict → {verts, blooms, blur} explicit. Falls back to the linear gradient if numpy/PIL
    is unavailable so the build never crashes on a thin environment."""
    if cfg is True:
        cfg = {}
    verts = cfg.get("verts") or GRAMMAR.get("bg_gradient") or ["canvas"]
    verts = [_resolve(c) for c in verts]
    blooms = [dict(b, color=_resolve(b.get("color", "#FFFFFF")))
              for b in (cfg.get("blooms") or [])]
    blur = float(cfg.get("blur", 0) or 0)
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    try:
        path = _bake_mesh(verts, blooms, W, H, blur)
        slide.shapes.add_picture(path, IN(0), IN(0), IN(W), IN(H))
    except Exception:
        _gradient_bg(slide, verts)        # graceful degrade (no numpy/PIL): linear sweep
    return None


# --- prism gradient primitive (grammar.prism) ----------------------------------
# The #1 identity of the prismatic-dark pack: shapes are "drawn with light, not
# filled" — a node has a near-transparent fill + a 1px CYAN→MAGENTA→AMBER gradient
# STROKE, and big KPI numbers / section watermarks render in gradient TEXT fill.
# Both are OOXML <a:gradFill> (verified to render in PowerPoint COM, 2026-06-13 PoC):
#   • line stroke  → <a:ln><a:gradFill><a:gsLst>3 stops</a:gsLst><a:lin/></a:ln>
#   • text run fill → run <a:rPr><a:gradFill>...</a:rPr>
# Prism is ALWAYS the 3-stop cyan→magenta→amber set (the pack: "프리즘은 시안·마젠타·
# 앰버 3스톱 1세트뿐 — 무지개로 늘리지 않는다"). Stops resolve from the look's accent
# roles (ACCENTS[0]=cyan, ACCENTS[1]=magenta) + the 3rd accent literal (amber), so a
# palette swap can't leak an off-prism hue. Every seam is gated on grammar.prism +
# DARK → hard no-op for the other 109 looks (byte-identical regression).
def _prism_stops():
    """The 3 prism hex stops from the active palette's accent roles. ACCENTS holds
    role names (blue/orange) + the look's literal 3rd accent (amber #FBBF24); each is
    resolved through the palette so the gradient always tracks the look's colors. Falls
    back to the canonical cyan/magenta/amber literals if a look defines fewer accents."""
    canon = ["#22D3EE", "#D946EF", "#FBBF24"]
    out = []
    for i in range(3):
        tok = ACCENTS[i] if i < len(ACCENTS) else canon[i]
        try:
            hexv = _resolve(tok)
        except Exception:
            hexv = canon[i]
        out.append(str(hexv).lstrip("#").upper())
    return out


def _grad_gs_lst(stops):
    """<a:gsLst> string for evenly-spaced stops (list of bare hex strings)."""
    n = max(2, len(stops))
    return "".join(
        '<a:gs pos="%d"><a:srgbClr val="%s"/></a:gs>'
        % (int(i * 100000 / (n - 1)), str(c).lstrip("#").upper())
        for i, c in enumerate(stops))


def _prism_cfg():
    """The active look's prism grammar as a dict, or None when not declared / light
    canvas → every prism seam is a hard no-op. Keys (all optional):
      node_alpha — node fill opacity 0..100 (default 20 = near-transparent)
      stroke_pt  — gradient stroke width (default 1.5)
      ang        — gradient angle in 60000ths-deg (default 0 = left→right)"""
    p = GRAMMAR.get("prism")
    if not p or not DARK:
        return None
    return {} if p is True else p


def gradient_stroke(sp, width_pt=None, stops=None, ang=0):
    """Paint a shape/connector OUTLINE with the prism linear gradient. Replaces the
    line fill with <a:ln><a:gradFill>3 stops</a:gradFill>. ang in 60000ths of a degree
    (0 = left→right). Verified to render in PowerPoint COM. Returns sp unchanged on any
    failure (fail-safe — falls back to whatever line was already set)."""
    try:
        cfg = _prism_cfg() or {}
        if stops is None:
            stops = _prism_stops()
        w = width_pt if width_pt is not None else cfg.get("stroke_pt", 1.5)
        spPr = sp._element.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = spPr.makeelement(qn("a:ln"), {})
            spPr.append(ln)
        ln.set("w", str(int(w * 12700)))
        ln.set("cap", "rnd")
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:pattFill"):
            e = ln.find(qn(tag))
            if e is not None:
                ln.remove(e)
        grad = parse_xml('<a:gradFill %s><a:gsLst>%s</a:gsLst>'
                         '<a:lin ang="%d" scaled="1"/></a:gradFill>'
                         % (_A_NS, _grad_gs_lst(stops), int(ang)))
        ln.insert(0, grad)   # fill group is the first child of a:ln
    except Exception:
        pass
    return sp


def _gradient_text_run(run, stops=None, ang=0):
    """Replace a text run's solid font fill with the prism gradient (run <a:rPr>
    <a:gradFill>). The numeric/letter glyphs are unchanged — only the fill. In the
    CT_TextCharacterProperties schema the fill group sits after a:ln and before
    a:effectLst, so we insert right after a:ln when present (PoC-verified order)."""
    try:
        if stops is None:
            stops = _prism_stops()
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill", "a:blipFill"):
            e = rPr.find(qn(tag))
            if e is not None:
                rPr.remove(e)
        grad = parse_xml('<a:gradFill %s><a:gsLst>%s</a:gsLst>'
                         '<a:lin ang="%d" scaled="1"/></a:gradFill>'
                         % (_A_NS, _grad_gs_lst(stops), int(ang)))
        ln = rPr.find(qn("a:ln"))
        if ln is not None:
            ln.addnext(grad)
        else:
            rPr.insert(0, grad)
    except Exception:
        pass
    return run


def gradient_text(slide, box, content, style_key, align="left", anchor="top",
                  upper=None, ang=0):
    """Add a textbox whose run(s) render in the prism gradient FILL — KPI big numbers
    and section watermarks. Mirrors text()'s sizing/face/prewrap but swaps the solid
    color for a gradient run fill. No glow on text (the pack forbids it for legibility).
    Caller gates on _prism_cfg(); this just draws."""
    x, y, w, h = box
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    pt, bold, lh, up = style(style_key)
    if upper is not None:
        up = upper
    is_disp = style_key in DISPLAY_TIERS
    lat = ACTIVE_FONT["latin_display"] if is_disp else ACTIVE_FONT["latin"]
    ea = ACTIVE_FONT["ea_display"] if is_disp else ACTIVE_FONT["ea"]
    al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT}[align]
    p = tf.paragraphs[0]
    p.alignment = al
    p.line_spacing = lh
    r = p.add_run()
    r.text = str(content).upper() if up else str(content)
    r.font.size = pt
    r.font.bold = bold
    r.font.color.rgb = C("ink")        # baseline fill (overridden by gradient)
    _set_run_font(r, lat, ea)
    _gradient_text_run(r, ang=ang)
    return tb


def _gradient_shape_fill(sp, stops=None, ang=2700000, alpha=None):
    """Replace a shape's solid FILL with the prism linear gradient (chart bars/area).
    ang in 60000ths-deg (default 2700000 = bottom→top so the bright cap is at the top,
    per the pack's '윗면 그라디언트 글로우 캡'). alpha (0..100) applies a uniform stop
    opacity so an area fill can fade. Fail-safe no-op on error."""
    try:
        if stops is None:
            stops = _prism_stops()
        a = '' if alpha is None else ('<a:alpha %s val="%d"/>'
                                      % (_A_NS, int(max(0, min(100, alpha)) * 1000)))
        n = max(2, len(stops))
        gs = "".join(
            '<a:gs pos="%d"><a:srgbClr val="%s">%s</a:srgbClr></a:gs>'
            % (int(i * 100000 / (n - 1)), str(c).lstrip("#").upper(), a)
            for i, c in enumerate(stops))
        spPr = sp._element.spPr
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"):
            e = spPr.find(qn(tag))
            if e is not None:
                spPr.remove(e)
        grad = parse_xml('<a:gradFill %s><a:gsLst>%s</a:gsLst>'
                         '<a:lin ang="%d" scaled="1"/></a:gradFill>'
                         % (_A_NS, gs, int(ang)))
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            ln.addprevious(grad)
        else:
            spPr.append(grad)
    except Exception:
        pass
    return sp


def _effect_lst(spPr):
    """Get (or create) the shape's <a:effectLst>, kept as the single shared effect
    container so glow + shadow can coexist on one shape (PowerPoint renders both).
    effectLst is the last child of CT_ShapeProperties, so appending is order-safe."""
    eff = spPr.find(qn("a:effectLst"))
    if eff is None:
        eff = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(eff)
    return eff


def _alpha_val(pct):
    """Clamp an opacity pct (0..100) to the OOXML a:alpha val range (0..100000).
    SINGLE seam for every a:alpha so a look setting outside 0..100 can never push the
    value past PowerPoint's bounds (out-of-range a:alpha → file refuses to open). For
    any already-valid pct the result equals the old int(pct*1000), so no regression."""
    return int(max(0, min(100, pct)) * 1000)


def _safe_hex(color):
    """Resolve a token/#hex to a 6-char uppercase hex, or None if it can't (fail-safe).
    PowerPoint rejects a non-hex srgbClr val and won't open the file, so every srgbClr
    seam gates through this — a bad color becomes a no-op instead of a corrupt deck."""
    try:
        hx = _resolve(color).lstrip("#").upper()
    except Exception:
        return None
    if len(hx) == 6 and all(c in "0123456789ABCDEF" for c in hx):
        return hx
    return None


# CT_EffectList requires its children in this schema order; an out-of-order list makes
# PowerPoint refuse the file. Effects arrive from several seams (glow inserts at front,
# neumorphic innerShdw inserts at front, shadow/chart-glow append/insert) so the final
# order can drift when two coexist on one shape — _sort_effects normalizes after every
# insert. With a single effect (every real look today) it is a no-op = byte-identical.
_EFFECT_ORDER = ("a:blur", "a:fillOverlay", "a:glow", "a:innerShdw",
                 "a:outerShdw", "a:prstShdw", "a:reflection", "a:softEdge")


def _sort_effects(eff):
    """Reorder an effectLst's children into CT_EffectList schema order. Stable (keeps
    same-tag siblings' relative order) and fail-safe (no-op on any error)."""
    try:
        rank = {qn(t): i for i, t in enumerate(_EFFECT_ORDER)}
        for c in sorted(list(eff), key=lambda e: rank.get(e.tag, len(_EFFECT_ORDER))):
            eff.append(c)        # re-appending an existing node MOVES it → reorders
    except Exception:
        pass


# --- hatch fill primitive (grammar.hatch_fill) ---------------------------------
# The hand-drawn / sketch identity is a SHADED HATCH, not a flat fill: card faces are
# filled with fine diagonal strokes (pencil shading) instead of a solid block. Built as
# OOXML <a:pattFill prst="ltUpDiag"> (a native PowerPoint preset pattern) so it renders
# deterministically with no image asset. fg = stroke color (ink), bg = the paper behind
# it (the card's own fill) — both via role tokens so the hatch adapts to canvas luminance
# (no hardcoded hex → never "vanishes on the opposite luminance"). Gated on
# grammar.hatch_fill; every other look adds zero pattFills (byte-identical). First
# fill-grammar primitive of the engine-gaps "needs-brush" class.
_HATCH_PRSTS = frozenset({
    "ltUpDiag", "dkUpDiag", "ltDnDiag", "dkDnDiag", "wdUpDiag", "wdDnDiag",
    "ltHorz", "ltVert", "smGrid", "diagCross",
})


def hatch_fill(sp, prst="ltUpDiag", fg="ink", bg=None, fg_alpha=None):
    """Replace a shape's fill with a diagonal hatch <a:pattFill>. fg/bg are role tokens
    or #hex; fg defaults to 'ink' (canvas-aware stroke), bg defaults to the shape's
    current solid fill (else 'canvas'). prst is validated against a hatch whitelist
    (falls back to ltUpDiag). No-op (fail-safe) when fg/bg can't resolve to hex, or when
    fg == bg (the hatch would be invisible — leave the existing fill)."""
    try:
        spPr = sp._element.spPr
        if bg is None:
            cur = spPr.find(qn("a:solidFill"))
            srgb = cur.find(qn("a:srgbClr")) if cur is not None else None
            bg = ("#" + srgb.get("val")) if (srgb is not None and srgb.get("val")) else "canvas"
        fg_hex = _safe_hex(fg)
        bg_hex = _safe_hex(bg)
        if fg_hex is None or bg_hex is None or fg_hex == bg_hex:
            return sp
        if prst not in _HATCH_PRSTS:
            prst = "ltUpDiag"
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                    "a:pattFill", "a:grpFill"):
            e = spPr.find(qn(tag))
            if e is not None:
                spPr.remove(e)
        fg_a = ('<a:alpha val="%d"/>' % _alpha_val(fg_alpha)) if fg_alpha is not None else ""
        patt = parse_xml(
            '<a:pattFill %s prst="%s">'
            '<a:fgClr><a:srgbClr val="%s">%s</a:srgbClr></a:fgClr>'
            '<a:bgClr><a:srgbClr val="%s"/></a:bgClr>'
            '</a:pattFill>' % (_A_NS, prst, fg_hex, fg_a, bg_hex))
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            ln.addprevious(patt)               # fill precedes a:ln in CT_ShapeProperties
        else:
            eff = spPr.find(qn("a:effectLst"))
            eff.addprevious(patt) if eff is not None else spPr.append(patt)
    except Exception:
        pass
    return sp


def _hatch_cfg():
    """Active look's hatch grammar (dict) or None when undeclared → every hatch seam is
    a hard no-op (byte-identical). Keys (optional): prst, fg, bg, alpha. Unlike glow,
    hatch is NOT dark-gated — the hand-drawn look is light."""
    h = GRAMMAR.get("hatch_fill")
    if not h:
        return None
    if h is True:
        h = {}
    return h


# --- wobble path primitive (grammar.wobble_path) -------------------------------
# The hand-drawn / sketch identity is a SHAKY HAND LINE, not a clean rect: card
# outlines look drawn by hand — edges bow slightly, corners overshoot. Built as a
# custGeom freeform (FreeformBuilder) whose rectangle perimeter is subdivided and
# each point nudged by a FIXED jitter sequence (no RNG → byte-stable, reproducible).
# Fill = the card color, stroke = ink (accent cards stroke in accent). Gated on
# grammar.wobble_path; every other look draws a clean rect (byte-identical no-op).
_WOBBLE_SEQ = (0.7, -1.0, 0.5, -0.8, 0.9, -0.6, 0.8, -0.7, 0.6, -0.9, 0.7, -0.5)


def _wobble_vertices(x, y, w, h, seg=3, amp=3.5, phase=0):
    """Deterministic wobbly-rectangle perimeter vertices (px). Each edge split into
    `seg` segments; every point offset by the fixed jitter sequence (perp to its edge)
    so the outline reads hand-drawn yet is byte-stable. `phase` shifts the sequence so
    a second overlaid stroke can wobble differently (sketch double-line)."""
    idx = [phase]

    def jit():
        v = _WOBBLE_SEQ[idx[0] % len(_WOBBLE_SEQ)] * amp
        idx[0] += 1
        return v
    pts = []
    for s in range(seg):
        pts.append((x + w * s / seg, y + jit()))            # top  L→R (perp = y)
    for s in range(seg):
        pts.append((x + w + jit(), y + h * s / seg))        # right T→B (perp = x)
    for s in range(seg):
        pts.append((x + w - w * s / seg, y + h + jit()))    # bottom R→L
    for s in range(seg):
        pts.append((x + jit(), y + h - h * s / seg))        # left  B→T
    return pts


def _freeform(slide, verts_px, fill=None, line_color=None, line_w=0):
    """Build a CLOSED freeform shape from px vertices (custGeom via FreeformBuilder).
    fill/line_color are tokens or #hex (None = no fill / no line). px → EMU at 9525/px."""
    fb = slide.shapes.build_freeform(float(verts_px[0][0]), float(verts_px[0][1]), scale=9525)
    fb.add_line_segments([(float(a), float(b)) for a, b in verts_px[1:]], close=True)
    sp = fb.convert_to_shape()
    if fill is not None:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    else:
        sp.fill.background()
    if line_color is not None and line_w:
        sp.line.color.rgb = C(line_color)
        sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    return sp


def _wobble_cfg():
    """Active look's wobble grammar (dict) or None → clean-rect no-op (byte-identical).
    Keys (optional): amp(px jitter), seg(segments/edge), sketch(bool double-stroke)."""
    w = GRAMMAR.get("wobble_path")
    if not w:
        return None
    if w is True:
        w = {}
    return w


def _wobble_card(slide, box, fillv, line_color, line_w, preset, cfg):
    """Hand-drawn card: a wobbly freeform rect (card fill + sketch ink stroke) carrying
    the preset's drop shadow. The wobble IS the look's identity (vs a clean rect). An
    optional second offset stroke (cfg.sketch) gives the 'drawn twice' pencil feel."""
    x, y, w, h = box
    amp = cfg.get("amp", 3.5)
    seg = cfg.get("seg", 3)
    lw = (line_w or 0) or 1.75
    lc = line_color or "ink"
    sp = _freeform(slide, _wobble_vertices(x, y, w, h, seg, amp),
                   fill=fillv, line_color=lc, line_w=lw)
    sh = preset.get("shadow")
    if sh:
        shadow(sp, sh.get("blur_pt", 9), sh.get("dist_pt", 3), sh.get("dir_deg", 90),
               sh.get("color", "#1D1B20"), sh.get("alpha", 82))
    if cfg.get("sketch"):
        # second, slightly-offset line-only stroke = sketchy double outline
        _freeform(slide, _wobble_vertices(x, y, w, h, seg, amp * 0.6, phase=3),
                  fill=None, line_color=lc, line_w=max(0.75, lw * 0.6))
    return sp


def shadow(sp, blur_pt=12, dist_pt=4, dir_deg=90, color="#000000", alpha=78):
    """Inject an outer drop shadow via OOXML — python-pptx's high-level API does
    not expose outerShdw (the same reason extract-style.py reads it via lxml).
    blur/dist in pt, dir in degrees, alpha 0..100 (opacity). Merges into the shared
    effectLst (a glow may already be present) and removes any prior outerShdw. A bad
    color or out-of-range alpha is skipped/clamped (fail-safe) — PowerPoint refuses a
    non-hex srgbClr or out-of-range a:alpha and won't open the file."""
    try:
        hx = _safe_hex(color)
        if hx is None:
            return sp
        spPr = sp._element.spPr
        eff = _effect_lst(spPr)
        for old in eff.findall(qn("a:outerShdw")):
            eff.remove(old)
        sh = eff.makeelement(qn("a:outerShdw"), {
            "blurRad": str(int(blur_pt * 12700)),
            "dist": str(int(dist_pt * 12700)),
            "dir": str(int(dir_deg * 60000)),
            "rotWithShape": "0",
        })
        clr = sh.makeelement(qn("a:srgbClr"), {"val": hx})
        a = clr.makeelement(qn("a:alpha"), {"val": str(_alpha_val(100 - alpha))})
        clr.append(a)
        sh.append(clr)
        eff.append(sh)
        _sort_effects(eff)
    except Exception:
        pass
    return sp


# --- glow primitive (grammar.glow) ---------------------------------------------
# A luminous outer halo via OOXML <a:glow> — PowerPoint renders this natively
# (verified by COM render; python-pptx has no high-level API, same as outerShdw).
# This is the visual primitive that defines the GLOW kit: cards/badges/KPI slabs/
# chart marks and connector lines on a DARK canvas get a soft same-hue halo, the
# neon-tech identity that brutalism/swiss/luxe never had. ⚠ glow reads as halo only
# on dark backgrounds — every caller is gated so it NEVER fires on a light look.
def glow(sp, rad_px=18, color=None, alpha=60):
    """Attach an <a:glow rad><srgbClr><alpha></a:glow> to a shape/connector.
    `color` defaults to the shape's own line color (else its fill), so a teal card
    halos teal and a violet badge halos violet — automatically on-palette, no per-
    call color list. alpha is the glow opacity 0..100. effectLst order: glow must
    precede outerShdw, so glow is inserted at the front of the shared list."""
    try:
        spPr = sp._element.spPr
        if color is None:
            color = _glow_color_of(sp)
        if color is None:
            return sp
        # color may be a TOKEN name (e.g. "blue"/"divider") — resolve to #hex.
        # PowerPoint rejects a non-hex srgbClr val and refuses to open the file.
        hexv = _resolve(color).lstrip("#").upper()
        if len(hexv) != 6 or any(c not in "0123456789ABCDEF" for c in hexv):
            return sp                          # not a valid hex → skip (fail-safe)
        eff = _effect_lst(spPr)
        for old in eff.findall(qn("a:glow")):
            eff.remove(old)
        g = eff.makeelement(qn("a:glow"), {"rad": str(int(rad_px * 12700))})
        clr = g.makeelement(qn("a:srgbClr"), {"val": hexv})
        a = clr.makeelement(qn("a:alpha"), {"val": str(_alpha_val(alpha))})
        clr.append(a)
        g.append(clr)
        eff.insert(0, g)               # glow before any outerShdw (schema order)
        _sort_effects(eff)
    except Exception:
        pass
    return sp


def _glow_color_of(sp):
    """Best hex for a shape's own halo: its outline color, else its solid fill.
    Returns a #hex string or None. Reads the rendered spPr so the halo always
    matches what's drawn (e.g. an accent-bordered dark card halos in the accent)."""
    try:
        spPr = sp._element.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            sf = ln.find(qn("a:solidFill"))
            srgb = sf.find(qn("a:srgbClr")) if sf is not None else None
            if srgb is not None and srgb.get("val"):
                return "#" + srgb.get("val")
        sf = spPr.find(qn("a:solidFill"))
        srgb = sf.find(qn("a:srgbClr")) if sf is not None else None
        if srgb is not None and srgb.get("val"):
            return "#" + srgb.get("val")
    except Exception:
        pass
    return None


def _glow_cfg():
    """The active look's glow grammar as a normalized dict, or None when the look
    declares no glow (→ every glow seam is a hard no-op = byte-identical regress).
    Keys (all optional): rad(px), alpha(0..100), cards, badges, lines, chart (bools,
    default True) gate WHICH element families halo; color (hex) forces a fixed hue
    instead of per-shape derivation. ⚠ Glow only ever applies on a DARK canvas —
    if the look isn't dark, returns None so light looks can't leak a halo."""
    g = GRAMMAR.get("glow")
    if not g or not DARK:
        return None
    if g is True:
        g = {}
    return g


def _glow_shape(sp, family, color=None):
    """Apply the active look's glow to a drawn shape IF its family is enabled.
    `family` in {cards, badges, lines, chart}. Central gate so every seam stays a
    one-liner and the on/off decision lives in one place (DRY)."""
    cfg = _glow_cfg()
    if cfg is None or not cfg.get(family, True):
        return sp
    return glow(sp, rad_px=cfg.get("rad", 18),
                color=color or cfg.get("color"),
                alpha=cfg.get("alpha", 60))


# --- neumorphic elevation primitive (card.elevation:"neumorphic") --------------
# Soft-UI depth is carried by an OPPOSED shadow pair, NOT by color: the fill == the
# canvas (contrast 0), so only light + shadow sculpt the form. Light is assumed from
# the TOP-LEFT (NW), so the two volume directions are exact mirrors:
#   CONVEX  (protrudes): OUTER dark cast SE(45°) + OUTER light highlight NW(225°).
#   CONCAVE (recedes):   INNER dark on the near NW(225°) wall.
# 🔑 Why this emits only ONE shadow per shape (and card() STACKS shapes for the pair):
# a single <a:effectLst> legally holds only one outerShdw, AND — verified by COM PNG
# render — a shape carrying an <a:effectDag> (the container that *would* allow a pair)
# has its solid FILL DROPPED to gray on export. So effectDag is unusable for any
# filled card. The showeet source itself stacks one-shadow shapes; we do the same,
# which preserves fill for ANY canvas/accent color. (engine-gaps.md G1.) First
# primitive grown by the visual-QA self-reinforcing loop.
def neumorphic_shadow(sp, kind="outer", dir_deg=45, color="#1A1B22", alpha=40,
                      blur_pt=20, dist_pt=15):
    """Emit ONE soft shadow into the shape's shared effectLst. kind='outer'(drop)|
    'inner'(recess). alpha = OPACITY 0..100. Colors may be TOKEN names → #hex; a bad
    hex is a no-op (fail-safe). Never uses effectDag (it drops the fill in COM render)
    — opposed pairs are built by STACKING shapes in card()/_neumorphic_card."""
    try:
        spPr = sp._element.spPr
        eff = _effect_lst(spPr)
        tag = "a:innerShdw" if kind == "inner" else "a:outerShdw"
        hx = _resolve(color).lstrip("#").upper()
        if len(hx) != 6 or any(ch not in "0123456789ABCDEF" for ch in hx):
            return sp                          # not a valid hex → skip (fail-safe)
        attrs = {"blurRad": str(int(blur_pt * 12700)),
                 "dist": str(int(dist_pt * 12700)),
                 "dir": str(int(dir_deg * 60000)), "rotWithShape": "0"}
        if tag == "a:outerShdw":
            attrs["algn"] = "tl"
        sh = eff.makeelement(qn(tag), attrs)
        clr = sh.makeelement(qn("a:srgbClr"), {"val": hx})
        a = clr.makeelement(qn("a:alpha"),
                            {"val": str(_alpha_val(alpha))})
        clr.append(a)
        sh.append(clr)
        # innerShdw must precede outerShdw in CT_EffectList schema order.
        eff.insert(0, sh) if tag == "a:innerShdw" else eff.append(sh)
        _sort_effects(eff)
    except Exception:
        pass
    return sp


def derive_neumorphic(canvas_hex):
    """Deterministically derive a neumorphic shadow recipe FROM THE CANVAS, so a look
    gets balanced soft-UI depth by declaring only its canvas (mirrors derive_dark_roles
    — measure/derive, never hand-tune per look). A look OVERRIDES any key via
    grammar.neumorphic, which is where extract-style.py writes the values MEASURED from
    an absorbed source. The STRUCTURE is fixed engine logic; only the magnitudes here,
    and those follow the canvas:
      dark  = the canvas hue driven toward black — a hue-true deep shade, so the shadow
              is the surface's OWN shadow, not a foreign gray (matches the source's
              tx1@low-lum darks).
      light = white, the highlight a sculpted surface catches.
      depth = dark_alpha scales with canvas luminance: a brighter canvas has more tonal
              room beneath it, so it must cast a deeper shadow or it looks flat (this is
              why a near-white deck needs MORE shadow %, not less). blur/dist seed from
              the measured source defaults (20/15pt) and read as a soft dish."""
    L = _rel_lum(canvas_hex)                         # 0..1
    dark = _mix(canvas_hex, "#000000", 0.86)         # canvas hue ~14% (deep, hue-true)
    t = max(0.0, min(1.0, (L - 0.45) / 0.5))         # 0 at L0.45 → 1 at L0.95
    dark_alpha = int(round(34 + 24 * t))             # brighter canvas → deeper (34..58)
    return {"dark": dark, "dark_alpha": dark_alpha,
            "light": "#FFFFFF", "light_alpha": 78, "blur": 22, "dist": 16}


def _neu_recipe():
    """Active neumorphic recipe = derive_neumorphic(canvas) overridden by any MEASURED
    values the look carries in grammar.neumorphic (a dict). So absorption is
    measurement-driven; a bare canvas still gets a correct derived recipe."""
    r = derive_neumorphic(_resolve("canvas"))
    g = GRAMMAR.get("neumorphic")
    if isinstance(g, dict):
        r.update({k: v for k, v in g.items() if k in r})
    return r


def _neumorphic_card(slide, box, rad, active, accent):
    """Draw a soft-UI card as STACKED single-shadow shapes (effectDag drops the fill
    in COM render, so a pair is built by stacking two one-shadow shapes — the showeet
    source's own method, fill-safe for any color). Recipe (color/alpha/blur/dist) comes
    from _neu_recipe() (canvas-derived + measured overrides) — NO hand-tuned constants.
    Light from the TOP-LEFT.

      DEFAULT = a CONCAVE depression carved into the SAME surface. Crucially NO outer
                shadow and NO border — so the card NEVER separates from the canvas; it
                reads as a recess in one continuous surface. The dent is built purely
                from OPPOSED-DIAGONAL inner shading: dark on the near NW edge + light
                on the far SE edge. (The canonical neumorphic concave.)
      ACTIVE  = a convex tile filled in the ACCENT that POPS OUT (outer pair).

    Two inner shadows can't share one effectLst (it holds a single innerShdw) and
    effectDag drops the fill, so the light SE rides a stacked TRANSPARENT overlay over
    the canvas shape. Verified by COM render: transparent + LIGHT inner renders, but
    transparent + DARK inner does NOT — dark MUST sit on the filled shape (not swappable)."""
    r = _neu_recipe()
    dark, light = r["dark"], r["light"]
    da, la, blur, dist = r["dark_alpha"], r["light_alpha"], r["blur"], r["dist"]
    fill = accent if (active and accent is not None) else "canvas"
    if active:
        # selected cell POPS OUT: convex raised tile (opposed OUTER pair), accent fill.
        base = rect(slide, box, fill=fill, radius=rad)
        neumorphic_shadow(base, "outer", 45, dark, da, blur, dist)      # dark cast SE
        top = rect(slide, box, fill=fill, radius=rad)
        neumorphic_shadow(top, "outer", 225, light, la, blur, dist)     # light highlight NW
        return top
    # DEFAULT cell: seamless concave — no outer, no border. Opposed-diagonal inner pair.
    base = rect(slide, box, fill="canvas", radius=rad)
    neumorphic_shadow(base, "inner", 225, dark, da, blur, dist)         # dark NEAR (NW) wall
    top = rect(slide, box, fill=None, radius=rad)                       # transparent overlay
    neumorphic_shadow(top, "inner", 45, light, la, blur, dist)          # light FAR (SE) wall
    return top


# --- mono-meta primitive (grammar.mono_meta) -----------------------------------
# The #1 identity cue of the dark-tech pack: a MONOSPACE META-LABEL LAYER —
# eyebrows/kickers (`// SECTION 02`), small labels (`[ DATA ]`), captions
# (`q1_results.log`), list/flow index numbers (`// 01`), and chart value/axis labels
# all render in an etched terminal face. This is a label SKIN, never a content
# rewrite: headlines, body copy, and data VALUES keep their proportional face and
# their exact text. Every seam below is a hard no-op unless the active look declares
# grammar.mono_meta (the other 109 looks are byte-identical).
#
# grammar.mono_meta (dict, all optional):
#   tiers     — meta tiers to mono-route (default _META_TIERS = eyebrow/label/caption)
#   prefix    — glyph prepended to eyebrow/kicker labels (e.g. "// "); default "// "
#   index_fmt — format for list/flow index numbers, must contain {n} (e.g. "// {n}")
#   color     — overrides the eyebrow/kicker color token (e.g. "blue" = neon cyan)
# A truthy bare `true` is treated as {} → all defaults.
def _mono_cfg():
    """Active look's mono_meta grammar as a normalized dict, or None when the look
    declares none → every mono-meta seam is a no-op (byte-identical regression)."""
    m = GRAMMAR.get("mono_meta")
    if not m:
        return None
    return {} if m is True else m


def _is_mono_tier(style_key):
    """True when the active look mono-routes this typography tier. The tier set is
    configurable so a look could, say, mono only captions — default = the eyebrow/
    label/caption meta layer. No mono_meta grammar → always False (no-op)."""
    cfg = _mono_cfg()
    if cfg is None:
        return False
    return style_key in tuple(cfg.get("tiers", _META_TIERS))


def _mono_face(is_disp):
    """Return (latin, ea) faces for a mono-routed run: the monospace face for BOTH
    so Latin labels AND any Korean glyphs render in the terminal face. (Most meta
    labels are Latin/numeric; Korean falls back via the OOXML ea slot — JetBrains
    Mono carries no Hangul, so PowerPoint substitutes a CJK face, which is the
    honest limit. Display tiers are never mono-routed so is_disp is just for parity.)"""
    mono = ACTIVE_FONT.get("mono") or ACTIVE_FONT["latin"]
    return mono, mono


def _mono_kicker_prefix(content):
    """Prepend the look's kicker glyph (default `// `) to an eyebrow/kicker string,
    idempotently (never double-prefixes). Pure label decoration — the words are
    unchanged, exactly like brutalism uppercasing a label. Returns the string
    untouched when no mono_meta grammar is active."""
    cfg = _mono_cfg()
    if cfg is None:
        return content
    pre = cfg.get("prefix", "// ")
    if not pre or not isinstance(content, str):
        return content
    return content if content.lstrip().startswith(pre.strip()) else pre + content


def _mono_index(n):
    """Format a 1-based sequence number as the look's mono index glyph (default
    `// 01`). Used by list/flow/agenda index badges so the terminal numbering reads
    as a code comment. No mono_meta grammar → plain zero-padded number (unchanged)."""
    cfg = _mono_cfg()
    if cfg is None:
        return f"{n:02d}"
    fmt = cfg.get("index_fmt", "// {n}")
    try:
        return fmt.format(n=f"{n:02d}")
    except Exception:
        return f"{n:02d}"


def _mono_kicker_color(default):
    """The eyebrow/kicker (and section index) color. A look can retint it with
    grammar.meta_color — used by a palette that has no 'blue' (the engine's default
    chrome accent) so the kicker stays IN palette instead of leaking an off-palette
    blue (bold-block-infographic: teal kicker, not blue). A mono_meta look uses its
    grammar.mono_meta.color. No grammar → the caller's default. byte-identical for
    every look that declares neither."""
    mc = GRAMMAR.get("meta_color")
    if mc:
        return mc
    cfg = _mono_cfg()
    if cfg is None:
        return default
    return cfg.get("color", default)


def _rotate_about(sp, deg, cx, cy):
    """Rotate a shape by `deg` (clockwise, PowerPoint convention) about an EXTERNAL
    pivot (cx, cy in EMU). Shapes rotated about the same pivot by the same angle
    stay rigid relative to each other — this is how a "card" (separate panel +
    accent bar + text boxes, all axis-aligned) is tilted as one sticker-like unit
    even though the engine has no grouping. A shape pivoted about its own center
    just spins in place (dx=dy=0)."""
    if not deg:
        return sp
    ccx = sp.left + sp.width / 2.0
    ccy = sp.top + sp.height / 2.0
    th = math.radians(deg)
    dx, dy = ccx - cx, ccy - cy
    nx = cx + dx * math.cos(th) - dy * math.sin(th)
    ny = cy + dx * math.sin(th) + dy * math.cos(th)
    sp.left = int(nx - sp.width / 2.0)
    sp.top = int(ny - sp.height / 2.0)
    sp.rotation = (sp.rotation or 0) + deg
    return sp


# Deterministic per-card tilt. A look opts in via components.card.tilt_deg (max
# degrees); each card gets an organic-but-fixed angle from this sequence keyed by
# draw order (no RNG → byte-stable). Absent/0 → no tilt → all 110 looks unchanged.
_TILT_SEQ = (-1.0, 0.75, -0.55, 1.0, -0.85, 0.6, -1.0, 0.5)
_CARD_SEQ = 0


def _next_tilt():
    """Next deterministic tilt angle (deg) for a card, or 0 when the active look
    declares no tilt."""
    global _CARD_SEQ
    t = (COMPONENTS.get("card") or {}).get("tilt_deg", 0) or 0
    if not t:
        return 0
    ang = t * _TILT_SEQ[_CARD_SEQ % len(_TILT_SEQ)]
    _CARD_SEQ += 1
    return ang


# ---------------------------------------------------------------------------
# Tilt decision = SEMANTIC, not stylistic.
# An element may tilt ONLY when its position/size carries no data meaning (a free
# block: cover title, Set cards, flow nodes, KPI cards). When geometry ENCODES a
# value — a bar's length, a point's coordinate, a matrix quadrant, a gauge angle —
# tilting would misread the data, so those archetypes stay axis-aligned. The pack's
# own prompt.md follows exactly this ("거친 미감은 회전·배치로, 막대 도형 자체는 정밀하게").
# GEOMETRY_LOCKED documents the data-bearing archetypes; they simply never call
# _tilt_group. (Among free blocks, a TOC/agenda still stays aligned for scanability
# — that secondary call is judgment, not encoded here.)
GEOMETRY_LOCKED = frozenset({
    "bars", "trend", "spread", "correlate", "share",      # charts
    "quadrant", "map", "matrix",                          # coordinate fields
    "gauge", "gantt", "timeline", "waterfall", "slope",   # value-positioned marks
    "tam_sam_som", "bullet",
})


def _tilt_group(slide, n0, cx_px, cy_px):
    """Rotate every shape added since index n0 about the pixel pivot (cx,cy) by the
    next deterministic tilt angle — the seam that tilts a 'free block' (panel +
    badge + text, drawn as separate axis-aligned shapes) as one rigid sticker.
    No-op when the active look declares no card.tilt_deg, so non-tilt looks are
    byte-identical. Free-block builders call this; GEOMETRY_LOCKED ones never do."""
    tilt = _next_tilt()
    if not tilt:
        return
    pcx, pcy = int(IN(cx_px)), int(IN(cy_px))
    for sp in list(slide.shapes)[n0:]:
        _rotate_about(sp, tilt, pcx, pcy)


def _accent_edge(slide, box, color, edge):
    """Draw a flat accent bar on one side of a card. The bar's SIDE+width come from
    the active style's `accent_edge` preset, so the *placement* of emphasis is a
    style property (top bar vs left rule) while the *color* stays the palette's
    accent role. Default = top/8px (reproduces the legacy hardcoded column bar)."""
    x, y, w, h = box
    t = edge.get("width_px", 8)
    side = edge.get("side", "top")
    bar = {"top": [x, y, w, t],
           "bottom": [x, y + h - t, w, t],
           "left": [x, y, t, h],
           "right": [x + w - t, y, t, h]}.get(side, [x, y, w, t])
    rect(slide, bar, fill=color, radius=0)


def card(slide, box, preset="card", fill=None, line_color=None, line_w=None,
         accent=None, active=False):
    """Draw a styled container from the active component preset (radius + border +
    shadow as one bundle). Falls back to a plain rect when the preset is absent,
    so this is safe even with a flat 'house' style. Per-call fill/line override
    the preset. This is the seam where the absorbed STYLE axis reaches the canvas.

    When `accent` (a color token) is given, an accent edge is drawn whose side+width
    come from the style preset's `accent_edge` (default top/8px). This lets an
    absorbed look move emphasis from a top bar to a left rule without touching
    layout code — the consulting (design-pick) look uses left/4px.

    A look whose card preset declares `elevation:"neumorphic"` takes a separate
    soft-UI BEZEL path (`_neumorphic_card`): a DEFAULT cell is a raised canvas frame
    (convex) holding a concave WELL — content recedes into the surface; an ACTIVE cell
    (`active=True`) POPS OUT as a convex card in the accent color. (Built by stacking
    single-shadow shapes — effectDag would drop the fill in COM render.)"""
    c = COMPONENTS.get(preset, {})
    if c.get("elevation") == "neumorphic":
        rad = c.get("radius", 0)
        rad = RAD.get(rad, rad) if isinstance(rad, str) else (rad or 0)
        # recipe (color/alpha/blur/dist) is canvas-derived + measured overrides inside
        # _neumorphic_card — NOT passed from the preset (no per-look hand-tuning).
        return _neumorphic_card(slide, box, rad, active, accent)
    rad = c.get("radius", 0)
    rad = RAD.get(rad, rad) if isinstance(rad, str) else (rad or 0)
    border = c.get("border") or {}
    lc = line_color if line_color is not None else border.get("color")
    lw = line_w if line_w is not None else border.get("width_pt", 0)
    fillv = fill if fill is not None else c.get("fill")
    # 다크 캔버스 카드 대비 보정: 룩이 카드 채움을 밝게(흡수 시 slate-50 등) 선언했어도
    # 다크 캔버스에선 surface(어둡게 도출)로 자동 리맵 → 밝은 카드 + 밝은(반전) 글자
    # 대비붕괴 방지. per-call fill override는 존중. 라이트 캔버스·이미 어두운 fill·무채움은
    # 무동작(무회귀 — 기존 다크룩은 surface가 이미 어두워 리맵 미발동).
    if DARK and fill is None and fillv is not None:
        try:
            if _rel_lum(_resolve(fillv)) > 0.7:
                fillv = "surface"
        except Exception:
            pass
    # wobble 문법: 손그림 룩 — 직선 카드 대신 삐뚤빼뚤 freeform(custGeom) 윤곽으로 교체.
    # 채움=카드색, 스트로크=ink(accent 카드는 accent색). 무선언 룩=미발동(byte-identical).
    _wob = _wobble_cfg()
    if _wob is not None and fillv is not None:
        return _wobble_card(slide, box, fillv,
                            accent if accent is not None else (lc or "ink"),
                            lw, c, _wob)
    sp = rect(slide, box, fill=fillv, line_color=lc, line_w=lw or 0, radius=rad)
    # 글래스 문법: 패널 채움을 반투명(기본 55%)으로 → 배경 그라디언트가 은은히 비침.
    # 텍스트는 솔리드 유지(빌더가 ink/white로 칠하므로 가독 보존). PPTX엔 backdrop-blur가
    # 없어 '프로스티드 블러'는 구현 불가 — 반투명+하이라이트 보더까지가 정직한 한계.
    if GRAMMAR.get("glass") and fillv is not None:
        _apply_alpha(sp, GRAMMAR.get("glass_alpha", 55))
    # 프리즘 문법: 노드는 "칠한 카드가 아니라 빛으로 그린 윤곽". 채움을 거의 투명
    # (기본 20% 알파)으로 낮추고 보더를 시안→마젠타→앰버 그라디언트 스트로크로 교체
    # → 다이어그램 노드가 빛 윤곽선으로 읽힌다. 프리즘에선 *모든* 서피스 채움 카드가
    # 글래스 노드 — accent 가 지정돼도(콘텐츠 카드 등) 단색 엣지 대신 그라디언트 스트로크로
    # 통일한다(프리즘의 강조는 단색 엣지가 아니라 보더 글로우/그라디언트). 단 block_fill
    # 원색 채움 카드(의미적 솔리드 타일)는 제외. 프리즘+다크일 때만(무선언/라이트=무동작).
    _prism = _prism_cfg()
    _is_node = (_prism is not None and fillv is not None
                and _resolve(fillv) in (_resolve("surface-2"), _resolve("surface-3"),
                                        _resolve("surface"), _resolve("surface-4")))
    if _is_node:
        if not GRAMMAR.get("glass"):
            _apply_alpha(sp, _prism.get("node_alpha", 20))
        gradient_stroke(sp)
        accent = None   # 프리즘 노드는 단색 accent 엣지를 그리지 않는다(그라디언트 스트로크가 강조)
    # 일반 가독성 레이어 — 자동 카드 반투명 (_bg_atmosphere 결정론 도출).
    # 룩이 배경 depth(그라디언트/라디얼글로우/이미지)를 선언했고, 카드에 채움이 있고,
    # 다크 캔버스이면 → 배경이 가장자리로 비치는 반투명 패널로 자동 전환한다. busy(다색
    # 메시) = 흰 프로스트 저알파, calm-depth(단색·동색조 어두운 그라디언트/단광원) = 옅은
    # 다크글래스. 룩이 명시 override(glass/glass_alpha)를 했으면 위 분기가 이미 처리했으니
    # 자동은 건너뛴다(override 우선). 프리즘 노드도 자기 경로가 처리하므로 제외(이중적용 방지).
    # depth 토큰이 없는 룩(brutalism·dark-tech 솔리드 카드 등)은 _bg_atmosphere depth=False
    # → 완전 미발동(슬라이드 XML byte-identical). 1차 가드 = depth 토큰 부재.
    elif (fillv is not None and DARK and not GRAMMAR.get("glass")
          and _bg_atmosphere()["depth"]):
        if _bg_atmosphere()["busy"]:
            _apply_frost(sp, _FROST_TINT, _FROST_ALPHA)
        else:
            _apply_alpha(sp, _DARKGLASS_ALPHA)
    # 해치 문법: 손그림 룩 — 카드 솔리드 채움을 사선 빗금(pattFill)으로 교체(빗금=ink, 바탕=
    # 카드 채움). 라이트 손그림 전용이라 위 glass/prism/atmosphere 분기와 상호배타(그 룩들은
    # 미발동). 무선언 룩·프리즘 노드·무채움이면 무동작(byte-identical).
    _hatch = _hatch_cfg()
    if _hatch is not None and fillv is not None and not _is_node:
        hatch_fill(sp, prst=_hatch.get("prst", "ltUpDiag"),
                   fg=_hatch.get("fg", "ink"), bg=fillv, fg_alpha=_hatch.get("alpha"))
    sh = c.get("shadow")
    if sh:
        shadow(sp, sh.get("blur_pt", 12), sh.get("dist_pt", 4),
               sh.get("dir_deg", 90), sh.get("color", "#000000"),
               sh.get("alpha", 78))
    # 글로우 문법: 다크 캔버스에서 카드/슬랩이 같은 색조의 부드러운 외곽 헤일로를 얻는다.
    # 헤일로 색 우선순위: ① accent 엣지가 있으면 그 액센트(가장 비비드 — 카드의 강조색과
    # 일치) ② 없으면 채움색(block_fill 카드 = 원색 채움) ③ 둘 다 없으면 보더색 → ACCENTS[0].
    # 스킨의 dim divider 보더만 쓰면 헤일로가 흐릿하므로, 강조색을 우선해 네온감을 살린다.
    # 무선언 룩·라이트 룩에선 _glow_shape 가 무동작(무회귀).
    if accent is not None:
        glow_col = accent
    elif fillv is not None and _resolve(fillv) != _resolve("surface-2") \
            and _resolve(fillv) != _resolve("surface-3"):
        glow_col = fillv                       # block_fill 원색 채움 카드
    else:
        glow_col = lc if lc else (ACCENTS[0] if ACCENTS else None)
    # 프리즘 노드(거의 투명 채움)에는 카드 글로우를 걸지 않는다 — PowerPoint a:glow 는
    # 반투명 도형의 내부까지 색을 번지게 해 인테리어를 시안으로 채워버린다(검증됨). 그라디언트
    # 스트로크 자체가 발광 윤곽 역할을 하므로 글로우 불필요. 글로우는 불투명/작은 마크
    # (배지·커넥터·차트)에만 남겨 깨끗한 헤일로로 읽히게 한다.
    # 글래스 카드(명시 glass=반투명 패널)도 같은 이유로 카드 글로우 제외 — 안 그러면 accent
    # 글로우가 반투명 카드 내부를 통째로 그 색으로 채워 "투명 글래스"가 "솔리드 색타일"로
    # 보인다(vivid 그리드 카드 버그). accent는 얇은 상단 엣지로만 남는다. glass 미선언 룩은
    # GRAMMAR.get("glass")=falsy → 조건 불변(byte-identical), glassmorphism은 glow 미선언이라
    # 본디 _glow_shape no-op → 무변화.
    if not _is_node and not GRAMMAR.get("glass"):
        _glow_shape(sp, "cards", color=glow_col)
    if accent is not None:
        _accent_edge(slide, box, accent, c.get("accent_edge") or {"side": "top", "width_px": 8})
    return sp


def line(slide, x1, y1, x2, y2, color="slate-300", w=1, glow_line=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    IN(x1), IN(y1), IN(x2), IN(y2))
    cn.line.color.rgb = C(color)
    cn.line.width = Pt(w)
    # 프리즘 커넥터: glow_line 컨텍스트(플로우·시스템 노드를 잇는 발광 커넥터)에서
    # 프리즘 룩이면 솔리드 색을 시안→마젠타→앰버 그라디언트 스트로크로 교체 → "좌→우
    # 색 전이 발광 라인"(스펙). 그라디언트 각은 라인의 좌→우 방향(0)으로 고정. 축·격자·
    # 베이스룰(glow_line=False)은 손대지 않는다(차트 정밀도 보존).
    if glow_line and _prism_cfg() is not None:
        gradient_stroke(cn, width_pt=max(w, 1.75))
    # 발광 라인: 커넥터만 opt-in(축·디바이더·차트 베이스룰은 glow_line=False 유지 →
    # 깨끗). glow 문법 선언 + 다크 + lines 활성일 때만 헤일로(네온 커넥터).
    if glow_line:
        _glow_shape(cn, "lines", color=_resolve(color))
    return cn


def _arrowhead(slide, cx, cy, size, color, deg=90):
    """Solid angular triangle centered at (cx,cy), pointing right by default
    (deg=90 rotates the base isoceles triangle apex from up→right). Brutalism's
    connector arrowheads ("크고 각진 삼각 화살촉") — no curves, flat fill."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                IN(cx - size / 2), IN(cy - size / 2),
                                IN(size), IN(size))
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = C(color)
    sp.line.fill.background()
    sp.rotation = deg
    return sp


def oval(slide, box, fill=None, line_color=None, line_w=0):
    x, y, w, h = box
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(x), IN(y), IN(w), IN(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line_color:
        sp.line.color.rgb = C(line_color)
        sp.line.width = Pt(line_w or 1)
    else:
        sp.line.fill.background()
    return sp


def shape_text(sp, content, style_key, color="white", align="center",
               anchor="middle", upper=None, wrap=True):
    tf = sp.text_frame
    tf.word_wrap = wrap
    # No-wrap chips (e.g. number badges) must keep the glyphs on one line: the
    # default ~0.1" L/R text margins shrink a small box's usable width enough to
    # wrap a 2-digit number like "02" onto two lines. Zero them so the run stays
    # single-line and optically centered. (bug E)
    if not wrap:
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    pt, bold, lh, up = style(style_key)
    if upper is not None:
        up = upper
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(content).upper() if up else str(content)
    r.font.size = pt
    r.font.bold = bold
    r.font.color.rgb = C(color)
    # mono-meta: a shape-anchored meta label (e.g. a chart legend chip's name) gets
    # the monospace face too, so the terminal label layer is consistent. No-op
    # otherwise (no grammar → _is_mono_tier False → default faces).
    if _is_mono_tier(style_key):
        mf, _ = _mono_face(False)
        _set_run_font(r, mf, mf)
    else:
        _set_run_font(r)
    _apply_tracking(r, style_key, _size_px(style_key))


def _num_badge(slide, box, txt, tier="h4", oval_fill="blue", oval_txt="on-accent"):
    """Sequence/rank number badge. Brutalism grammar (`badge: square`) → black
    square + white tabular number; default → accent oval. One seam for list/rank
    (agenda/flow inline their own size/format variants). No grammar → unchanged."""
    # mono-meta grammar: the index is a TERMINAL code-comment number (`// 01`), not
    # a filled chip — drawn as a bare mono label, left-aligned & middle-anchored, in
    # neon cyan. The numeric VALUE is preserved (only the `// ` glyph + mono face are
    # added). No mono_meta → falls through to the normal chip badges below.
    if _mono_cfg() is not None and str(txt).strip().isdigit():
        text(slide, box, _mono_index(int(str(txt).strip())), "label",
             color=_mono_kicker_color("blue"), align="left", anchor="middle",
             upper=True)
        return None
    if GRAMMAR.get("badge") == "square":
        b = rect(slide, box, fill="navy", radius=0)
        shape_text(b, txt, tier, color="white", wrap=False)
    elif GRAMMAR.get("badge") == "bare":
        # 럭셔리 에디토리얼 세리프: 칩(원/사각) 없음 — 가는 세리프 숫자만, 잉크 차콜.
        # "노드는 박스가 아니라 헤어라인 룰과 활자의 정렬로 암시한다"(prompt.md).
        text(slide, box, str(txt), tier, color="ink", align="center", anchor="middle")
        return None
    else:
        b = oval(slide, box, fill=oval_fill)
        shape_text(b, txt, tier, color=oval_txt)
        # 글로우: 액센트 원형 배지가 채움색으로 발광(네온 칩). 무선언/라이트=무동작.
        _glow_shape(b, "badges", color=_resolve(oval_fill))
    return b


_SLIDE_SEQ = 0       # deterministic slide counter for per-slide effects (radial corner cycle)
_CARD_ENUM = 0       # per-slide card enumerator (hyundai giant cyan 01/02/… numbers); reset in new_slide


def _bg_panel(slide):
    """Large inset 'sheet' panel (grammar.bg_panel) — the blueprint deck's framed drawing
    sheet that sits on EVERY slide: a big rounded rectangle inset from the edges with a
    subtle fill + a thin technical border; the grid is drawn OVER it so the panel reads as
    a sheet on a darker mat. No-op unless grammar.bg_panel is declared (byte-identical).

    grammar.bg_panel (dict, all optional):
      inset      — px inset from each edge (default 40)
      radius     — corner radius px (default 18)
      fill       — panel fill token/hex (default None = no fill, border-only frame)
      fill_alpha — fill opacity 0..100 (default opaque)
      border     — frame line token/hex (default 'blue')
      border_pt  — frame line weight pt (default 1.0)"""
    cfg = GRAMMAR.get("bg_panel")
    if not cfg:
        return
    if cfg is True:
        cfg = {}
    inset = cfg.get("inset", 40)
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    fillv = cfg.get("fill")
    sp = rect(slide, [inset, inset, W - 2 * inset, H - 2 * inset],
              fill=(_resolve(fillv) if fillv else None),
              line_color=cfg.get("border", "blue"),
              line_w=cfg.get("border_pt", 1.0), radius=cfg.get("radius", 18))
    if fillv and cfg.get("fill_alpha") is not None:
        _apply_alpha(sp, int(cfg["fill_alpha"]))
    return sp


def slide_bg(slide, color):
    # 메시 배경 룩(글래스모피즘·비비드) → 풀블리드 4정점 베이크 메시(무선언 무회귀). bg_gradient
    # 보다 우선 — 같은 정점색을 1D 선형이 아니라 2D 색면으로 깐다. 무선언이면 아래 분기로.
    m = GRAMMAR.get("bg_mesh")
    if m:
        _mesh_bg(slide, m)
        _bg_panel(slide)
        _grid_bg(slide)
        _radial_glow(slide, _SLIDE_SEQ)
        return
    # 글래스모피즘 등 grammar.bg_gradient 선언 룩 → 풀블리드 다색 그라디언트(무선언 무회귀).
    g = GRAMMAR.get("bg_gradient")
    if g:
        _gradient_bg(slide, [_resolve(c) for c in g])
        _bg_panel(slide)
        _grid_bg(slide)
        _radial_glow(slide, _SLIDE_SEQ)
        return
    # Skip when the (resolved) color is white — a white canvas needs no fill.
    # Resolve first so the role token "canvas" (=white in light mode) is a no-op.
    if _resolve(color).upper() in ("#FFFFFF", "WHITE"):
        return
    rect(slide, [0, 0, SLIDE["width_px"], SLIDE["height_px"]], fill=color)
    _bg_panel(slide)
    _grid_bg(slide)
    # Engineered-dark corner radial bloom (no-op for every other look). Drawn right on
    # the canvas so content sits above it. Corner cycles by the slide counter so the
    # bloom isn't pinned to one corner deck-wide. cover/section repaint via slide_bg →
    # the bloom redraws on the same corner over their bg (one bloom, same spot).
    _radial_glow(slide, _SLIDE_SEQ)


def _grid_bg(slide):
    """Faint micro-grid overlay — the dark-tech 'cyber depth' background. Draws 1px
    hairlines on a fixed pitch in a dim near-canvas color so the charcoal reads as
    graph paper without competing with content. Gated HARD: only when the active
    look declares grammar.grid_bg AND the canvas is dark (a grid on a light deck
    would be noise). Every other look → no shapes added (byte-identical).

    grammar.grid_bg (dict, all optional):
      pitch — line spacing in px (default 48 ≈ 0.5in at 96dpi, per the pack)
      color — hairline color token (default 'surface-2' = the code-block #16181D)
      width — line weight in pt (default 0.5)"""
    cfg = GRAMMAR.get("grid_bg")
    if not cfg or not DARK:
        return
    if cfg is True:
        cfg = {}
    pitch = cfg.get("pitch", 48)
    col = cfg.get("color", "surface-2")
    wpt = cfg.get("width", 0.5)
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    x = pitch
    while x < W:
        line(slide, x, 0, x, H, color=col, w=wpt)
        x += pitch
    y = pitch
    while y < H:
        line(slide, 0, y, W, y, color=col, w=wpt)
        y += pitch


# --- chevron motif primitive (grammar.chevron) ---------------------------------
# The #1 identity cue of the HYUNDAI cinematic investor-day pack — and the one mark
# that none of the other four dark looks own (dark-tech=grid, prismatic=prism stroke,
# engineered=corner radial, vivid=mesh glass). The deck's whole visual signature is a
# NESTED RIGHT-POINTING CHEVRON (`>` stack) with a bright white-cyan "eye" glow at the
# inner V-vertex, fading out through medium blue to the navy canvas. It appears in four
# roles, all gated on grammar.chevron + DARK (hard no-op for every other look):
#   • cover/closing  → big LEFT 5-stack + eye (bookend symmetry, prompt.md p-01/p-71)
#   • section        → RIGHT 4-stack + small left starburst (chapter slate, p-02/p-08)
#   • Mode-B header   → tiny 3-stack header mark + cyan eyebrow label (every body page)
#   • chart/table top → chevron-TAB stage band (trapezoidal CHEVRON tabs, p-48/p-52)
# A chevron RING is drawn as two thick connector legs meeting at the apex (a clean `>`
# outline, no block fill), in a cyan→blue stroke; the eye is an oval with a strong glow.
# This is pure decoration (carries no data) → never collides with GEOMETRY_LOCKED marks.
def _chevron_cfg():
    """Active look's chevron grammar as a normalized dict, or None when the look
    declares none / the canvas is light → every chevron seam is a no-op."""
    c = GRAMMAR.get("chevron")
    if not c or not DARK:
        return None
    return {} if c is True else c


# cyan-blue family stops for the chevron rings: bright core → header blue → fade.
# Resolved through the palette so a swap can't leak an off-family hue.
_CHEV_CORE = "#F5FBFF"      # inner V-vertex white-cyan peak (prompt.md chevron glow peak)
_CHEV_MID = "#00B0F0"       # medium blue (header label hue)
_CHEV_FAR = "#0060A8"       # outer ring → fades toward the navy canvas


def _chevron_ring(slide, apex_x, apex_y, half_h, depth, w_pt, color, glow_rad=0):
    """One `>` outline: an upper leg (apex → top-left) + a lower leg (apex → bottom-
    left), mirrored about the horizontal through the apex. half_h = vertical reach of
    each leg (px); depth = how far left the open ends sit (px). Right-pointing vertex
    at (apex_x, apex_y). Returns the two connectors. Optional halo for the bright rings."""
    top = (apex_x - depth, apex_y - half_h)
    bot = (apex_x - depth, apex_y + half_h)
    l1 = line(slide, apex_x, apex_y, top[0], top[1], color=color, w=w_pt)
    l2 = line(slide, apex_x, apex_y, bot[0], bot[1], color=color, w=w_pt)
    if glow_rad:
        for ln in (l1, l2):
            glow(ln, rad_px=glow_rad, color=_resolve(color), alpha=70)
    return l1, l2


def _chevron_stack(slide, apex_x, apex_y, reach_h, n=5, gap=None, w_pt=9.0,
                   eye=True):
    """A nested set of `n` right-pointing chevron rings sharing the apex, growing
    outward, in the cyan→blue family (inner rings bright/glowing, outer rings dim and
    fading to navy). With `eye`, a bright white-cyan oval + strong halo sits at the
    apex — the deck's signature 'eye' at the V-vertex. reach_h = the OUTERMOST ring's
    half-height (px); rings step inward by `gap`."""
    gap = gap or reach_h / (n + 0.6)
    # outer (dim, far) → inner (bright, glowing); draw far first so bright sits on top.
    # The ramp keeps EVERY ring legible (the pack reads as a full layered `>`, not a
    # lone bright stroke): inner rings = white-cyan core + halo, mid rings = bright
    # #00B0F0, outer rings = #66FFFF-leaning blue (still clearly cyan vs the navy bg).
    for i in range(n - 1, -1, -1):
        half = (i + 1) * gap
        depth = half * 0.88          # `>` aspect ≈ vertical:horizontal 1 : 0.88
        t = i / max(1, n - 1)        # 0 = innermost, 1 = outermost
        if t < 0.25:
            col, gr = _CHEV_CORE, 13
            wp = w_pt
        elif t < 0.55:
            col, gr = "#5AD8FF", 9
            wp = w_pt * 0.92
        else:
            col, gr = _CHEV_MID, 5
            wp = w_pt * 0.85
        _chevron_ring(slide, apex_x, apex_y, half, depth, wp, col, glow_rad=gr)
    if eye:
        r = max(10.0, reach_h * 0.07)
        ov = oval(slide, [apex_x - r, apex_y - r, 2 * r, 2 * r], fill=_CHEV_CORE)
        glow(ov, rad_px=int(r * 1.9), color=_CHEV_CORE, alpha=85)
        # a second, larger soft cyan bloom around the eye for the cinematic flare
        r2 = r * 0.55
        ov2 = oval(slide, [apex_x - r2, apex_y - r2, 2 * r2, 2 * r2], fill="#66FFFF")
        glow(ov2, rad_px=int(r * 1.3), color="#66FFFF", alpha=80)


def _starburst(slide, cx, cy, size):
    """A small white-cyan starburst glow — the chevron's reciprocal accent that
    punctuates a multi-word section title (prompt.md p-08). A bright dot + halo +
    two faint crossing rays."""
    r = size * 0.16
    dot = oval(slide, [cx - r, cy - r, 2 * r, 2 * r], fill=_CHEV_CORE)
    glow(dot, rad_px=int(size * 0.5), color="#66FFFF", alpha=80)
    line(slide, cx - size / 2, cy, cx + size / 2, cy, color="#66FFFF", w=1.0)
    line(slide, cx, cy - size / 2, cx, cy + size / 2, color="#66FFFF", w=1.0)


def _chevron_cover(slide, side="left"):
    """Big 5-stack chevron on one edge (cover/closing bookend). Left edge = right-
    pointing stack whose apex sits ~32% in, reaching ~80% canvas height. No-op
    unless grammar.chevron is active."""
    if _chevron_cfg() is None:
        return
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    if side == "left":
        apex_x = W * 0.30
    else:
        apex_x = W * 0.70
    apex_y = H * 0.5
    _chevron_stack(slide, apex_x, apex_y, reach_h=H * 0.40, n=5, w_pt=11.0,
                   eye=_chevron_cfg().get("eye_glow", True))


def _chevron_section(slide, multiword=False):
    """Section chapter motif: a RIGHT-edge 4-stack chevron (+ a small left starburst
    for multi-word titles). No eye (the cover owns the hero eye); rings still glow."""
    if _chevron_cfg() is None:
        return
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    apex_x = W * 0.86
    apex_y = H * 0.52
    _chevron_stack(slide, apex_x, apex_y, reach_h=H * 0.30, n=4, w_pt=9.0, eye=False)
    if multiword:
        # reciprocal accent: a small starburst just LEFT of and above the title's
        # baseline (the section_title slot starts at x≈96 / y≈427), punctuating it
        # without sitting on the glyphs.
        _starburst(slide, W * 0.035, H * 0.36, 86)


def _chevron_header_mark(slide):
    """The Mode-B permanent header mark: a tiny 3-stack chevron glow at top-left +
    the cyan '2024 CEO Investor Day'-style eyebrow label. Fixed position, every body
    page. The label text comes from the look's chevron.header_label (default a generic
    investor-day tag); drawn only when grammar.chevron.header_mark is set."""
    cfg = _chevron_cfg()
    if cfg is None or not cfg.get("header_mark", True):
        return
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    ax, ay = W * 0.058, H * 0.052
    _chevron_stack(slide, ax, ay, reach_h=16, n=3, w_pt=2.6, eye=False)
    label = cfg.get("header_label", "Investor Day")
    text(slide, [W * 0.072, H * 0.030, W * 0.30, H * 0.045], label, "caption",
         color="#00B0F0", anchor="middle", upper=True)


def _stage_band(slide, box, stages, active=None):
    """Chevron-TAB stage band (prompt.md primitive #4): trapezoidal CHEVRON tabs
    flowing left→right (e.g. 초기 > 성장기 > 성숙기 / Upstream > Midstream > Downstream).
    Only the active tab is cyan-filled (#66FFFF + dark text); the rest are surface fill
    + blue-grey text. box = [x,y,w,h]. Used by charts and the comparison table when
    grammar.chevron.stage_band is on. Returns nothing (draws in place)."""
    x, y, w, h = box
    n = max(1, len(stages))
    if active is None:
        active = n // 2
    gap = 6
    tab_w = (w - gap * (n - 1)) / n
    overlap = h * 0.34          # how far each chevron's point bites into the next
    for i, label in enumerate(stages):
        tx = x + i * (tab_w + gap)
        on = (i == active)
        sp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                    IN(tx), IN(y), IN(tab_w + overlap), IN(h))
        sp.shadow.inherit = False
        sp.fill.solid()
        sp.fill.fore_color.rgb = C("blue" if on else "surface-2")
        sp.line.fill.background()
        try:
            sp.adjustments[0] = min(0.5, overlap / (tab_w + overlap))
        except Exception:
            pass
        if on:
            glow(sp, rad_px=12, color="#66FFFF", alpha=55)
        shape_text(sp, str(label), "label",
                   color="#101030" if on else "muted", wrap=False)


def header(slide, data, dark=None):
    """헤더밴드: kicker(eyebrow) + title + divider. 좌표는 grid.json 파생
    (_meta.header_band). 슬롯명은 archetypes.md/grid.json 의 kicker/title/divider.
    data 는 'kicker' 또는 'eyebrow'(별칭) 둘 다 허용한다.
    dark=None 이면 활성 룩의 DARK 플래그를 따른다 → 다크 캔버스에서 제목·eyebrow·
    디바이더 색을 한 곳에서 자동 플립(17 헤더 아키타입 공용 시맨)."""
    if dark is None:
        dark = DARK
    hb = LAYOUTS["_meta"]["header_band"]
    ti = "ink"
    eb = _mono_kicker_color("blue" if not dark else "blue-light")
    kick = data.get("kicker") or data.get("eyebrow")
    # Hyundai permanent header mark (Mode B): a fixed tiny chevron glow + cyan
    # investor-day label sits in the eyebrow slot on EVERY body page. It IS the
    # eyebrow line for this look, so the page kicker is suppressed (the pack's top-
    # left is always the constant investor-day mark, never a per-page label). No-op
    # for every other look → page kicker renders as before.
    if _chevron_cfg() is not None and _chevron_cfg().get("header_mark", True):
        _chevron_header_mark(slide)
    elif kick:
        text(slide, hb["kicker"], _mono_kicker_prefix(kick), "eyebrow",
             color=eb, upper=True)
    if data.get("title"):
        text(slide, hb["title"], data["title"], "h2", color=ti)
    dv = hb["divider"]
    # Header rule: a hairline by default; a look can declare `grammar.header_rule`
    # (e.g. brutalism's thick black bar) so the divider matches the deck language
    # across all 17 header archetypes from one seam.
    hr = GRAMMAR.get("header_rule") or {}
    line(slide, dv[0], dv[1], dv[0] + dv[2], dv[1],
         color=hr.get("color", "divider"), w=hr.get("width_pt", 1.5))


def footer(slide, data):
    """선택적 푸터밴드: caption(좌) + pageno(우). 좌표는 grid.json 파생."""
    fb = LAYOUTS["_meta"]["footer_band"]
    if data.get("caption"):
        text(slide, fb["caption"], data["caption"], "caption", color="subtle")
    if data.get("pageno") is not None:
        text(slide, fb["pageno"], str(data["pageno"]), "caption",
             color="subtle", align="right")


def new_slide(prs):
    global _SLIDE_SEQ, _CARD_ENUM
    _SLIDE_SEQ = len(prs.slides._sldIdLst)   # 0-based index of the slide about to be added
    _CARD_ENUM = 0                           # restart the per-slide card enumerator
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Canvas seam: paint the content background ONCE here so every archetype
    # (statement·duo·trio·grid·bars·share·…) sits on the look's canvas without each
    # builder calling slide_bg. slide_bg no-ops on white, so normal light looks are
    # unchanged; a look declaring a bold canvas (dark navy OR bright brutalist
    # yellow) gets it painted. Frame archetypes draw their own bg over this.
    slide_bg(s, "canvas")
    return s


# ---------------------------------------------------------------------------
# Archetype builders (24 슬러그) — 좌표는 전부 layouts.json(grid.json 파생).
# 매직넘버 금지: 슬롯 박스에서 패딩만 상수로 빼서 파생한다.
# ---------------------------------------------------------------------------
PAD = 32          # 카드 내부 패딩(px) — 박스 안에서만 쓰는 로컬 상수
GAP = 24          # 카드/셀 사이 간격(px)


def _slots(slug):
    return LAYOUTS["archetypes"][slug]["slots"]


def _split_cols(box, n, gap=GAP):
    """box[x,y,w,h] 를 n 개 가로 컬럼으로 등분. 각 컬럼 박스 리스트 반환."""
    x, y, w, h = box
    if n <= 0:
        return []
    cw = (w - gap * (n - 1)) / n
    return [[x + i * (cw + gap), y, cw, h] for i in range(n)]


def _split_rows(box, n, gap=12):
    """box 를 n 개 가로 행으로 등분."""
    x, y, w, h = box
    if n <= 0:
        return []
    rh = (h - gap * (n - 1)) / n
    return [[x, y + i * (rh + gap), w, rh] for i in range(n)]


def _split_grid(box, ncol, nrow, gap=GAP):
    """box 를 ncol×nrow 셀로 분할. 행우선(row-major) 셀 박스 리스트."""
    x, y, w, h = box
    cw = (w - gap * (ncol - 1)) / ncol
    ch = (h - gap * (nrow - 1)) / nrow
    out = []
    for r in range(nrow):
        for c in range(ncol):
            out.append([x + c * (cw + gap), y + r * (ch + gap), cw, ch])
    return out


# ---------------------------------------------------------------------------
# 미디어 슬롯 헬퍼 — 이미지 또는 플레이스홀더.
# ---------------------------------------------------------------------------
def _resolve_src(src):
    """이미지 경로 해소. 절대경로면 그대로, 아니면 spec 디렉토리 기준.
    WSL 마운트 경로(/mnt/c/...) 권장. 존재하지 않으면 None."""
    if not src:
        return None
    if not os.path.isabs(src):
        src = os.path.join(SPEC_DIR or os.getcwd(), src)
    return src if os.path.exists(src) else None


def media(slide, box, spec, fit=None):
    """미디어 슬롯 채우기. spec={src,fit,caption,focal}.
    - src 있고 파일 존재 → add_picture, box 에 맞춤(cover=크롭 / contain=레터박스).
    - 없으면 → 플레이스홀더(slate-50 + 1px 테두리 + 중앙 라벨 + 비율힌트, 이모지 금지).
    - caption 있으면 박스 하단에 작은 캡션.
    좌표는 호출자가 layouts.json 슬롯에서 넘긴 box 그대로 사용(매직넘버 없음)."""
    x, y, w, h = box
    spec = spec or {}
    fit = fit or spec.get("fit", "cover")
    src = _resolve_src(spec.get("src"))
    cap = spec.get("caption")
    cap_h = 28 if cap else 0
    img_h = h - cap_h
    if src:
        try:
            _place_picture(slide, [x, y, w, img_h], src, fit, spec.get("focal", "center"))
        except Exception as e:
            _media_placeholder(slide, [x, y, w, img_h],
                               f"이미지 로드 실패: {os.path.basename(src)}")
    else:
        _media_placeholder(slide, [x, y, w, img_h], spec.get("label", "이미지 자리"),
                           ratio=f"{w:.0f}×{img_h:.0f}")
    if cap:
        text(slide, [x, y + h - cap_h, w, cap_h], cap, "caption",
             color="muted")


def _place_picture(slide, box, src, fit, focal):
    """실제 이미지 삽입. cover=박스 꽉 채우고 크롭, contain=비율유지 레터박스.
    PIL 로 원본 비율을 읽어 결정론적으로 크롭/레터박스 계산한다."""
    x, y, w, h = box
    try:
        from PIL import Image
        with Image.open(src) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = w, h
    box_ar = w / h if h else 1
    img_ar = iw / ih if ih else 1
    if fit == "contain":
        # 비율 유지 레터박스: 박스 안에 들어가도록 축소, 가운데 정렬
        if img_ar > box_ar:
            nw = w; nh = w / img_ar
        else:
            nh = h; nw = h * img_ar
        px = x + (w - nw) / 2
        py = y + (h - nh) / 2
        slide.shapes.add_picture(src, IN(px), IN(py), IN(nw), IN(nh))
    else:
        # cover: 박스를 꽉 채우고 넘치는 부분 크롭. add_picture 후 crop_* 로 잘라낸다.
        pic = slide.shapes.add_picture(src, IN(x), IN(y), IN(w), IN(h))
        if img_ar > box_ar:
            # 이미지가 더 넓음 → 좌우 크롭
            crop = (1 - box_ar / img_ar) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_ar < box_ar:
            # 이미지가 더 높음 → 상하 크롭(focal 로 미세조정)
            crop = (1 - img_ar / box_ar) / 2
            if focal == "top":
                pic.crop_top = 0; pic.crop_bottom = crop * 2
            elif focal == "bottom":
                pic.crop_top = crop * 2; pic.crop_bottom = 0
            else:
                pic.crop_top = crop; pic.crop_bottom = crop
        return pic


def _media_placeholder(slide, box, label="이미지 자리", ratio=None):
    """이미지 없을 때 자리표시. slate-50 채움 + 1px 테두리 + 중앙 라벨 + 단순 도형
    아이콘(이모지 금지: 산 모양 삼각형 + 해 원 으로 '사진' 픽토그램)."""
    x, y, w, h = box
    rect(slide, box, fill="surface-2", line_color="divider", line_w=1, radius=0)
    # 픽토그램(도형만): 작은 원(해) + 삼각형(산) — 박스 중앙 위쪽
    icon = min(w, h) * 0.18
    cx, cy = x + w / 2, y + h / 2 - icon * 0.4
    oval(slide, [cx - icon * 0.55, cy - icon * 0.7, icon * 0.5, icon * 0.5],
         fill="hairline")
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 IN(cx - icon * 0.6), IN(cy - icon * 0.2),
                                 IN(icon * 1.4), IN(icon * 0.9))
    tri.shadow.inherit = False
    tri.fill.solid(); tri.fill.fore_color.rgb = C("hairline")
    tri.line.fill.background()
    # 라벨 + 비율 힌트
    lbl = label + (f"\n{ratio}" if ratio else "")
    text(slide, [x, y + h / 2 + icon * 0.6, w, 50], lbl, "caption",
         color="subtle", align="center", anchor="top")


# ---------------------------------------------------------------------------
# 차트 헬퍼 — 약 20종 디스패치(native python-pptx + 도형 합성).
# ---------------------------------------------------------------------------
# Categorical series color = the deck's DECLARED accents first (ACCENTS roles),
# then the primary accent's own tint ladder as padding. So:
#   - mono palette (1 accent)  -> [blue, blue-light, blue-2, blue-pale]  (= old
#     behavior exactly: one coherent tint ladder, strong 2-series contrast)
#   - dual+ palette/look       -> [blue, orange, blue-light, …]  (the look's 2nd
#     accent differentiates series 2 — e.g. ppt-dark-tech mint vs violet)
# This honors a look's deliberate multi-accent identity (esp. line/slope charts
# where color is the ONLY differentiator) while never leaking OFF-palette tokens:
# only ACCENTS (palette-defined) + primary-accent tints are ever used — the old
# green/slate/navy bleed bug came from raw semantic tokens, not from this.
# C() resolves each token against the look-mutated COLORS. >ramp series cycle.
def chart_color(i):
    """i-th chart series color token (accent-first, tint-padded). See note above."""
    pad = ["blue-light", "blue-2", "blue-pale"]
    return ACCENTS[i] if i < len(ACCENTS) else pad[(i - len(ACCENTS)) % len(pad)]

_NATIVE = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column_stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "radar": XL_CHART_TYPE.RADAR,
}


def _theme_axes(ch):
    """Flip native-chart axis ink for canvas contrast. Tick labels -> ink,
    axis lines + gridlines -> divider. Resolves the active role tokens, so in
    light mode this reproduces the dark-on-white default and in dark mode it
    lifts labels/lines to readable values. python-pptx's high-level API exposes
    these, so no OOXML surgery needed. Best-effort per axis (some chart types
    lack a value/category axis)."""
    ink = C("ink")
    div = C("divider")
    for getter in ("category_axis", "value_axis"):
        try:
            ax = getattr(ch, getter)
        except Exception:
            continue
        try:
            ax.tick_labels.font.color.rgb = ink
        except Exception:
            pass
        try:
            # 축 제목(예: '건수')은 다크 캔버스에서 기본색이 어두워 묻힌다 → ink로.
            # 다크 게이트라 라이트는 미발동(무회귀 — 기존 baseline엔 축제목 색 XML 없음).
            if DARK and ax.has_title:
                for _p in ax.axis_title.text_frame.paragraphs:
                    _p.font.color.rgb = ink
                    for _r in _p.runs:
                        _r.font.color.rgb = ink
        except Exception:
            pass
        try:
            ax.format.line.color.rgb = div
        except Exception:
            pass
        try:
            if ax.has_major_gridlines:
                ax.major_gridlines.format.line.color.rgb = div
        except Exception:
            pass
    try:
        # 차트 제목도 다크에서 ink로(라이트 미발동).
        if DARK and ch.has_title:
            for _p in ch.chart_title.text_frame.paragraphs:
                _p.font.color.rgb = ink
                for _r in _p.runs:
                    _r.font.color.rgb = ink
    except Exception:
        pass


def _theme_legend(ch):
    """Legend text -> ink so series names stay readable on a dark canvas."""
    try:
        ch.legend.font.color.rgb = C("ink")
    except Exception:
        pass


_BRUTAL_BAR_TYPES = {"column_clustered", "column_stacked", "column_stacked_100",
                     "bar_clustered"}


def _chart_brutal_bars(slide, box, spec):
    """Neo-brutalist bar chart, drawn (not native), per the pack's chart grammar:
    every bar = flat primary fill + thick black border + hard offset shadow, value
    in a black square badge above the bar, category label directly below, a single
    4px black baseline (no gridlines, no legend, zero-based). One emphasis bar gets
    a different primary; the rest share one color. Skin (border weight/color, shadow)
    is INHERITED from the active look's card preset so the chart matches its cards.

    Renders the FIRST series only — brutalism's chart language is single-series with
    one highlighted bar ("전부 파랑, 핵심만 빨강"); multi-series clustering is not part
    of the family. Highlight index via spec['highlight'] (default: max value)."""
    x, y, w, h = box
    cats, vals = _series_values(spec)
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    n = len(vals)
    # Skin from the look's card preset (DRY — chart shares the deck's border/shadow).
    cc = COMPONENTS.get("card", {})
    bd = cc.get("border") or {}
    bcol = bd.get("color", "navy")
    bwid = bd.get("width_pt", 4)
    sh = cc.get("shadow")
    acc = ACCENTS[0] if ACCENTS else "blue"
    hot = (ACCENTS[2] if len(ACCENTS) > 2 else
           (ACCENTS[-1] if len(ACCENTS) > 1 else acc))
    hl = spec.get("highlight")
    if hl is None:
        hl = max(range(n), key=lambda i: vals[i])
    # Geometry: width:gap = 4:3, zero baseline, headroom for value badge + labels.
    TOP, BOT, BADGE_H = 52, 48, 34
    base_y = y + h - BOT
    plot_h = h - TOP - BOT
    u = w / (7 * n - 3) if n > 1 else w / 4.0
    barw = 4 * u
    gap = 3 * u
    vmax = max(vals) or 1.0
    # 4px black baseline rule (no gridlines).
    line(slide, x, base_y, x + w, base_y, color=bcol, w=3)
    for i, v in enumerate(vals):
        bx = x + i * (barw + gap)
        bh = max(2.0, (v / vmax) * plot_h)
        by = base_y - bh
        col = hot if i == hl else acc
        sp = rect(slide, [bx, by, barw, bh], fill=col, line_color=bcol,
                  line_w=bwid, radius=0)
        if sh:
            shadow(sp, sh.get("blur_pt", 0), sh.get("dist_pt", 6),
                   sh.get("dir_deg", 135), sh.get("color", "#0A0A0A"),
                   sh.get("alpha", 0))
        # value badge: black square above the bar, white tabular number
        bw_badge = min(barw, 96)
        try:
            vtxt = spec.get("value_fmt", "{:g}").format(v)
        except Exception:
            vtxt = str(v)
        badge = rect(slide, [bx + (barw - bw_badge) / 2, by - BADGE_H - 6,
                             bw_badge, BADGE_H], fill="navy", radius=0)
        shape_text(badge, vtxt, "label", color="white", wrap=False)
        # category label directly below the bar
        if i < len(cats):
            text(slide, [bx - gap / 2, base_y + 8, barw + gap, 32], str(cats[i]),
                 "label", color="ink", align="center", upper=True)


# --- brutal chart skin helpers (shared by the drawn-shape brutal renderers) ---
def _brutal_skin():
    """(border_color, border_width_pt, shadow_dict) from the active look's card
    preset — so brutal charts share the deck's exact border/shadow."""
    cc = COMPONENTS.get("card", {})
    bd = cc.get("border") or {}
    return bd.get("color", "navy"), bd.get("width_pt", 4), cc.get("shadow")


def _brutal_shadow(sp, sh):
    if sh:
        shadow(sp, sh.get("blur_pt", 0), sh.get("dist_pt", 6),
               sh.get("dir_deg", 135), sh.get("color", "#0A0A0A"),
               sh.get("alpha", 0))


def _brutal_palette(n):
    """n flat PRIMARY colors, cycling the look's accents (원색 충돌, no tints).
    Excludes any accent equal to the canvas — a canvas-colored data mark/segment
    would vanish into the background (brutalism's bright canvas IS one of its
    accents, e.g. yellow). Falls back to all accents if filtering empties it."""
    canvas = _resolve("canvas").upper()
    base = [a for a in (ACCENTS or ["blue"]) if _resolve(a).upper() != canvas] \
        or (ACCENTS or ["blue"])
    return [base[i % len(base)] for i in range(n)]


def _chart_brutal_line(slide, box, spec):
    """Neo-brutalist line/trend, drawn: thick black polyline (straight segments,
    no curve) + black-bordered flat-primary square markers at each vertex, 4px
    baseline, direct category labels, no gridlines/legend (multi-series → small
    black-bordered name badges). Per the pack's line rule."""
    x, y, w, h = box
    cats = spec.get("categories", [])
    series = spec.get("series") or []
    allv = [float(v) for s in series for v in s.get("values", [])]
    n = max((len(s.get("values", [])) for s in series), default=0)
    if not allv or n == 0:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    bcol, bwid, _sh = _brutal_skin()
    vmax = max(allv); vmin = min(0.0, min(allv)); rng = (vmax - vmin) or 1.0
    TOP, BOT = 56, 52
    base_y = y + h - BOT; plot_h = h - TOP - BOT
    left = x + 50; right = x + w - 30; span = right - left
    xs = [left + (span * (i / (n - 1)) if n > 1 else span / 2) for i in range(n)]
    yf = lambda v: base_y - (v - vmin) / rng * plot_h
    line(slide, left, base_y, right, base_y, color=bcol, w=3)   # 4px baseline
    cols = _brutal_palette(len(series))
    for si, s in enumerate(series):
        vs = [float(v) for v in s.get("values", [])]
        for i in range(len(vs) - 1):
            line(slide, xs[i], yf(vs[i]), xs[i + 1], yf(vs[i + 1]), color=bcol, w=4)
        for i, v in enumerate(vs):
            rect(slide, [xs[i] - 10, yf(v) - 10, 20, 20], fill=cols[si],
                 line_color=bcol, line_w=bwid, radius=0)
    for i in range(n):
        if i < len(cats):
            text(slide, [xs[i] - 50, base_y + 10, 100, 28], str(cats[i]),
                 "label", color="ink", align="center", upper=True)
    if len(series) > 1:
        bx0 = left
        for si, s in enumerate(series):
            nm = str(s.get("name", "") or f"S{si+1}")
            rect(slide, [bx0, y + 6, 24, 24], fill=cols[si], line_color=bcol,
                 line_w=bwid, radius=0)
            text(slide, [bx0 + 32, y + 4, 30 + len(nm) * 16, 28], nm, "label",
                 color="ink", upper=True)
            bx0 += 62 + len(nm) * 16


def _chart_brutal_scatter(slide, box, spec):
    """Neo-brutalist scatter: thick black L-axes + black-bordered flat-primary
    square markers (one emphasis marker via point.highlight). No smooth dots."""
    x, y, w, h = box
    pts = spec.get("points", [])
    if not pts:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    bcol, bwid, _sh = _brutal_skin()
    xv = [float(p.get("x", 0)) for p in pts]; yv = [float(p.get("y", 0)) for p in pts]
    xmn, xmx = min(xv), max(xv); ymn, ymx = min(yv), max(yv)
    xr = (xmx - xmn) or 1.0; yr = (ymx - ymn) or 1.0
    L = x + 50; R = x + w - 30; B = y + h - 52; T = y + 24
    line(slide, L, T, L, B, color=bcol, w=4)
    line(slide, L, B, R, B, color=bcol, w=4)
    base = ACCENTS[0] if ACCENTS else "blue"
    hot = ACCENTS[-1] if len(ACCENTS) > 1 else base
    for p in pts:
        px = L + (float(p.get("x", 0)) - xmn) / xr * (R - L)
        py = B - (float(p.get("y", 0)) - ymn) / yr * (B - T)
        col = hot if p.get("highlight") else base
        sz = 30
        rect(slide, [px - sz / 2, py - sz / 2, sz, sz], fill=col,
             line_color=bcol, line_w=bwid, radius=0)
        if p.get("label"):
            text(slide, [px + sz / 2 + 8, py - 15, 240, 30], str(p["label"]),
                 "label", color="ink", upper=True, anchor="middle")


def _chart_brutal_pie(slide, box, spec):
    """Neo-brutalist 'share' = a thick horizontal PROPORTION bar: flat-primary
    segments sized by value, a single hard shadow, black borders dividing the
    segments, value labels inside each (text color flips by segment luminance).
    A bold proportion bar is the deterministic, reliable brutalist take on share
    (PIE-wedge donuts render unreliably); the caller (build_share) still lists
    categories with matching square chips."""
    x, y, w, h = box
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    bcol, bwid, sh = _brutal_skin()
    cols = _brutal_palette(len(vals))
    total = sum(vals) or 1.0
    barh = min(200, h * 0.5)
    by = y + (h - barh) / 2
    # one hard shadow under the whole bar (per-segment shadows would look busy)
    base = rect(slide, [x, by, w, barh], fill=cols[0], line_color=bcol,
                line_w=bwid, radius=0)
    _brutal_shadow(base, sh)
    cx = x
    for i, v in enumerate(vals):
        seg = (v / total) * w
        rect(slide, [cx, by, seg, barh], fill=cols[i], line_color=bcol,
             line_w=bwid, radius=0)
        if seg > 56:   # only label segments wide enough to hold the number
            text(slide, [cx, by + barh / 2 - 20, seg, 40], f"{v:g}", "h4",
                 color=_ink_on(cols[i]), align="center", anchor="middle")
        cx += seg


def _chart_swiss_bars(slide, box, spec):
    """Swiss-editorial bars: flat ACHROMATIC (ink) bars + ONE spot-color emphasis
    bar, a single 3pt ink baseline rule (no gridlines/legend/border/shadow),
    uppercase direct value(위)+category(아래) labels. bar:gap = 3:2, zero baseline.
    맥킨지식 절제(강조 1개·나머지 무채색)를 스위스 그리드 정밀성으로."""
    x, y, w, h = box
    cats, vals = _series_values(spec)
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    n = len(vals)
    ink = "navy"                              # #111111 — 무채색 베이스
    spot = ACCENTS[0] if ACCENTS else "blue"  # 슬라이드 스팟 1색
    hl = spec.get("highlight")
    if hl is None:
        hl = max(range(n), key=lambda i: vals[i])
    TOP, BOT, LBL = 48, 44, 28
    base_y = y + h - BOT
    plot_h = h - TOP - BOT
    u = w / (5 * n - 2) if n > 1 else w / 3.0   # 막대:간격 = 3:2
    barw, gap = 3 * u, 2 * u
    vmax = max(vals) or 1.0
    for i, v in enumerate(vals):
        bx = x + i * (barw + gap)
        bh = max(2.0, (v / vmax) * plot_h)
        by = base_y - bh
        rect(slide, [bx, by, barw, bh], fill=(spot if i == hl else ink), radius=0)
        try:
            vtxt = spec.get("value_fmt", "{:g}").format(v)
        except Exception:
            vtxt = str(v)
        # 값 레이블 위 · 카테고리 아래 — label 티어(대문자+자간은 grammar가 적용)
        text(slide, [bx - gap / 2, by - LBL - 6, barw + gap, LBL], vtxt, "label",
             color="ink", align="center", upper=True)
        if i < len(cats):
            text(slide, [bx - gap / 2, base_y + 8, barw + gap, LBL], str(cats[i]),
                 "label", color="ink", align="center", upper=True)
    line(slide, x, base_y, x + w, base_y, color=ink, w=3)   # x축 = 3pt 잉크 룰


# --- swiss chart skin helpers (shared by the drawn swiss renderers) ---
def _swiss_skin():
    """(ink, spot) for swiss charts: achromatic ink (#111 via 'navy') base +
    the slide's single spot accent. Mirror of _brutal_skin but for the
    무채색+스팟1 language (no borders/shadows)."""
    return "navy", (ACCENTS[0] if ACCENTS else "blue")


def _swiss_emphasis(vals, spec):
    """Index of the ONE spot-color element (explicit spec['highlight'] or max).
    Swiss/McKinsey restraint = everything ink, one thing spot."""
    if not vals:
        return 0
    hl = spec.get("highlight")
    return hl if hl is not None else max(range(len(vals)), key=lambda i: vals[i])


def _swiss_share_colors(vals, spec):
    """Per-segment colors for a swiss share: ink everywhere, spot on the one
    emphasis segment. Reused by _chart_swiss_pie AND build_share's chips so the
    bar and its legend always agree."""
    ink, spot = _swiss_skin()
    hl = _swiss_emphasis([float(v) for v in (vals or [])], spec)
    return [spot if i == hl else ink for i in range(len(vals or []))]


def _chart_swiss_line(slide, box, spec):
    """Swiss line/trend: thin (2.5pt) achromatic polylines + ONE spot series
    (emphasis = spec['highlight'] series, else the last; single series = spot),
    small precise square vertex markers in the line color, one 3pt ink baseline
    rule, uppercase direct category labels, no gridlines/legend (multi-series →
    small uppercase name labels). Mirror of _chart_brutal_line, swiss-styled."""
    x, y, w, h = box
    cats = spec.get("categories", [])
    series = spec.get("series") or []
    allv = [float(v) for s in series for v in s.get("values", [])]
    n = max((len(s.get("values", [])) for s in series), default=0)
    if not allv or n == 0:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    ink, spot = _swiss_skin()
    emph = spec.get("highlight")
    if emph is None:
        emph = len(series) - 1          # 마지막 시리즈 강조(흔한 "올해" 라인)
    vmax = max(allv); vmin = min(0.0, min(allv)); rng = (vmax - vmin) or 1.0
    TOP, BOT = 48, 44
    base_y = y + h - BOT; plot_h = h - TOP - BOT
    left = x + 50; right = x + w - 30; span = right - left
    xs = [left + (span * (i / (n - 1)) if n > 1 else span / 2) for i in range(n)]
    yf = lambda v: base_y - (v - vmin) / rng * plot_h
    line(slide, left, base_y, right, base_y, color=ink, w=3)   # 3pt 잉크 베이스룰
    one = len(series) == 1
    for si, s in enumerate(series):
        vs = [float(v) for v in s.get("values", [])]
        col = spot if (one or si == emph) else ink
        lw = 3 if (one or si == emph) else 2
        for i in range(len(vs) - 1):
            line(slide, xs[i], yf(vs[i]), xs[i + 1], yf(vs[i + 1]), color=col, w=lw)
        for i, v in enumerate(vs):
            rect(slide, [xs[i] - 7, yf(v) - 7, 14, 14], fill=col, radius=0)
    for i in range(n):
        if i < len(cats):
            text(slide, [xs[i] - 50, base_y + 8, 100, 26], str(cats[i]),
                 "label", color="ink", align="center", upper=True)
    if len(series) > 1:
        bx0 = left
        for si, s in enumerate(series):
            nm = str(s.get("name", "") or f"S{si+1}")
            col = spot if si == emph else ink
            rect(slide, [bx0, y + 8, 20, 20], fill=col, radius=0)
            text(slide, [bx0 + 28, y + 6, 40 + len(nm) * 16, 26], nm, "label",
                 color="ink", upper=True)
            bx0 += 56 + len(nm) * 16


def _chart_swiss_scatter(slide, box, spec):
    """Swiss scatter: thin (3pt) ink L-axes + small achromatic square markers,
    one spot marker via point.highlight, uppercase labels. No gridlines, no
    smooth dots. Mirror of _chart_brutal_scatter, swiss-styled."""
    x, y, w, h = box
    pts = spec.get("points", [])
    if not pts:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    ink, spot = _swiss_skin()
    xv = [float(p.get("x", 0)) for p in pts]; yv = [float(p.get("y", 0)) for p in pts]
    xmn, xmx = min(xv), max(xv); ymn, ymx = min(yv), max(yv)
    xr = (xmx - xmn) or 1.0; yr = (ymx - ymn) or 1.0
    L = x + 50; R = x + w - 30; B = y + h - 48; T = y + 20
    line(slide, L, T, L, B, color=ink, w=3)
    line(slide, L, B, R, B, color=ink, w=3)
    for p in pts:
        px = L + (float(p.get("x", 0)) - xmn) / xr * (R - L)
        py = B - (float(p.get("y", 0)) - ymn) / yr * (B - T)
        col = spot if p.get("highlight") else ink
        sz = 18
        rect(slide, [px - sz / 2, py - sz / 2, sz, sz], fill=col, radius=0)
        if p.get("label"):
            text(slide, [px + sz / 2 + 8, py - 14, 240, 28], str(p["label"]),
                 "label", color="ink", upper=True, anchor="middle")


def _chart_swiss_pie(slide, box, spec):
    """Swiss 'share' = a thin horizontal proportion bar: achromatic ink segments +
    ONE spot segment (the emphasis/largest), hairline ink dividers between segments
    (no thick border, no shadow), uppercase value labels (color flips per segment
    luminance). The deterministic swiss take on share (caller build_share lists
    categories with matching chips). Mirror of _chart_brutal_pie, swiss-styled."""
    x, y, w, h = box
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    cols = _swiss_share_colors(vals, spec)
    total = sum(vals) or 1.0
    barh = min(150, h * 0.42)
    by = y + (h - barh) / 2
    cx = x
    for i, v in enumerate(vals):
        seg = (v / total) * w
        rect(slide, [cx, by, seg, barh], fill=cols[i], radius=0)
        if i > 0:                                   # hairline ink divider
            line(slide, cx, by, cx, by + barh, color="navy", w=1.5)
        if seg > 50:
            text(slide, [cx, by + barh / 2 - 18, seg, 36], f"{v:g}", "h4",
                 color=_ink_on(cols[i]), align="center", anchor="middle")
        cx += seg
    line(slide, x, by + barh, x + cx - x, by + barh, color="navy", w=3)  # base rule


# --- luxe (dark-luxury keynote) chart skin — canvas-aware hairline charts ---
# 럭셔리 헤어라인 미감: 무채색 막대는 0.75pt 헤어라인 *아웃라인*(채움 없음), 강조 1개만
# 샴페인 골드 솔리드, 넓은 간격(막대:간격=1:1), 격자/범례 박스 제거, 직접 라벨. swiss와
# 같은 "무채색+스팟1" 절제지만 (a) 다크 캔버스 대응 위해 하드코딩 색 대신 role 토큰
# (ink/muted/accent → resolve_dark가 캔버스 휘도로 자동 도출) (b) 솔리드 막대 대신 헤어라인
# 윤곽. 다크/라이트 어디서나 가독성 유지(캔버스-인지).
def _luxe_skin():
    """(neutral, spot, hair) role tokens for luxe charts — 캔버스-인지.
    neutral=recessive 무채색 데이터선: 다크=중간 그레이(muted, 어두운 배경 위 가시),
    라이트=룩 본연의 잉크 차콜(navy, 따뜻한 어두운 선). 'muted'는 라이트 캔버스에서
    룩이 재정의 안 하면 쿨 슬레이트 디폴트로 폴백해 웜 골드 팔레트와 충돌하므로
    라이트에선 navy로 플립(accent_cycle mono_accent 채움색 픽스와 같은 거울상 규칙).
    spot=슬라이드 스팟(골드), hair=가는 축선(다크=divider 미세 / 라이트=subtle 옅게)."""
    neutral = "muted" if DARK else "navy"
    hair = "divider" if DARK else "subtle"
    return neutral, (ACCENTS[0] if ACCENTS else "blue"), hair


def _chart_luxe_bars(slide, box, spec):
    """럭셔리 막대: 무채색 헤어라인 *아웃라인* 막대(채움 없음) + 강조 1개 골드 솔리드,
    넓은 간격(1:1), 가는 축 룰, 직접 값/카테고리 라벨. 모서리 0px, 격자/범례 없음."""
    x, y, w, h = box
    cats, vals = _series_values(spec)
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    n = len(vals)
    neutral, spot, hair = _luxe_skin()
    hl = spec.get("highlight")
    if hl is None:
        hl = max(range(n), key=lambda i: vals[i])
    TOP, BOT, LBL = 48, 44, 28
    base_y = y + h - BOT
    plot_h = h - TOP - BOT
    u = w / (2 * n - 1) if n > 1 else w / 3.0      # 막대:간격 = 1:1 (넓은 럭셔리 여백)
    barw, gap = u, u
    vmax = max(vals) or 1.0
    for i, v in enumerate(vals):
        bx = x + i * (barw + gap)
        bh = max(2.0, (v / vmax) * plot_h)
        by = base_y - bh
        if i == hl:
            rect(slide, [bx, by, barw, bh], fill=spot, radius=0)        # 강조 골드 솔리드
        else:
            rect(slide, [bx, by, barw, bh], fill=None,                  # 무채색 헤어라인 윤곽
                 line_color=neutral, line_w=1.0, radius=0)
        try:
            vtxt = spec.get("value_fmt", "{:g}").format(v)
        except Exception:
            vtxt = str(v)
        text(slide, [bx - gap / 2, by - LBL - 6, barw + gap, LBL], vtxt, "label",
             color="ink", align="center", upper=True)
        if i < len(cats):
            text(slide, [bx - gap / 2, base_y + 8, barw + gap, LBL], str(cats[i]),
                 "label", color="ink", align="center", upper=True)
    line(slide, x, base_y, x + w, base_y, color=hair, w=1)     # 가는 축 룰(헤어라인)


def _chart_luxe_line(slide, box, spec):
    """럭셔리 라인: 0.75pt 골드 헤어라인 폴리라인(강조/단일) + 무채색 보조선, 면적 채움
    없음, 작은 골드 윤곽 마커, 가는 베이스룰. _chart_swiss_line의 럭셔리·캔버스-인지판."""
    x, y, w, h = box
    cats = spec.get("categories", [])
    series = spec.get("series") or []
    allv = [float(v) for s in series for v in s.get("values", [])]
    n = max((len(s.get("values", [])) for s in series), default=0)
    if not allv or n == 0:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    neutral, spot, hair = _luxe_skin()
    emph = spec.get("highlight")
    if emph is None:
        emph = len(series) - 1
    vmax = max(allv); vmin = min(0.0, min(allv)); rng = (vmax - vmin) or 1.0
    TOP, BOT = 48, 44
    base_y = y + h - BOT; plot_h = h - TOP - BOT
    left = x + 50; right = x + w - 30; span = right - left
    xs = [left + (span * (i / (n - 1)) if n > 1 else span / 2) for i in range(n)]
    yf = lambda v: base_y - (v - vmin) / rng * plot_h
    line(slide, left, base_y, right, base_y, color=hair, w=1)      # 가는 베이스룰
    one = len(series) == 1
    for si, s in enumerate(series):
        vs = [float(v) for v in s.get("values", [])]
        col = spot if (one or si == emph) else neutral
        lw = 2.5 if (one or si == emph) else 1.5
        for i in range(len(vs) - 1):
            line(slide, xs[i], yf(vs[i]), xs[i + 1], yf(vs[i + 1]), color=col, w=lw)
        for i, v in enumerate(vs):
            rect(slide, [xs[i] - 6, yf(v) - 6, 12, 12], fill=None,
                 line_color=col, line_w=1.25, radius=12)          # 골드 윤곽 마커(원)
    for i in range(n):
        if i < len(cats):
            text(slide, [xs[i] - 50, base_y + 8, 100, 26], str(cats[i]),
                 "label", color="ink", align="center", upper=True)
    if len(series) > 1:
        bx0 = left
        for si, s in enumerate(series):
            nm = str(s.get("name", "") or f"S{si+1}")
            col = spot if si == emph else neutral
            line(slide, bx0, y + 18, bx0 + 22, y + 18, color=col, w=2.5)
            text(slide, [bx0 + 28, y + 6, 40 + len(nm) * 16, 26], nm, "label",
                 color="ink", upper=True)
            bx0 += 60 + len(nm) * 16


def _chart_luxe_scatter(slide, box, spec):
    """럭셔리 산점도: 가는 헤어라인 L축 + 작은 무채색 윤곽 원 마커, 강조 1개 골드 솔리드."""
    x, y, w, h = box
    pts = spec.get("points", [])
    if not pts:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    neutral, spot, hair = _luxe_skin()
    xv = [float(p.get("x", 0)) for p in pts]; yv = [float(p.get("y", 0)) for p in pts]
    xmn, xmx = min(xv), max(xv); ymn, ymx = min(yv), max(yv)
    xr = (xmx - xmn) or 1.0; yr = (ymx - ymn) or 1.0
    L = x + 50; R = x + w - 30; B = y + h - 48; T = y + 20
    line(slide, L, T, L, B, color=hair, w=1)
    line(slide, L, B, R, B, color=hair, w=1)
    for p in pts:
        px = L + (float(p.get("x", 0)) - xmn) / xr * (R - L)
        py = B - (float(p.get("y", 0)) - ymn) / yr * (B - T)
        sz = 16
        if p.get("highlight"):
            rect(slide, [px - sz / 2, py - sz / 2, sz, sz], fill=spot, radius=sz)
        else:
            rect(slide, [px - sz / 2, py - sz / 2, sz, sz], fill=None,
                 line_color=neutral, line_w=1.25, radius=sz)
        if p.get("label"):
            text(slide, [px + sz / 2 + 8, py - 14, 240, 28], str(p["label"]),
                 "label", color="ink", upper=True, anchor="middle")


def _chart_luxe_share(slide, box, spec):
    """럭셔리 비중: 가는 비례 막대 — 무채색 윤곽 세그먼트 + 강조 1개 골드 솔리드,
    헤어라인 디바이더. _chart_swiss_pie의 럭셔리·캔버스-인지판."""
    x, y, w, h = box
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    neutral, spot, hair = _luxe_skin()
    hl = spec.get("highlight")
    if hl is None:
        hl = max(range(len(vals)), key=lambda i: vals[i])
    total = sum(vals) or 1.0
    barh = min(140, h * 0.40)
    by = y + (h - barh) / 2
    cx = x
    for i, v in enumerate(vals):
        seg = (v / total) * w
        if i == hl:
            rect(slide, [cx, by, seg, barh], fill=spot, radius=0)
            tcol = _ink_on(spot)
        else:
            rect(slide, [cx, by, seg, barh], fill=None,
                 line_color=neutral, line_w=1.0, radius=0)
            tcol = "ink"
        if seg > 50:
            text(slide, [cx, by + barh / 2 - 18, seg, 36], f"{v:g}", "h4",
                 color=tcol, align="center", anchor="middle")
        cx += seg
    line(slide, x, by + barh, cx, by + barh, color=hair, w=1)


# --- dark-tech composed charts (grammar.chart_style == "tech") -----------------
# The pack's chart identity: "차콜 위 발광 데이터 + 모노스페이스 수치 라벨". Drawn
# (not native) so EVERY value/axis label is the monospace meta face and the emphasis
# mark gets a neon glow. Bars = flat fill, ONE cyan emphasis bar (+1px top-edge glow),
# the rest dim neutral #3A3D45; line = glowing polyline with a 10%-alpha area fade,
# straight segments, end-point glow dots. Routed only for the tech look (others keep
# their renderer); value labels read mono via the active mono_meta grammar (the bars
# pass the SERIES color so the number glows in cyan/violet, per the spec).
def _tech_neutral():
    """Dim neutral fill for non-emphasis tech bars. The pack names #3A3D45; use the
    token if the look defines one, else the literal (kept on-palette by the look)."""
    return "#3A3D45"


def _chart_tech_bars(slide, box, spec):
    x, y, w, h = box
    cats, vals = _series_values(spec)
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    n = len(vals)
    spot = ACCENTS[0] if ACCENTS else "blue"     # 네온 시안 강조
    neutral = _tech_neutral()
    hl = spec.get("highlight")
    if hl is None:
        hl = max(range(n), key=lambda i: vals[i])
    TOP, BOT, LBL = 52, 44, 30
    base_y = y + h - BOT
    plot_h = h - TOP - BOT
    u = w / (5 * n - 2) if n > 1 else w / 3.0    # 막대:간격 = 3:2 (스펙)
    barw, gap = 3 * u, 2 * u
    vmax = max(vals) or 1.0
    # 어두운 가로 격자 헤어라인 1개(중간값) — 발광 없음, 차트는 가라앉힘.
    line(slide, x, base_y - plot_h * 0.5, x + w, base_y - plot_h * 0.5,
         color="divider", w=0.75)
    for i, v in enumerate(vals):
        bx = x + i * (barw + gap)
        bh = max(2.0, (v / vmax) * plot_h)
        by = base_y - bh
        emph = (i == hl)
        sp = rect(slide, [bx, by, barw, bh], fill=(spot if emph else neutral),
                  radius=2)
        if emph:
            _glow_shape(sp, "chart", color=_resolve(spot))   # 네온 막대 헤일로
            # 1px 네온 시안 상단 엣지(스펙: "막대 상단 1px 네온 글로우 엣지")
            edge = rect(slide, [bx, by, barw, 2], fill=spot, radius=0)
            _glow_shape(edge, "chart", color=_resolve(spot))
        try:
            vtxt = spec.get("value_fmt", "{:g}").format(v)
        except Exception:
            vtxt = str(v)
        # 값 레이블 = 모노(강조는 시리즈색, 그 외 밝은 ink), 막대 끝 바로 위에 직접
        text(slide, [bx - gap / 2, by - LBL - 4, barw + gap, LBL], vtxt, "label",
             color=(spot if emph else "ink"), align="center", upper=True)
        if i < len(cats):
            text(slide, [bx - gap / 2, base_y + 8, barw + gap, LBL], str(cats[i]),
                 "label", color="muted", align="center", upper=True)
    line(slide, x, base_y, x + w, base_y, color="divider", w=0.75)   # x 베이스라인(0 기준)


def _chart_tech_line(slide, box, spec):
    x, y, w, h = box
    cats = spec.get("categories", [])
    series = spec.get("series") or []
    allv = [float(v) for s in series for v in s.get("values", [])]
    n = max((len(s.get("values", [])) for s in series), default=0)
    if not allv or n == 0:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    spot = ACCENTS[0] if ACCENTS else "blue"
    violet = ACCENTS[1] if len(ACCENTS) > 1 else spot
    emph = spec.get("highlight")
    if emph is None:
        emph = len(series) - 1
    vmax = max(allv); vmin = min(0.0, min(allv)); rng = (vmax - vmin) or 1.0
    TOP, BOT = 52, 44
    base_y = y + h - BOT; plot_h = h - TOP - BOT
    left = x + 50; right = x + w - 30; span = right - left
    xs = [left + (span * (i / (n - 1)) if n > 1 else span / 2) for i in range(n)]
    yf = lambda v: base_y - (v - vmin) / rng * plot_h
    line(slide, left, base_y - plot_h * 0.5, right, base_y - plot_h * 0.5,
         color="divider", w=0.75)                              # 어두운 격자 1개
    line(slide, left, base_y, right, base_y, color="divider", w=0.75)  # 베이스룰
    for si, s in enumerate(series):
        vs = [float(v) for v in s.get("values", [])]
        col = spot if (len(series) == 1 or si == emph) else violet
        for i in range(len(vs) - 1):
            ln = line(slide, xs[i], yf(vs[i]), xs[i + 1], yf(vs[i + 1]),
                      color=col, w=2)
            _glow_shape(ln, "chart", color=_resolve(col))      # 발광 데이터 라인
        # 끝점 발광 마커 + 시리즈명 모노 라벨(직접 레이블링)
        if vs:
            mk = oval(slide, [xs[-1] - 6, yf(vs[-1]) - 6, 12, 12], fill=col)
            _glow_shape(mk, "chart", color=_resolve(col))
            nm = str(s.get("name", "") or "")
            if nm:
                text(slide, [xs[-1] + 12, yf(vs[-1]) - 14, 220, 28], nm, "label",
                     color=col, anchor="middle", upper=True)
    for i in range(n):
        if i < len(cats):
            text(slide, [xs[i] - 50, base_y + 8, 100, 28], str(cats[i]),
                 "label", color="muted", align="center", upper=True)


# --- prismatic-dark composed charts (grammar.chart_style == "prism") -----------
# The pack's chart identity: "발광 그라디언트 차트 — 정밀도 우선". Same data-ink-minimal
# geometry as the tech charts, but the DATA marks render in the prism cyan→magenta→amber
# gradient: bars = 35%-alpha gradient fill + a bright top glow cap; line = gradient
# polyline + a fade area (gradient 18%→0); donut = gradient arc with the key figure in
# the centre. Value labels stay crisp mono (the pack forbids glow on text). Routed only
# for the prism look; every other deck keeps its own renderer.
def _chart_prism_bars(slide, box, spec):
    x, y, w, h = box
    cats, vals = _series_values(spec)
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    n = len(vals)
    spot = ACCENTS[0] if ACCENTS else "blue"
    hl = spec.get("highlight")
    if hl is None:
        hl = max(range(n), key=lambda i: vals[i])
    TOP, BOT, LBL = 52, 44, 30
    base_y = y + h - BOT
    plot_h = h - TOP - BOT
    u = w / (5 * n - 2) if n > 1 else w / 3.0    # 막대:간격 = 3:2 (스펙)
    barw, gap = 3 * u, 2 * u
    vmax = max(vals) or 1.0
    line(slide, x, base_y - plot_h * 0.5, x + w, base_y - plot_h * 0.5,
         color="divider", w=0.75)                 # 어두운 격자 헤어라인 1개
    for i, v in enumerate(vals):
        bx = x + i * (barw + gap)
        bh = max(2.0, (v / vmax) * plot_h)
        by = base_y - bh
        emph = (i == hl)
        # 모든 막대 = 프리즘 그라디언트 채움(강조는 불투명, 그 외 35% 알파 — 스펙).
        sp = rect(slide, [bx, by, barw, bh], fill=spot, radius=2)
        _gradient_shape_fill(sp, alpha=None if emph else 35)
        # 윗면 0.04in 밝은 그라디언트 글로우 캡(스펙). 강조 막대만 발광 헤일로.
        cap = rect(slide, [bx, by, barw, max(3, plot_h * 0.04)], fill=spot, radius=0)
        _gradient_shape_fill(cap)
        if emph:
            _glow_shape(cap, "chart", color=_resolve(spot))
        try:
            vtxt = spec.get("value_fmt", "{:g}").format(v)
        except Exception:
            vtxt = str(v)
        text(slide, [bx - gap / 2, by - LBL - 4, barw + gap, LBL], vtxt, "label",
             color="ink", align="center", upper=True)   # 값 = 또렷한 모노(글로우 금지)
        if i < len(cats):
            text(slide, [bx - gap / 2, base_y + 8, barw + gap, LBL], str(cats[i]),
                 "label", color="muted", align="center", upper=True)
    line(slide, x, base_y, x + w, base_y, color="divider", w=0.75)   # x 베이스라인


def _chart_prism_line(slide, box, spec):
    x, y, w, h = box
    cats = spec.get("categories", [])
    series = spec.get("series") or []
    allv = [float(v) for s in series for v in s.get("values", [])]
    n = max((len(s.get("values", [])) for s in series), default=0)
    if not allv or n == 0:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    stops = _prism_stops()
    spot = ACCENTS[0] if ACCENTS else "blue"
    vmax = max(allv); vmin = min(0.0, min(allv)); rng = (vmax - vmin) or 1.0
    TOP, BOT = 52, 44
    base_y = y + h - BOT; plot_h = h - TOP - BOT
    left = x + 50; right = x + w - 30; span = right - left
    xs = [left + (span * (i / (n - 1)) if n > 1 else span / 2) for i in range(n)]
    yf = lambda v: base_y - (v - vmin) / rng * plot_h
    line(slide, left, base_y - plot_h * 0.5, right, base_y - plot_h * 0.5,
         color="divider", w=0.75)                              # 어두운 격자 1개
    line(slide, left, base_y, right, base_y, color="divider", w=0.75)  # 베이스룰
    for si, s in enumerate(series):
        vs = [float(v) for v in s.get("values", [])]
        # 발광 그라디언트 폴리라인(2pt) — 세그먼트마다 x 위치 비율로 스톱을 잘라 좌→우
        # 전체에 걸쳐 프리즘 색 전이를 만든다(세그먼트 단색 회피).
        for i in range(len(vs) - 1):
            ln = line(slide, xs[i], yf(vs[i]), xs[i + 1], yf(vs[i + 1]),
                      color=spot, w=2)
            t0, t1 = i / max(1, n - 1), (i + 1) / max(1, n - 1)
            seg = [_mix("#" + stops[0], "#" + stops[-1], t)
                   for t in (t0, (t0 + t1) / 2, t1)]
            gradient_stroke(ln, width_pt=2, stops=seg)
            _glow_shape(ln, "chart", color=_resolve(spot))
        # 끝점 발광 마커
        if vs:
            mk = oval(slide, [xs[-1] - 6, yf(vs[-1]) - 6, 12, 12], fill=stops[-1])
            _gradient_shape_fill(mk)
            _glow_shape(mk, "chart", color=stops[-1])
            nm = str(s.get("name", "") or "")
            if nm:
                text(slide, [xs[-1] + 12, yf(vs[-1]) - 14, 220, 28], nm, "label",
                     color="ink", anchor="middle", upper=True)
    for i in range(n):
        if i < len(cats):
            text(slide, [xs[i] - 50, base_y + 8, 100, 28], str(cats[i]),
                 "label", color="muted", align="center", upper=True)


def _chart_prism_share(slide, box, spec):
    """비중 = 그라디언트 도넛 아크(조각 3개 이하) + 중앙 핵심 수치. PIE 도넛은 COM 렌더가
    불안정해(brutal 의 교훈) 비례 가로 막대 세그먼트로 근사하되, 각 세그먼트를 프리즘
    스톱색으로 칠하고 강조 글로우 — '그라디언트 아크'의 정직한 근사."""
    x, y, w, h = box
    cats = spec.get("categories", [])
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    stops = _prism_stops()
    total = sum(vals) or 1.0
    # 도넛 대신 굵은 비례 링 바(높이 0.32in 캡슐), 세그먼트별 프리즘 스톱.
    bar_h = min(60, h * 0.22)
    by = y + h * 0.42
    cx = x
    for i, v in enumerate(vals[:6]):
        seg_w = (v / total) * w
        col = "#" + stops[i % len(stops)]
        sp = rect(slide, [cx, by, max(4, seg_w - 4), bar_h], fill=col, radius=8)
        _glow_shape(sp, "chart", color=col)
        if i < len(cats):
            text(slide, [cx, by + bar_h + 10, max(40, seg_w), 26],
                 str(cats[i]), "label", color="muted", align="left", upper=True)
            text(slide, [cx, by - 30, max(40, seg_w), 26],
                 "%g" % v, "label", color="ink", align="left", upper=True)
        cx += seg_w


# --- hyundai cinematic charts (grammar.chart_style == "cinematic") -------------
# The pack's chart identity: "다크 네이비 무대 위에 사이안 한 톤으로 발광하며, 항목별
# 무지개 색칠 없이 톤 레이어링만 사용하는 절제된 키노트 차트". The opposite of the other
# four dark looks' chart grammars: NO mono value labels (the pack forbids monospace),
# NO per-series rainbow, NO gradient bar FILL. Instead a single cyan-blue family:
#   • bars  → tonal stack (#66FFFF / #00D3FD / #0060A8), ONE emphasis bar cyan-glowing
#   • line  → #00D3FD 2.5pt line + a 15-30% cyan AREA fill + white-fill/cyan-RING dots
#             (the pack: 풀-사이안 dot 절대 금지) + an optional chevron-tab STAGE band
#   • share → cyan tonal proportion bars (single family, no rainbow)
# Value labels stay proportional white (emphasis value = cyan). Drawn, not native, so
# the dots, area, and stage band all read. Routed only for the cinematic look.
_CINE_TONES = ("#66FFFF", "#00D3FD", "#0060A8")     # top / mid / base cyan stack


def _cine_area_fill(slide, xs, ys, base_y, color="#00D3FD", alpha=22):
    """Translucent cyan area under a line — a freeform polygon (xs across, ys up,
    closed along the baseline). Gives the cinematic 'linear-gradient → 0' wash the
    pack draws below its hydrogen-market line. Solid cyan at low alpha (PPTX area
    gradient under a freeform is unreliable; flat low-alpha is the honest version)."""
    if len(xs) < 2:
        return
    try:
        fb = slide.shapes.build_freeform(IN(xs[0]), IN(base_y), scale=Emu(1))
        pts = [(IN(xs[i]), IN(ys[i])) for i in range(len(xs))]
        pts.append((IN(xs[-1]), IN(base_y)))
        fb.add_line_segments(pts, close=True)
        sp = fb.convert_to_shape()
        sp.shadow.inherit = False
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(color)
        _apply_alpha(sp, alpha)
        sp.line.fill.background()
    except Exception:
        pass


def _cine_stage_band(slide, box, spec):
    """Draw the chevron-tab stage band ABOVE the plot when the look enables it and
    the spec supplies `stage` = list of labels (+ optional `stage_active` index).
    Returns the y the plot should start below. No band → returns box y unchanged."""
    cfg = _chevron_cfg()
    stages = spec.get("stage")
    if cfg is None or not cfg.get("stage_band", True) or not stages:
        return box[1]
    x, y, w, h = box
    band_h = 56
    _stage_band(slide, [x, y, w, band_h], stages, spec.get("stage_active"))
    return y + band_h + 24


def _chart_cinematic_bars(slide, box, spec):
    x, y, w, h = box
    cats, vals = _series_values(spec)
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    y = _cine_stage_band(slide, box, spec); h = box[1] + box[3] - y
    n = len(vals)
    hl = spec.get("highlight")
    if hl is None:
        hl = max(range(n), key=lambda i: vals[i])
    TOP, BOT, LBL = 48, 44, 30
    base_y = y + h - BOT
    plot_h = h - TOP - BOT
    u = w / (5 * n - 2) if n > 1 else w / 3.0
    barw, gap = 3 * u, 2 * u
    vmax = max(vals) or 1.0
    line(slide, x, base_y, x + w, base_y, color="divider", w=0.75)
    for i, v in enumerate(vals):
        bx = x + i * (barw + gap)
        bh = max(2.0, (v / vmax) * plot_h)
        by = base_y - bh
        emph = (i == hl)
        # single cyan family: emphasis = bright cyan + glow; rest = mid/base tonal
        col = _CINE_TONES[0] if emph else (_CINE_TONES[1] if i % 2 == 0 else _CINE_TONES[2])
        sp = rect(slide, [bx, by, barw, bh], fill=col, radius=2)
        if emph:
            _glow_shape(sp, "chart", color="#66FFFF")
        try:
            vtxt = spec.get("value_fmt", "{:g}").format(v)
        except Exception:
            vtxt = str(v)
        text(slide, [bx - gap / 2, by - LBL - 4, barw + gap, LBL], vtxt, "label",
             color=("blue" if emph else "ink"), align="center")
        if i < len(cats):
            text(slide, [bx - gap / 2, base_y + 8, barw + gap, LBL], str(cats[i]),
                 "label", color="muted", align="center")


def _chart_cinematic_line(slide, box, spec):
    x, y, w, h = box
    cats = spec.get("categories", [])
    series = spec.get("series") or []
    allv = [float(v) for s in series for v in s.get("values", [])]
    n = max((len(s.get("values", [])) for s in series), default=0)
    if not allv or n == 0:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    y = _cine_stage_band(slide, box, spec); h = box[1] + box[3] - y
    vmax = max(allv); vmin = min(0.0, min(allv)); rng = (vmax - vmin) or 1.0
    TOP, BOT = 48, 44
    base_y = y + h - BOT; plot_h = h - TOP - BOT
    left = x + 50; right = x + w - 30; span = right - left
    xs = [left + (span * (i / (n - 1)) if n > 1 else span / 2) for i in range(n)]
    yf = lambda v: base_y - (v - vmin) / rng * plot_h
    line(slide, left, base_y, right, base_y, color="divider", w=0.75)
    for si, s in enumerate(series):
        vs = [float(v) for v in s.get("values", [])]
        ys = [yf(v) for v in vs]
        # area fill under the (first/emphasis) series only — the cinematic wash
        if si == 0:
            _cine_area_fill(slide, xs[:len(vs)], ys, base_y,
                            color=_CINE_TONES[1], alpha=22)
        col = _CINE_TONES[1] if (len(series) == 1 or si == 0) else _CINE_TONES[2]
        for i in range(len(vs) - 1):
            ln = line(slide, xs[i], ys[i], xs[i + 1], ys[i + 1], color=col, w=2.5)
            _glow_shape(ln, "chart", color="#66FFFF")
        # white-fill + cyan-RING data dots at every point (pack: 풀-사이안 dot 금지)
        for i in range(len(vs)):
            oval(slide, [xs[i] - 5, ys[i] - 5, 10, 10], fill="white",
                 line_color=col, line_w=1.75)
        nm = str(s.get("name", "") or "")
        if nm and vs:
            text(slide, [xs[-1] + 12, ys[-1] - 14, 240, 28], nm, "label",
                 color="blue", anchor="middle")
    # direct value labels above each emphasis point (white; cyan would clutter)
    if series:
        vs0 = [float(v) for v in series[0].get("values", [])]
        for i, v in enumerate(vs0):
            try:
                vt = spec.get("value_fmt", "{:g}").format(v)
            except Exception:
                vt = str(v)
            text(slide, [xs[i] - 50, yf(v) - 34, 100, 26], vt, "label",
                 color="ink", align="center")
    for i in range(n):
        if i < len(cats):
            text(slide, [xs[i] - 50, base_y + 8, 100, 26], str(cats[i]),
                 "label", color="muted", align="center")


def _chart_cinematic_share(slide, box, spec):
    """비중 = 사이안 단톤 비례 막대(무지개 금지). 강조 세그먼트만 발광."""
    x, y, w, h = box
    cats = spec.get("categories", [])
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    vals = [float(v) for v in (vals or [])]
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    total = sum(vals) or 1.0
    bar_h = min(60, h * 0.22)
    by = y + h * 0.42
    cx = x
    hl = spec.get("highlight", 0)
    for i, v in enumerate(vals[:6]):
        seg_w = (v / total) * w
        col = _CINE_TONES[i % len(_CINE_TONES)]
        sp = rect(slide, [cx, by, max(4, seg_w - 4), bar_h], fill=col, radius=8)
        if i == hl:
            _glow_shape(sp, "chart", color="#66FFFF")
        if i < len(cats):
            text(slide, [cx, by + bar_h + 10, max(40, seg_w), 26],
                 str(cats[i]), "label", color="muted", align="left")
            text(slide, [cx, by - 30, max(40, seg_w), 26],
                 "%g" % v, "label", color="ink", align="left")
        cx += seg_w


# --- memphis decorative scatter (grammar.decor == "memphis") -------------------
# 멤피스 정체성의 핵심이자 brutalism/swiss/luxe엔 없던 '장식 레이어' 프리미티브:
# 정원·삼각형·검정 도트그리드·지그재그 띠가 슬라이드 가장자리에 부유한다. build 루프가
# 콘텐츠 렌더 직후 슬라이드별로 1회 호출 — idx 결정론 배치(같은 배치 반복 금지),
# 좌상단 제목 영역을 피해 우/하단 코너·엣지에만(블리드 허용). 채움 도형은 원색 평면 +
# 검정 3pt 외곽선 + 6pt 하드 오프셋 섀도(차트/노드와 동일한 멤피스 어휘).
_MEMPHIS_KINDS = ("circle", "triangle", "dots", "zigzag")
_MEMPHIS_ZONES = (              # (x, y, w, h) px — 좌상단(제목) 회피, 일부 캔버스 밖 블리드
    (1740, 36, 150, 150),       # 우상단
    (1780, 470, 130, 130),      # 우엣지 중앙(블리드)
    (1730, 880, 150, 150),      # 우하단
    (40, 900, 130, 130),        # 좌하단
    (560, 980, 150, 150),       # 하단 중앙(블리드)
)


def _memphis_one(slide, kind, box, fill):
    x, y, w, h = box
    if kind == "dots":                       # 검정 도트 그리드(보더·섀도 없음)
        d, gap = 14, 30
        for r in range(3):
            for c in range(3):
                oval(slide, [x + c * gap, y + r * gap, d, d], fill="navy")
        return
    if kind == "zigzag":                     # 검정 지그재그 띠(굵은 꺾은선)
        seg = w / 4.0
        pts = [(x + seg * j, y + (0 if j % 2 == 0 else h * 0.5)) for j in range(5)]
        for j in range(len(pts) - 1):
            line(slide, pts[j][0], pts[j][1], pts[j + 1][0], pts[j + 1][1],
                 color="navy", w=5)
        return
    shp = MSO_SHAPE.ISOSCELES_TRIANGLE if kind == "triangle" else MSO_SHAPE.OVAL
    sp = slide.shapes.add_shape(shp, IN(x), IN(y), IN(w), IN(h))
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = C(fill)
    sp.line.color.rgb = C("navy")
    sp.line.width = Pt(3)
    shadow(sp, blur_pt=0, dist_pt=6, dir_deg=135, color="#1A1A1A", alpha=0)
    return sp


def _memphis_decor(slide, idx):
    """슬라이드 가장자리에 멤피스 장식 도형 3개를 idx 결정론으로 흩뿌린다."""
    cols = ACCENTS or ["blue"]
    nz = len(_MEMPHIS_ZONES)
    for k in range(3):
        z = _MEMPHIS_ZONES[(idx * 2 + k) % nz]
        kind = _MEMPHIS_KINDS[(idx + k) % len(_MEMPHIS_KINDS)]
        col = cols[(idx + k) % len(cols)]
        _memphis_one(slide, kind, z, col)


# ---------------------------------------------------------------------------
# 일반 가독성 레이어 — 자동 차트 백드롭 (busy 배경 위 fragile 차트 보호)
# ---------------------------------------------------------------------------
# busy(다색 메시) 배경 위에서 *얇은* 데이터 마크(폴리라인·산점·기울기선)는 메시색과
# 충돌해 가독성이 깨진다. _chart_backdrop() 는 차트 영역 뒤에 calm 다크 패널 1장을 깔아
# 마크를 차분한 바탕 위로 올린다. 이 패널은 장식 카드 frost(흰 틴트 저알파)와 *다르다* —
# 카드 frost 는 패널을 밝게만 해 다크 캔버스에선 메시를 더 비치게 만들어 얇은 마크를
# 오히려 묻는다. 차트 백드롭의 존재 의의는 가독성이라, 배경 변동을 *억제*해야 한다 →
# 캔버스를 검정쪽으로 살짝 내린 어두운 채움을 메시를 거의 누를 만큼 높은 alpha
# (_CHART_BACKDROP_ALPHA, busy 면 _ALPHA_BUSY)로 깐다. frost 레짐과 채움색·alpha 모두 분리.
#   FRAGILE(백드롭 O): 폴리라인·산점도·버블·기울기선 + (스펙트) 조밀 라벨 테이블.
#   ROBUST(백드롭 X): 막대/컬럼(굵은 마크 자체가 바탕) · share 도넛/파이 · 면적 채움 ·
#                     quadrant/matrix/kpi(자기 셀·슬랩 바탕이 이미 대비 제공).
# 호출 조건 = busy AND ctype∈FRAGILE. calm 배경(hyundai·engineered)은 busy=False라
# 자동 미발동 → 원본(차트 패널 없음) 일치. 룩은 grammar.chart_backdrop 로 강제 on/off
# 가능(True=항상 / False=절대 / 리스트=해당 타입만).
_FRAGILE_CHART_TYPES = frozenset({
    "line", "line_markers", "xy_scatter", "bubble", "slope",
})
_CHART_BACKDROP_PAD = 14   # 패널이 차트 box 보다 살짝 넉넉히 — 축 라벨까지 감싸는 여백
# 차트 백드롭 틴트는 그 룩의 *카드 프로스트 방향*을 따라간다(resolve_dark 철학 — 룩 일관성을
# 엔진이 보장). 카드 frost 레짐(card() 1383~1386)이 busy=frost(흰 틴트 저알파)·calm-depth=
# 다크글래스로 갈리듯, 차트 백드롭도 같은 레짐 분기를 쓰되 *더 불투명*하게 깐다(얇은 데이터
# 마크가 메시 위로 또렷이 올라오게).
#
# 과거(다크 슬랩) 패착: vivid 는 라이트 프로스트 카드(surface 틴트 16% alpha) 룩인데, 차트
# 백드롭만 캔버스를 검정쪽으로 내린 near-opaque 다크(86~90%)로 깔아 *룩이 깨졌다* — 하단
# KPI/비교 카드(16% 라이트)와 같은 슬라이드에서 차트 패널만 다크 슬랩으로 떠 패밀리룩·opacity
# 가 어긋났다. 원본 design-pick vivid 는 라인 차트도 "라이트 프로스트 베일"(메시가 은은히
# 비치되 차분) 위에 발광 라인을 얹는다. → 백드롭도 카드와 같은 라이트 프로스트 패밀리로.
#
# 수정: busy(frost) 레짐 룩이면 백드롭도 *라이트 프로스트*(_FROST_TINT 흰 틴트)로 깔되, 카드
# 16%보다 더 불투명(_CHART_BACKDROP_FROST_ALPHA)하게 메시를 눌러 마크 가독성을 확보한다.
# calm-depth(다크글래스) 레짐 룩이 향후 busy 가 아닌 채로 fragile 차트를 띄우면(현재 vivid 만
# busy 라 실질 미발생) 카드 다크글래스 방향을 따라 캔버스 파생 다크 틴트를 깐다. 알파 시작값은
# 명세대로 ~35%(렌더 육안으로 튜닝 — 옅으면 ↑, 너무 불투명하면 ↓). 라인/마크 가독성은 발광
# (_glow_shape) 유지로 확보 — 라이트 프로스트 위에서도 시안·마젠타 라인이 헤일로로 읽힌다.
_CHART_BACKDROP_FROST_ALPHA = 48      # busy(frost) 레짐: 라이트 흰 프로스트 불투명도(%) — 카드 16%보다
                                      # 불투명해 메시를 눌러 마크가 읽힘. 시작값 35%는 마젠타 메시존에서
                                      # 마젠타 라인·옅은 산점이 약했다(figure-ground 부족) → 48%로 올려
                                      # 흰 프로스트를 더 깔아 ground 를 밝고 균일하게(채도 라인 대비 ↑).
                                      # 너무 옅으면 ↑/너무 불투명(메시 사라짐)이면 ↓(렌더 튜닝)
_CHART_BACKDROP_CANVAS_DARKEN = 0.30  # calm-depth(다크글래스) 레짐: 캔버스를 검정쪽으로 내려 차분한 다크 그라운드
_CHART_BACKDROP_ALPHA = 86            # calm-depth: 다크 백드롭 기본 불투명도(%)
_CHART_BACKDROP_ALPHA_BUSY = 90       # (calm 레짐 내 다채 변동 대비 여유 — 현재 busy=frost 분기로 미사용)
_CHART_BACKDROP_INSET = 6             # 플롯 영역을 살짝 안으로 — 패널이 차트를 프레임


def _chart_is_fragile(ctype, spec):
    """ctype(+spec 힌트)이 busy 배경 위에서 백드롭이 필요한 *얇은* 마크 종류인가.
    spread(히스토그램)는 ctype=column_clustered 라 타입만으론 막대와 구별 안 되므로
    빌더가 spec['_fragile']=True 힌트를 박는다(명시 > 추론)."""
    if spec.get("_fragile"):
        return True
    return ctype in _FRAGILE_CHART_TYPES


def _chart_backdrop_regime():
    """차트 백드롭의 (tint, alpha, frost) = 카드 프로스트 *방향*을 따라간다.
    card() 의 자동 카드 반투명 분기(1383~1386)와 동일 레짐 판정:
      busy(다색 메시) 룩  → frost: 흰 틴트(_FROST_TINT) + _CHART_BACKDROP_FROST_ALPHA.
                            카드 frost(흰 16%)와 같은 라이트 프로스트 패밀리, 단 더 불투명
                            (메시를 눌러 얇은 데이터 마크가 읽히게).
      calm-depth(단광원·동색조 어두운 그라디언트) 룩 → 다크글래스: 캔버스를 검정쪽으로
                            내린 다크 틴트 + 높은 alpha(차분한 다크 그라운드).
    이 함수는 색·알파만 고르고, frost 여부(True=흰 틴트로 *교체*, False=원색 유지)를
    함께 돌려준다(_apply_frost vs _apply_alpha 분기용). 백드롭 적용부에서 DARK 게이트는
    이미 _maybe_chart_backdrop 가 건다."""
    if _bg_atmosphere().get("busy"):
        tint, alpha, frost = _FROST_TINT, _CHART_BACKDROP_FROST_ALPHA, True
    else:
        canvas_hex = _resolve("canvas")
        dark_tint = _mix(canvas_hex, "#000000", _CHART_BACKDROP_CANVAS_DARKEN)
        tint, alpha, frost = dark_tint, _CHART_BACKDROP_ALPHA, False
    # HITL 노브: 일관성↔가독성 균형은 결정론으로 못 박는 아트디렉팅 판단 → 룩 토큰
    # grammar.chart_backdrop_opacity 또는 CLI --chart-backdrop-opacity 가 있으면 가독
    # 바닥선 기본을 덮어쓴다(0~100). 엔진은 "바닥선"만 기본 제공, 그 위 균형은 사람이 정함.
    ov = GRAMMAR.get("chart_backdrop_opacity")
    if ov is not None:
        try:
            alpha = max(0.0, min(100.0, float(ov)))
        except (TypeError, ValueError):
            pass
    return (tint, alpha, frost)


def _chart_backdrop(slide, box):
    """차트 영역 뒤 프로스트 패널 1장. 틴트·알파·레짐은 _chart_backdrop_regime() 이
    카드 프로스트 방향에서 도출(busy=라이트 흰 프로스트 / calm-depth=다크글래스). frost
    레짐이면 채움을 흰 틴트로 교체(_apply_frost) — 어떤 메시색이 비쳐도 균일한 라이트
    프로스트 베일이 되어 발광 라인이 그 위로 또렷이 올라온다. 패널은 box 보다
    _CHART_BACKDROP_PAD 만큼 넉넉히 그려 축/라벨까지 프레임하고, 라운드는 카드 프리셋
    반경을 따른다(덱 일관 — KPI/비교 카드와 같은 라운드·같은 프로스트 패밀리)."""
    x, y, w, h = box
    p = _CHART_BACKDROP_PAD
    c = COMPONENTS.get("card", {})
    rad = c.get("radius", 0)
    rad = RAD.get(rad, rad) if isinstance(rad, str) else (rad or 0)
    tint, alpha, frost = _chart_backdrop_regime()
    sp = rect(slide, [x - p, y - p, w + 2 * p, h + 2 * p],
              fill=tint, radius=rad)
    if frost:
        _apply_frost(sp, tint, alpha)   # 흰 틴트로 교체 + 저알파 = 라이트 프로스트 베일
    else:
        _apply_alpha(sp, alpha)         # 다크 틴트 원색 유지 + 높은 alpha = 다크글래스
    return sp


def _maybe_chart_backdrop(slide, box, ctype, spec):
    """차트 그리기 직전 호출 — 자동/override 정책에 따라 백드롭을 깐다(차트 도형보다
    먼저 그려 뒤에 깔리게). 다크 캔버스에서만(라이트=반투명 패널 무의미).
    백드롭을 깔았으면 플롯이 패널 안으로 _CHART_BACKDROP_INSET 만큼 들어간 box 를
    반환(패널이 차트를 프레임). 안 깔았으면 원본 box 그대로(무회귀)."""
    if not DARK:
        return box
    ov = GRAMMAR.get("chart_backdrop")
    if ov is True:
        draw = True
    elif ov is False:
        draw = False
    elif isinstance(ov, (list, tuple, set)):
        draw = ctype in ov or (spec.get("_fragile") and "spread" in ov)
    else:  # 미선언 → 자동(_bg_atmosphere)
        draw = _bg_atmosphere()["busy"] and _chart_is_fragile(ctype, spec)
    if not draw:
        return box
    _chart_backdrop(slide, box)
    i = _CHART_BACKDROP_INSET
    x, y, w, h = box
    return [x + i, y + i, w - 2 * i, h - 2 * i]


def chart(slide, box, spec):
    """차트 디스패처. spec={type, categories, series:[{name,values}], values, ...}.
    native(python-pptx) 12종 + 도형 합성 8종. 팔레트 램프로 시리즈 색 지정.
    빈 데이터면 안내 플레이스홀더. 좌표는 호출자가 넘긴 box."""
    spec = spec or {}
    ctype = spec.get("type", "column_clustered")
    cats = spec.get("categories", [])
    series = spec.get("series")
    if not series and "values" in spec:
        series = [{"name": spec.get("series_name", ""), "values": spec["values"]}]
    series = series or []
    # Normalized series must be written back: _chart_native reads spec["series"]
    # directly, so a chart given a bare `values` list (e.g. spread/histogram)
    # would otherwise render with ZERO series — an empty plot area. (bug C)
    spec["series"] = series
    has_data = bool(series) and any(s.get("values") for s in series)

    # 일반 가독성 레이어: busy 배경 위 fragile 차트 백드롭(차트 도형보다 먼저 = 뒤에 깔림).
    # depth/calm 배경 룩은 busy=False라 무발동 → 원본(맨 차트) 일치. depth 토큰 없는 룩은
    # DARK 게이트 + busy=False 로 이중 무발동(무회귀). 렌더 데이터가 있을 때만(빈
    # 플레이스홀더 제외) — 산점/버블은 series 가 아니라 points 에 데이터가 있으므로
    # has_data 만으론 누락된다(디스패처의 data 조건과 정합: series·points·values 중 하나).
    # 백드롭을 깔면 플롯이 패널 안으로 살짝 inset 된 box 를 돌려받아(패널이 차트를
    # 프레임) 이후 모든 차트 렌더러가 그 box 를 쓴다. 무발동이면 원본 box 그대로(무회귀).
    _renderable = has_data or bool(spec.get("points")) or bool(spec.get("values"))
    if _renderable:
        box = _maybe_chart_backdrop(slide, box, ctype, spec)

    # Composition grammar: a look can declare `chart_style: "brutal"` so bar/column
    # charts render as drawn shapes in the deck's visual language (thick black
    # border + hard shadow + value badges + 4px baseline + direct labels) instead
    # of a smooth native chart — the difference between "absorbed the skin" and
    # "absorbed the composition". Only bar/column families route here; other chart
    # types keep their (themed) native/composed renderers.
    _cstyle = GRAMMAR.get("chart_style")
    if _cstyle == "brutal":
        if has_data and ctype in _BRUTAL_BAR_TYPES:
            return _chart_brutal_bars(slide, box, spec)
        if has_data and ctype in ("line", "line_markers", "area"):
            return _chart_brutal_line(slide, box, spec)
        if ctype in ("xy_scatter", "bubble") and spec.get("points"):
            return _chart_brutal_scatter(slide, box, spec)
        if ctype in ("pie", "doughnut") and (has_data or spec.get("values")):
            return _chart_brutal_pie(slide, box, spec)
    elif _cstyle == "swiss":
        # 스위스 완성형: 막대·라인·산점도·share 전부 swiss 렌더러(무채색+스팟1·얇은 룰).
        if has_data and ctype in _BRUTAL_BAR_TYPES:
            return _chart_swiss_bars(slide, box, spec)
        if has_data and ctype in ("line", "line_markers", "area"):
            return _chart_swiss_line(slide, box, spec)
        if ctype in ("xy_scatter", "bubble") and spec.get("points"):
            return _chart_swiss_scatter(slide, box, spec)
        if ctype in ("pie", "doughnut") and (has_data or spec.get("values")):
            return _chart_swiss_pie(slide, box, spec)
    elif _cstyle == "luxe":
        # 다크 럭셔리: 헤어라인 아웃라인 막대·골드 라인·윤곽 산점도·비례 막대(캔버스-인지).
        if has_data and ctype in _BRUTAL_BAR_TYPES:
            return _chart_luxe_bars(slide, box, spec)
        if has_data and ctype in ("line", "line_markers", "area"):
            return _chart_luxe_line(slide, box, spec)
        if ctype in ("xy_scatter", "bubble") and spec.get("points"):
            return _chart_luxe_scatter(slide, box, spec)
        if ctype in ("pie", "doughnut") and (has_data or spec.get("values")):
            return _chart_luxe_share(slide, box, spec)
    elif _cstyle == "tech":
        # 다크 테크: 평면 네온 막대(강조 1개 시안 + 1px 상단 글로우 엣지, 나머지 중립
        # #3A3D45) + 발광 폴리라인(선 아래 페이드) + 모노 값/축 라벨. 막대/라인만 라우팅,
        # 그 외(산점도·share)는 themed native/composed 유지.
        if has_data and ctype in _BRUTAL_BAR_TYPES:
            return _chart_tech_bars(slide, box, spec)
        if has_data and ctype in ("line", "line_markers", "area"):
            return _chart_tech_line(slide, box, spec)
    elif _cstyle == "prism":
        # 프리즈매틱 다크: 프리즘 그라디언트 데이터 마크 — 막대(그라디언트 채움+윗면 글로우
        # 캡)·라인(좌→우 색전이 폴리라인+끝점 발광)·비중(프리즘 세그먼트 링). 모노 값 라벨.
        # 막대/라인/도넛 라우팅, 산점도는 themed native 유지.
        if has_data and ctype in _BRUTAL_BAR_TYPES:
            return _chart_prism_bars(slide, box, spec)
        if has_data and ctype in ("line", "line_markers", "area"):
            return _chart_prism_line(slide, box, spec)
        if ctype in ("pie", "doughnut") and (has_data or spec.get("values")):
            return _chart_prism_share(slide, box, spec)
    elif _cstyle == "cinematic":
        # 현대 시네마틱: 사이안 단톤 패밀리 절제 차트 — 막대(톤 스택, 강조 1개 시안 글로우)·
        # 라인(시안 라인+면적 채움+흰fill/시안ring 닷+선택 chevron-tab stage band)·비중(시안
        # 톤 비례 막대). 모노 라벨 없음(prompt.md 금지), 무지개 금지. 막대/라인/도넛 라우팅.
        if has_data and ctype in _BRUTAL_BAR_TYPES:
            return _chart_cinematic_bars(slide, box, spec)
        if has_data and ctype in ("line", "line_markers", "area"):
            return _chart_cinematic_line(slide, box, spec)
        if ctype in ("pie", "doughnut") and (has_data or spec.get("values")):
            return _chart_cinematic_share(slide, box, spec)

    if ctype in ("xy_scatter", "bubble"):
        return _chart_xy(slide, box, spec, ctype)
    composed = {
        "waterfall": _chart_waterfall, "funnel": _chart_funnel,
        "gauge": _chart_gauge, "progress": _chart_gauge,
        "gantt": _chart_gantt, "timeline": _chart_gantt,
        "tam_sam_som": _chart_tam_sam_som, "bullet": _chart_bullet,
        "slope": _chart_slope, "kpi_cards": _chart_kpi_cards,
    }
    if ctype in composed:
        return composed[ctype](slide, box, spec)
    if ctype in ("pie", "doughnut"):
        return _chart_pie(slide, box, spec, ctype)
    if not has_data:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    return _chart_native(slide, box, spec, ctype)


def _chart_native(slide, box, spec, ctype):
    """CategoryChartData 기반 native 차트. 시리즈색=팔레트 램프."""
    x, y, w, h = box
    xl = _NATIVE.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
    cd = CategoryChartData()
    cd.categories = spec.get("categories", [])
    for i, s in enumerate(spec.get("series", [])):
        cd.add_series(s.get("name", f"S{i+1}"), s.get("values", []),
                      number_format=spec.get("number_format"))
    gf = slide.shapes.add_chart(xl, IN(x), IN(y), IN(w), IN(h), cd)
    ch = gf.chart
    multi = len(spec.get("series", [])) > 1
    ch.has_legend = spec.get("legend", multi)
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = PT_PX(TYPO["caption"]["size_px"])
        _theme_legend(ch)
    # Series color = brand ramp. Bars/areas read .format.fill; lines read
    # .format.line; markers need their own fill/line. Set ALL so EVERY chart type
    # follows the active palette (the off-brand steel-blue/red line defaults were
    # Office's, not ours). Each in its own try so one unsupported facet never
    # aborts the rest.
    _gcfg = _glow_cfg()
    for i, ps in enumerate(ch.series):
        col = C(chart_color(i))
        try:
            ps.format.fill.solid()
            ps.format.fill.fore_color.rgb = col
        except Exception:
            pass
        try:
            ps.format.line.color.rgb = col
        except Exception:
            pass
        try:
            ps.marker.format.fill.solid()
            ps.marker.format.fill.fore_color.rgb = col
            ps.marker.format.line.color.rgb = col
        except Exception:
            pass
        # 글로우: 발광 룩 + chart 활성이면 시리즈 도형(막대/면/마커)에 같은 색조 헤일로를
        # c:spPr 에 주입(네이티브 차트 마크에도 네온감). 차트 part 의 spPr 이라 best-effort
        # — 실패해도 차트는 정상(무선언/라이트=무동작).
        if _gcfg is not None and _gcfg.get("chart", True):
            try:
                spPr = ps.format._element.find(qn("c:spPr"))
                if spPr is not None:
                    hexv = _safe_hex(_gcfg.get("color") or chart_color(i))
                    if hexv is not None:
                        eff = _effect_lst(spPr)
                        for old in eff.findall(qn("a:glow")):
                            eff.remove(old)          # never stack 2 glows (CT_EffectList: max 1)
                        g = eff.makeelement(qn("a:glow"),
                                            {"rad": str(int(_gcfg.get("rad", 18) * 12700))})
                        clr = g.makeelement(qn("a:srgbClr"), {"val": hexv})
                        a = clr.makeelement(qn("a:alpha"),
                                            {"val": str(_alpha_val(_gcfg.get("alpha", 60)))})
                        clr.append(a)
                        g.append(clr)
                        eff.insert(0, g)
                        _sort_effects(eff)
            except Exception:
                pass
    try:
        ch.category_axis.tick_labels.font.size = PT_PX(13)
        va = ch.value_axis
        va.tick_labels.font.size = PT_PX(13)
        va.has_major_gridlines = True
    except Exception:
        pass
    # 단일 시리즈는 PowerPoint가 시리즈명을 자동 차트제목으로 띄우는데(명시 title 요소가
    # 아니라 python-pptx로 직접 색 지정 불가) 다크 캔버스에선 어두워 묻힌다 → 다크일 때만
    # 명시 제목(시리즈명)으로 만들어 ink로 칠한다. 라이트는 auto 유지(무회귀).
    if DARK:
        _series = spec.get("series", [])
        if len(_series) == 1 and _series[0].get("name"):
            try:
                ch.has_title = True
                ch.chart_title.text_frame.text = _series[0]["name"]
                for _p in ch.chart_title.text_frame.paragraphs:
                    _p.font.color.rgb = C("ink")
                    for _r in _p.runs:
                        _r.font.color.rgb = C("ink")
            except Exception:
                pass
    _theme_axes(ch)
    return ch


def _chart_pie(slide, box, spec, ctype):
    x, y, w, h = box
    cats = spec.get("categories", [])
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    cd = CategoryChartData()
    cd.categories = cats or [str(i + 1) for i in range(len(vals))]
    cd.add_series("", vals)
    xl = XL_CHART_TYPE.DOUGHNUT if ctype == "doughnut" else XL_CHART_TYPE.PIE
    gf = slide.shapes.add_chart(xl, IN(x), IN(y), IN(w), IN(h), cd)
    ch = gf.chart
    # Respect spec legend: callers that draw their OWN legend (e.g. build_share's
    # color-chip list) pass legend=False to avoid a duplicate native legend. (bug B)
    ch.has_legend = spec.get("legend", True)
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.RIGHT
        ch.legend.include_in_layout = False
        ch.legend.font.size = PT_PX(14)
        _theme_legend(ch)
    try:
        pts = ch.plots[0].series[0].points
        for i, p in enumerate(pts):
            p.format.fill.solid()
            p.format.fill.fore_color.rgb = C(chart_color(i))
    except Exception:
        pass
    return ch


def _chart_xy(slide, box, spec, ctype):
    """스캐터/버블. XyChartData / BubbleChartData 사용.
    spec.points=[{x,y,(size),label}] 또는 series 형식."""
    x, y, w, h = box
    pts = spec.get("points", [])
    if not pts:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    if ctype == "bubble":
        from pptx.chart.data import BubbleChartData
        cd = BubbleChartData()
        sr = cd.add_series("")
        for p in pts:
            sr.add_data_point(p.get("x", 0), p.get("y", 0), p.get("size", 1))
        xl = XL_CHART_TYPE.BUBBLE
    else:
        from pptx.chart.data import XyChartData
        cd = XyChartData()
        sr = cd.add_series("")
        for p in pts:
            sr.add_data_point(p.get("x", 0), p.get("y", 0))
        xl = XL_CHART_TYPE.XY_SCATTER
    gf = slide.shapes.add_chart(xl, IN(x), IN(y), IN(w), IN(h), cd)
    ch = gf.chart
    ch.has_legend = False
    try:
        ch.series[0].format.fill.solid()
        ch.series[0].format.fill.fore_color.rgb = C("blue")
    except Exception:
        pass
    return ch


# ----- 도형 합성 차트 (native 로 안 되는 8종, box 에서 결정론적으로 파생) -----
def _series_values(spec):
    """첫 시리즈 (cats, vals) 추출."""
    cats = spec.get("categories", [])
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    return cats, (vals or [])


def _chart_waterfall(slide, box, spec):
    """워터폴: 시작값 + 증감 막대 누적. spec.values 의 각 값이 증감,
    첫/마지막은 total 로 본다(spec.totals=[idx] 로 표시 가능)."""
    x, y, w, h = box
    cats, vals = _series_values(spec)
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    totals = set(spec.get("totals", [0, len(vals) - 1]))
    n = len(vals)
    bw = (w - GAP * (n - 1)) / n
    # 누적 범위 계산
    cum = 0
    lows, run = [], []
    for i, v in enumerate(vals):
        if i in totals:
            run.append((0, v)); cum = v
        else:
            start = cum; cum += v
            run.append((min(start, cum), max(start, cum)))
    vmax = max(hi for _, hi in run) or 1
    base_y = y + h - 28
    plot_h = h - 40
    for i, (lo, hi) in enumerate(run):
        bx = x + i * (bw + GAP)
        bh = (hi - lo) / vmax * plot_h
        by = base_y - (hi / vmax * plot_h)
        col = "blue" if i in totals else ("green" if vals[i] >= 0 else "orange")
        rect(slide, [bx, by, bw, max(2, bh)], fill=col, radius=RAD["sm"])
        if i < len(cats):
            text(slide, [bx - 4, base_y + 4, bw + 8, 24], str(cats[i]), "caption",
                 color="muted", align="center")


def _chart_funnel(slide, box, spec):
    """퍼널: 위→아래로 좁아지는 사다리꼴 막대(값 비례 폭)."""
    x, y, w, h = box
    cats, vals = _series_values(spec)
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    vmax = max(vals) or 1
    n = len(vals)
    rh = (h - 12 * (n - 1)) / n
    for i, v in enumerate(vals):
        bw = v / vmax * w
        bx = x + (w - bw) / 2
        by = y + i * (rh + 12)
        rect(slide, [bx, by, bw, rh], fill=chart_color(i),
             radius=RAD["sm"])
        lbl = (str(cats[i]) + "  " if i < len(cats) else "") + str(v)
        text(slide, [x, by + (rh - 20) / 2, w, 24], lbl, "body-md",
             color="on-accent", align="center")


def _chart_gauge(slide, box, spec):
    """게이지/진행률: 트랙 + 채움 막대 + 중앙 %텍스트. spec.percent(0~100)."""
    x, y, w, h = box
    pct = float(spec.get("percent", spec.get("value", 0)) or 0)
    pct = max(0, min(100, pct))
    # 도넛 게이지 대신 명시적 수평 진행바(읽기 쉬움)
    bar_h = min(48, h * 0.18)
    by = y + h / 2 + 10
    rect(slide, [x, by, w, bar_h], fill="band", radius=RAD["pill"])
    if pct > 0:
        rect(slide, [x, by, w * pct / 100, bar_h], fill="blue", radius=RAD["pill"])
    text(slide, [x, y + h / 2 - 90, w, 100], f"{pct:.0f}%", "display",
         color="blue", align="center")
    if spec.get("caption"):
        text(slide, [x, by + bar_h + 10, w, 40], spec["caption"], "body-lg",
             color="muted", align="center")


def _chart_gantt(slide, box, spec):
    """간트/타임라인: 행별 막대(start,end 0~1 또는 일수). spec.tasks=[{label,start,end}]."""
    x, y, w, h = box
    tasks = spec.get("tasks", [])
    if not tasks:
        cats, vals = _series_values(spec)
        tasks = [{"label": c, "start": 0, "end": v} for c, v in zip(cats, vals)]
    if not tasks:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    mx = max((t.get("end", 1) for t in tasks), default=1) or 1
    label_w = w * 0.25
    track_x = x + label_w
    track_w = w - label_w
    n = len(tasks)
    rh = (h - 10 * (n - 1)) / n
    for i, t in enumerate(tasks):
        ry = y + i * (rh + 10)
        text(slide, [x, ry, label_w - 12, rh], t.get("label", ""), "body-md",
             color="ink", anchor="middle")
        s = t.get("start", 0) / mx
        e = t.get("end", 1) / mx
        rect(slide, [track_x + s * track_w, ry + rh * 0.2,
                     max(4, (e - s) * track_w), rh * 0.6],
             fill=chart_color(i), radius=RAD["sm"])


def _chart_tam_sam_som(slide, box, spec):
    """TAM-SAM-SOM: 동심원 3겹 + 라벨. spec.bubbles=[{label,value}] (큰→작은)."""
    x, y, w, h = box
    items = spec.get("bubbles") or spec.get("layers") or []
    if not items:
        cats, vals = _series_values(spec)
        items = [{"label": c, "value": v} for c, v in zip(cats, vals)]
    items = items[:3]
    if not items:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    d = min(w * 0.55, h)
    cx = x + d / 2
    cy = y + h / 2
    fills = ["blue-pale", "blue-light", "blue"]
    sizes = [d, d * 0.66, d * 0.36]
    for i in range(len(items)):
        sz = sizes[i]
        oval(slide, [cx - sz / 2, cy - sz / 2, sz, sz], fill=fills[i])
    lx = x + d + 40
    lh = h / max(1, len(items))
    for i, it in enumerate(items):
        ly = y + i * lh
        text(slide, [lx, ly + 8, w - d - 60, 36],
             it.get("label", ""), "h4",
             color="blue" if fills[i] != "blue-pale" else "ink")
        if it.get("value") is not None:
            text(slide, [lx, ly + 50, w - d - 60, 44], str(it["value"]), "h2",
                 color="ink")


def _chart_bullet(slide, box, spec):
    """불릿 차트: 행별 (실측 막대 vs 목표 마커 + 정성 밴드). spec.rows=[{label,value,target,max}]."""
    x, y, w, h = box
    rows = spec.get("rows", [])
    if not rows:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    label_w = w * 0.22
    track_x = x + label_w
    track_w = w - label_w
    n = len(rows)
    rh = (h - 16 * (n - 1)) / n
    for i, r in enumerate(rows):
        ry = y + i * (rh + 16)
        mx = r.get("max", max(r.get("value", 0), r.get("target", 0)) * 1.25) or 1
        text(slide, [x, ry, label_w - 12, rh], r.get("label", ""), "body-md",
             color="ink", anchor="middle")
        # 정성 밴드(연한 배경)
        rect(slide, [track_x, ry + rh * 0.25, track_w, rh * 0.5],
             fill="band", radius=0)
        val = r.get("value", 0) / mx
        rect(slide, [track_x, ry + rh * 0.35, val * track_w, rh * 0.3],
             fill="blue", radius=0)
        if r.get("target") is not None:
            tx = track_x + r["target"] / mx * track_w
            rect(slide, [tx - 2, ry + rh * 0.2, 4, rh * 0.6], fill="ink")


def _chart_slope(slide, box, spec):
    """슬로프(기울기) 차트: 좌→우 두 시점 값 연결선. spec.series=[{name,values:[a,b]}]."""
    x, y, w, h = box
    series = spec.get("series", [])
    if not series:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    allv = [v for s in series for v in s.get("values", [])[:2]]
    if not allv:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    vmin, vmax = min(allv), max(allv)
    span = (vmax - vmin) or 1
    lx, rx = x + w * 0.18, x + w * 0.82
    top, bot = y + 20, y + h - 30

    def py(v):
        return bot - (v - vmin) / span * (bot - top)
    cats = spec.get("categories", ["전", "후"])
    text(slide, [x, y + h - 26, w * 0.4, 24], str(cats[0] if cats else ""),
         "caption", color="muted", align="left")
    text(slide, [x + w * 0.6, y + h - 26, w * 0.4, 24],
         str(cats[1] if len(cats) > 1 else ""), "caption",
         color="muted", align="right")
    for i, s in enumerate(series):
        vals = s.get("values", [])
        if len(vals) < 2:
            continue
        col = chart_color(i)
        line(slide, lx, py(vals[0]), rx, py(vals[1]), color=col, w=2.5)
        oval(slide, [lx - 6, py(vals[0]) - 6, 12, 12], fill=col)
        oval(slide, [rx - 6, py(vals[1]) - 6, 12, 12], fill=col)
        text(slide, [x, py(vals[0]) - 14, w * 0.16, 28],
             f"{s.get('name','')} {vals[0]}", "caption", color="ink",
             align="right")
        text(slide, [rx + 8, py(vals[1]) - 14, w * 0.16, 28], str(vals[1]),
             "caption", color="ink")


def _chart_kpi_cards(slide, box, spec):
    """빅넘버 KPI 카드 N개(가로 분할). spec.kpis=[{label,value,note}]."""
    kpis = spec.get("kpis") or spec.get("stats") or []
    if not kpis:
        return _media_placeholder(slide, box, "데이터 없음")
    boxes = _split_cols(box, len(kpis))
    accents = accent_cycle(len(kpis))
    # KPI 카드는 자유 블록(수치가 카드 위치가 아니라 카드 '안'에 있음) → 틸트 대상.
    # block_fill 문법이면 카드를 원색으로 채우고 텍스트색을 휘도로 플립.
    filled = bool(GRAMMAR.get("block_fill"))
    # KPI value color: a look may render the big number in an accent via
    # grammar.kpi_value_accent. Two forms (no grammar → plain ink, unchanged for
    # every other look; only the accent-edge / non-filled variant):
    #   • True  → the dark-tech '44pt 시안 KPI 숫자' spec: each card's CYCLING accent
    #             (card 0 = accent[0], card 1 = accent[1], …).
    #   • "#hex"/token → a FIXED color for every big number (vivid-gradient-future:
    #             all key metrics in glow-cyan #7DF9FF, per its preview — distinct
    #             from prismatic's per-card gradient numbers).
    _kva = GRAMMAR.get("kpi_value_accent")
    val_accent = bool(_kva) and not filled
    val_fixed = _kva if (val_accent and isinstance(_kva, str)) else None
    # 프리즘: KPI 빅넘버를 프리즘 그라디언트 텍스트 채움으로(스펙 #1 식별 단서). 카드는
    # 글래스 패널(card 의 prism/glass 처리). 채움 카드(block_fill)면 그라디언트 텍스트
    # 불가(대비) → 기존 _ink_on 유지.
    prism_val = (_prism_cfg() is not None) and not filled
    for i, k in enumerate(kpis):
        bx, by, bw, bh = boxes[i]
        n0 = len(slide.shapes)
        acc = accents[i % len(accents)]
        on = _ink_on(acc) if filled else None
        if filled:
            card(slide, boxes[i], "card", fill=acc)
        elif prism_val:
            # 프리즘 KPI = 다크 글래스 노드(거의 투명 + 그라디언트 헤어라인 보더). accent
            # 엣지 대신 노드 처리 — card() 의 _is_node 경로가 투명채움+그라디언트 스트로크.
            card(slide, boxes[i], "card", fill="surface-2")
        else:
            card(slide, boxes[i], "card", accent=acc)
        text(slide, [bx + PAD, by + bh * 0.18, bw - 2 * PAD, 36],
             k.get("label", ""), "label", color=on if filled else "muted",
             align="center", upper=True)
        if prism_val:
            gradient_text(slide, [bx + PAD, by + bh * 0.38, bw - 2 * PAD, bh * 0.34],
                          k.get("value", ""), "h1", align="center")
        else:
            text(slide, [bx + PAD, by + bh * 0.38, bw - 2 * PAD, bh * 0.34],
                 k.get("value", ""), "h1",
                 color=on if filled
                 else ((val_fixed or acc) if val_accent else "ink"),
                 align="center")
        if k.get("note"):
            text(slide, [bx + PAD, by + bh * 0.74, bw - 2 * PAD, bh * 0.2],
                 k["note"], "body", color=on if filled else "muted",
                 align="center")
        _tilt_group(slide, n0, bx + bw / 2.0, by + bh / 2.0)


# ---------------------------------------------------------------------------
# 카드 콘텐츠 헬퍼 (Set/Contrast 공용) — 이미지 카드 변종 지원.
# ---------------------------------------------------------------------------
def _content_card(slide, box, item, accent="blue"):
    """카드: (선택적 이미지 밴드 상단 ~55%) + label/heading/body 하단.
    item={image:{src,...}, label, heading, body|bullets}. 이미지 있으면 이미지 카드.
    룩이 components.card.tilt_deg 를 선언하면 카드(패널+accent+텍스트) 전체를 카드
    중심 기준으로 결정론적 소각도 회전 → 네오브루탈리즘 '스티커' 느낌(무선언=무회전)."""
    x, y, w, h = box
    n0 = len(slide.shapes)
    # 'image' 키가 있으면(빈 dict {} 라도) 이미지 카드 — src 없으면 media()가
    # 플레이스홀더를 그린다. ({} 는 falsy 이므로 get() 대신 키 존재로 판정.)
    # fill/border 는 활성 style/look 의 components.card 프리셋이 결정한다
    # (하드코딩 금지). accent 엣지는 이미지 카드엔 그리지 않는다(이미지가 상단을
    # 덮으므로) — 텍스트 카드일 때만 프리셋의 accent_edge 로 그린다.
    if "image" in item:
        card(slide, box, "card")
        img = item.get("image") or {}
        band_h = h * 0.55
        media(slide, [x, y, w, band_h], img, fit=img.get("fit", "cover"))
        ty = y + band_h + 16
        th = h - band_h - 24
    else:
        # `active:true` marks the SELECTED cell. Under a neumorphic look it pops OUT
        # as a convex accent-filled card (no-op for every other look, which has no
        # elevation:neumorphic, so `active` is ignored downstream).
        card(slide, box, "card", accent=accent, active=bool(item.get("active")))
        ty = y + 24
        th = h - 40
    # Active neumorphic cell is filled in the accent → its text must read on-accent.
    neu_active = bool(GRAMMAR.get("neumorphic")) and bool(item.get("active")) \
        and "image" not in item
    inx = x + PAD
    inw = w - 2 * PAD
    cy = ty
    # Hyundai card enumerator: a GIANT cyan SOLID number (01, 02, …) at the top of
    # each text card — the pack's cards anchor (p-64) and the ONLY place a big cyan
    # numeral appears (outline strictly banned). Counts per slide. No-op without the
    # grammar key OR on image cards (the image owns the top band).
    if GRAMMAR.get("card_enumerator") and "image" not in item:
        global _CARD_ENUM
        _CARD_ENUM += 1
        enum_h = min(120, h * 0.34)
        text(slide, [inx, cy - 6, inw, enum_h], f"{_CARD_ENUM:02d}", "display",
             color="blue", align="center", anchor="top")
        cy += enum_h + 10
    _enum = bool(GRAMMAR.get("card_enumerator")) and "image" not in item
    _al = "center" if _enum else "left"
    if item.get("label"):
        text(slide, [inx, cy, inw, 26], item["label"], "label",
             color=("on-accent" if neu_active else accent), upper=True, align=_al)
        cy += 32
    if item.get("heading"):
        text(slide, [inx, cy, inw, 44], item["heading"], "h3",
             color=("on-accent" if neu_active else "ink"), align=_al)
        cy += 56
    rest = max(40, (ty + th) - cy)
    if item.get("bullets"):
        bullets(slide, [inx, cy, inw, rest], item["bullets"], "body",
                color=("on-accent" if neu_active else "body"),
                marker_color=("on-accent" if neu_active else accent))
    elif item.get("body"):
        text(slide, [inx, cy, inw, rest], item["body"], "body",
             color=("on-accent" if neu_active else "body"), align=_al)
    _tilt_group(slide, n0, x + w / 2.0, y + h / 2.0)


# ===========================================================================
# Family 1 — Frame
# ===========================================================================
def _cover_title_card(s, d, L):
    """Brutalism cover: the headline lives INSIDE a flat bordered block tilted a
    few degrees (the pack's "흰 틸트 타이틀 박스"), with the subtitle in an accent
    block colliding at the opposite tilt. Generalizes the pack's cover composition
    (which our skin-only absorb flattened to plain text on the canvas). The block
    is measured to HUG the text so it reads as a sticker, not a banner."""
    up = "display" in (GRAMMAR.get("upper_tiers") or ())
    lat = ACTIVE_FONT["latin_display"]
    ea = ACTIVE_FONT["ea_display"]
    pad = 28
    tx, ty, tw, _th = L["title"]
    inw = tw - 2 * pad
    title = d.get("title", "")
    lines = _prewrap(title, inw, TYPO["display"]["size_px"], True, lat, ea, upper=up)
    size_px = TYPO["display"]["size_px"]
    lh_px = size_px * TYPO["display"]["line"]
    linew = max((_measure(ln.upper() if up else ln, lat, ea, True, int(size_px))
                 for ln in lines), default=inw)
    card_w = min(tw, linew + 2 * pad)
    card_h = len(lines) * lh_px + 2 * pad
    # title block (white card + ink headline), tilted as one rigid unit
    n0 = len(s.shapes)
    card(s, [tx - pad, ty, card_w, card_h], "card", fill="surface")
    text(s, [tx, ty + pad, card_w - 2 * pad, card_h - 2 * pad], lines, "display",
         color="ink", anchor="top")
    cx, cy = int(IN(tx - pad + card_w / 2.0)), int(IN(ty + card_h / 2.0))
    for sp in list(s.shapes)[n0:]:
        _rotate_about(sp, -2.0, cx, cy)
    # subtitle in an accent block, colliding at the opposite tilt
    if d.get("subtitle"):
        sx, sy, sw, _sh = L["subtitle"]
        acc = ACCENTS[0] if ACCENTS else "blue"
        sub_w = min(sw, _measure(d["subtitle"], lat, ea, True,
                                 int(TYPO["h4"]["size_px"])) + 2 * pad)
        sub_h = TYPO["h4"]["size_px"] * TYPO["h4"]["line"] + 2 * pad
        n1 = len(s.shapes)
        blk = card(s, [sx, sy, sub_w, sub_h], "card", fill=acc)
        shape_text(blk, d["subtitle"], "h4", color=_ink_on(acc))
        scx, scy = int(IN(sx + sub_w / 2.0)), int(IN(sy + sub_h / 2.0))
        for sp in list(s.shapes)[n1:]:
            _rotate_about(sp, 1.5, scx, scy)


def _cover_glass_card(s, d, L):
    """Glassmorphism cover (grammar.cover == 'glass_card'): kicker + title + subtitle +
    meta sit inside ONE COMPACT frosted glass panel (upper-left) that HUGS the text, so the
    cover carries the look's second identity pillar (frosted glass), not just the mesh (the
    QA panel flagged the bare-text cover). The panel is a normal card() → the look's glass
    alpha auto-frosts it (mesh shows through) and the white card border is the highlight rim.
    No tilt (that is the brutalist title_card). Heights are measured from TYPO so the panel
    hugs the real text block (no full-height empty banner)."""
    up = "display" in (GRAMMAR.get("upper_tiers") or ())
    lat = ACTIVE_FONT["latin_display"]
    ea = ACTIVE_FONT["ea_display"]
    pad, gap = 40, 16
    tx, ty, tw, _th = L["title"]
    inw = tw - 2 * pad
    title = d.get("title", "")
    dsize = TYPO["display"]["size_px"]
    dlh = dsize * TYPO["display"]["line"]
    lines = _prewrap(title, inw, dsize, True, lat, ea, upper=up)
    title_h = max(1, len(lines)) * dlh
    has_kick = bool(d.get("kicker") or d.get("eyebrow"))
    has_sub = bool(d.get("subtitle"))
    has_meta = bool(d.get("meta"))
    kh_ = TYPO["caption"]["size_px"] * 1.7 if has_kick else 0
    sh_ = TYPO["h4"]["size_px"] * TYPO["h4"]["line"] if has_sub else 0
    mh_ = TYPO["caption"]["size_px"] * TYPO["caption"]["line"] if has_meta else 0
    inner = (kh_ + gap if has_kick else 0) + title_h \
        + (gap + sh_ if has_sub else 0) + (gap + mh_ if has_meta else 0)
    px = tx - pad
    py = ty - pad - (kh_ + gap if has_kick else 0)
    pw, ph = tw, inner + 2 * pad
    # The reference COVER panel is a DARK frosted glass anchor (light text), even though the
    # body cards are light frost — the cover's dark plate gives the deck its authoritative
    # contrast. So draw the cover panel with a dark translucent fill + light highlight rim +
    # light text, tunable via grammar.cover_card. (Body cards stay light via glass_alpha.)
    cc = GRAMMAR.get("cover_card") or {}
    rad = (COMPONENTS.get("card") or {}).get("radius", 20)
    grad = cc.get("gradient")
    if grad:
        # Gradient-block cover (vivid-gradient infographic): an OPAQUE vivid gradient panel
        # (purple→coral etc.) holding white title text — the look's #1 cue.
        sp = rect(s, [px, py, pw, ph], fill=_resolve(grad[0]),
                  line_color=cc.get("border"), line_w=cc.get("border_pt", 0) or 0, radius=rad)
        _gradient_shape_fill(sp, [_resolve(c) for c in grad], ang=int(cc.get("ang", 2700000)))
    else:
        # Frosted-glass cover (glassmorphism): translucent dark plate over the mesh.
        sp = rect(s, [px, py, pw, ph], fill=cc.get("fill", "navy"),
                  line_color=cc.get("border", "#FFFFFF"), line_w=cc.get("border_pt", 1.0),
                  radius=rad)
        _apply_alpha(sp, int(cc.get("alpha", 56)))
    p_ink, p_sub, p_kick = (cc.get("ink", "white"), cc.get("sub", "blue-pale"),
                            cc.get("kicker", "blue-light"))
    cy = py + pad
    if has_kick:
        text(s, [tx, cy, inw, kh_], _mono_kicker_prefix(d.get("kicker") or d["eyebrow"]),
             "eyebrow", color=_mono_kicker_color(p_kick), upper=True)
        cy += kh_ + gap
    text(s, [tx, cy, inw, title_h], lines, "display", color=p_ink, anchor="top")
    cy += title_h
    if has_sub:
        cy += gap
        text(s, [tx, cy, inw, sh_], d["subtitle"], "h4", color=p_sub)
        cy += sh_
    if has_meta:
        cy += gap
        text(s, [tx, cy, inw, mh_], d["meta"], "caption", color=p_sub)
    return s


def _cover_color_block(s, d, L):
    """Confident-color-block cover (grammar.cover == 'color_block'): a big LEFT
    full-height SOLID color COLUMN (ACCENTS[0]) carrying the headline in WHITE,
    top-weighted and left-aligned — the pack's #1 identity cue (캔버스 ~43% 비대칭
    색면 + 그 위 화이트 헤드라인, 직각 0 radius). The right half stays the white
    canvas for subtitle/meta. No-op for every other look (gated in build_cover)."""
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    M = SLIDE["margin_px"]
    acc = ACCENTS[0] if ACCENTS else "blue"
    col_w = int(W * 0.43)
    rect(s, [0, 0, col_w, H], fill=acc, radius=0)
    on = _ink_on(acc)                       # white on a saturated block
    inx = L["kicker"][0]                     # 96 — align with the default slot x
    inw = col_w - inx - 56                   # right inner pad inside the column
    if d.get("kicker") or d.get("eyebrow"):
        text(s, [inx, L["kicker"][1], inw, L["kicker"][3]],
             d.get("kicker") or d["eyebrow"], "eyebrow", color=on, upper=True)
    text(s, [inx, L["title"][1], inw, H - L["title"][1] - 120],
         d.get("title", ""), "display", color=on, anchor="top")
    rx = col_w + M
    rw = W - rx - M
    if d.get("subtitle"):
        text(s, [rx, L["subtitle"][1], rw, L["subtitle"][3]], d["subtitle"],
             "h4", color="muted")
    if d.get("meta"):
        text(s, [rx, H - M - L["meta"][3], rw, L["meta"][3]], d["meta"],
             "caption", color="muted")
    return s


def _accent_for_index(index):
    """Map a section index ('01','2',…) to an ACCENTS slot so consecutive color-block
    dividers cycle colors; 1-based index → 0-based slot. Non-numeric → slot 0."""
    n = len(ACCENTS) or 1
    try:
        digits = "".join(ch for ch in str(index) if ch.isdigit())
        return (int(digits) - 1) % n
    except Exception:
        return 0


def _section_color_block(s, d, L):
    """Confident-color-block divider (grammar.cover == 'color_block'): a full-bleed
    SOLID accent block (cycled by the section index for variety) carrying the index
    and section name in WHITE — the pack's '디바이더 = 화면을 채우는 거대 색블록 + 화이트
    섹션 번호'. No-op for every other look (gated in build_section)."""
    acc = ACCENTS[_accent_for_index(d.get("index"))] if ACCENTS else "blue"
    slide_bg(s, acc)
    on = _ink_on(acc)
    # Giant section number (grammar.section_watermark) draws ON the block — no-op
    # unless the look declares it. When present it IS the section number, so the
    # small index label is suppressed to avoid a duplicate (bold-block-infographic's
    # 200pt 섹션 번호 디바이더). When absent (confident-color-block) the small index
    # stands as before.
    wm = _section_watermark(s, d.get("index"))
    if d.get("index") and wm is None:
        text(s, L["index"], str(d["index"]), "h1", color=on)
    text(s, L["section_title"], d.get("title", d.get("section_title", "")),
         "h1", color=on, anchor="top")
    if d.get("caption"):
        text(s, L["caption"], d["caption"], "body-lg", color=on)
    return s


def _statement_color_block(s, d):
    """Confident-color-block statement (grammar.cover == 'color_block'): the whole
    slide is ONE solid accent block with the lead set big in WHITE — the pack's
    confident 색면 스테이트먼트. No header chrome; the statement IS the message.
    No-op for every other look (gated in build_statement)."""
    L = LAYOUTS["archetypes"]["statement"]["slots"]
    acc = ACCENTS[0] if ACCENTS else "blue"
    slide_bg(s, acc)
    on = _ink_on(acc)
    lx, ly, lw, lh = L["lead"]
    text(s, [lx, ly, lw, lh], d.get("lead", ""), "h1", color=on, anchor="middle")
    if d.get("support"):
        sx, sy, sw, sh = L["support"]
        if isinstance(d["support"], list):
            bullets(s, [sx, sy, sw, sh], d["support"], "body-lg",
                    color=on, marker_color=on)
        else:
            text(s, [sx, sy, sw, sh], d["support"], "body-lg", color=on)
    return s


def build_cover(prs, d):
    s = new_slide(prs)
    A = LAYOUTS["archetypes"]["cover"]
    L = A["slots"]
    # Cover background precedence: an explicit per-slide spec `bg` wins; otherwise
    # the look's declared `cover_bg` (a dark title slide over a light body — navy
    # IR decks, full-bleed keynotes); otherwise the body `canvas` (derived, so a
    # light look gets a light cover and a dark look a dark one). The forced literal
    # used to live in the demo spec, which painted EVERY look's cover dark.
    bg = d.get("bg") or _LOOK_COVER_BG or "canvas"
    slide_bg(s, bg)
    _orb_decor(s)   # 소프트 3D 구 데코(grammar.orb_decor); 배경 위·콘텐츠 아래. 무선언=무동작.
    # dark = 실제 배경 휘도로 판정(이름목록 대신) → 컬러 캔버스(예: 브루탈리즘 옐로)는
    # 밝으므로 어두운 잉크, 네이비/블랙은 흰 잉크. 흰 배경은 그대로 light.
    dark = is_dark(_resolve(bg))
    # Hyundai cover: the big LEFT 5-stack chevron + eye glow (bookend with closing).
    # No-op without grammar.chevron. Drawn after bg, before the title (right side).
    if not d.get("background"):
        _chevron_cover(s, "left")
    if d.get("background"):
        media(s, [0, 0, SLIDE["width_px"], SLIDE["height_px"]], d["background"],
              fit="cover")
    ink = "white" if dark else "navy"
    sub = "slate-300" if dark else "slate-500"
    # Glassmorphism cover: wrap kicker/title/subtitle/meta in one frosted glass panel so
    # the cover carries the look's second identity pillar (frosted glass), not just the
    # mesh. No-op for every other look (gated on grammar.cover == 'glass_card'); a full-
    # bleed bg image still owns the composition.
    if GRAMMAR.get("cover") == "glass_card" and not d.get("background"):
        return _cover_glass_card(s, d, L)
    if GRAMMAR.get("cover") == "color_block" and not d.get("background"):
        return _cover_color_block(s, d, L)
    # Hyundai cover: title block lives on the RIGHT half (chevron owns the left),
    # right-aligned WHITE (the pack: 표지·종료만 영문 타이틀 흰색). The chevron grammar
    # reuses the slot x to push the text into the right 55-95% band. No-op otherwise.
    chev = _chevron_cfg() is not None and not d.get("background")
    if chev:
        W = SLIDE["width_px"]
        rx = W * 0.50
        rw = W * 0.93 - rx
        if d.get("kicker") or d.get("eyebrow"):
            text(s, [rx, L["kicker"][1], rw, L["kicker"][3]],
                 d.get("kicker") or d["eyebrow"], "eyebrow",
                 color="blue-light", align="right", upper=True)
        text(s, [rx, L["title"][1], rw, L["title"][3]], d.get("title", ""),
             "display", color="white", anchor="top", align="right")
        if d.get("subtitle"):
            text(s, [rx, L["subtitle"][1], rw, L["subtitle"][3]], d["subtitle"],
                 "h4", color=sub, align="right")
        if d.get("meta"):
            text(s, [rx, L["meta"][1], rw, L["meta"][3]], d["meta"], "caption",
                 color="slate-300", align="right")
        return s
    if d.get("kicker") or d.get("eyebrow"):
        text(s, L["kicker"], _mono_kicker_prefix(d.get("kicker") or d["eyebrow"]),
             "eyebrow", color=_mono_kicker_color("blue-light" if dark else "blue"),
             upper=True)
    # Grammar cover variant: headline-in-a-block. Only when the look declares it
    # AND there is no full-bleed background image (which owns the composition).
    if GRAMMAR.get("cover") == "title_card" and not d.get("background"):
        _cover_title_card(s, d, L)
        if d.get("meta"):
            text(s, L["meta"], d["meta"], "caption", color=ink)
        return s
    text(s, L["title"], d.get("title", ""), "display", color=ink, anchor="top")
    if d.get("subtitle"):
        text(s, L["subtitle"], d["subtitle"], "h4", color=sub)
    if d.get("meta"):
        text(s, L["meta"], d["meta"], "caption", color="slate-400")
    return s


def _section_watermark(slide, index):
    """Giant low-contrast MONO section number behind a divider — the dark-tech
    '02' 120pt watermark. Drawn right after the bg so the index/title/caption sit
    on top. Gated on grammar.section_watermark (any canvas); no grammar → no-op.

    grammar.section_watermark (dict, all optional):
      size  — glyph size in pt (default 120)
      color — watermark color token (default 'divider' = the dim #2A2D35 hairline hue)
      box   — [x,y,w,h] px override (default top-right, bleeding off the right edge)"""
    cfg = GRAMMAR.get("section_watermark")
    if not cfg or index in (None, ""):
        return
    if cfg is True:
        cfg = {}
    size_pt = cfg.get("size", 120)
    # Default watermark hue follows the canvas: dim 'divider' on a dark deck (the
    # original dark-tech behavior — byte-identical), a faint 'line' tint on a light
    # deck so a light look (bold-block-infographic's 200pt 섹션 번호) can opt in.
    col = cfg.get("color") or ("divider" if DARK else "line")
    W, H = SLIDE["width_px"], SLIDE["height_px"]
    box = cfg.get("box", [W - 640, 60, 620, 280])
    tb = slide.shapes.add_textbox(IN(box[0]), IN(box[1]), IN(box[2]), IN(box[3]))
    tf = tb.text_frame
    tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(index)
    r.font.size = Pt(size_pt)
    r.font.bold = True
    r.font.color.rgb = C(col)
    mf, _ = _mono_face(False)
    _set_run_font(r, mf, mf)
    # 프리즘: 거대 섹션 번호를 프리즘 그라디언트 텍스트로(스펙 "중앙 거대 그라디언트 텍스트
    # 섹션 번호 120pt"). cfg.gradient False 면 기존 dim 워터마크 유지(다른 다크 룩 무회귀).
    if _prism_cfg() is not None and cfg.get("gradient", True):
        _gradient_text_run(r)
    return tb


def build_section(prs, d):
    s = new_slide(prs)
    L = LAYOUTS["archetypes"]["section"]["slots"]
    if GRAMMAR.get("cover") == "color_block" and not d.get("background"):
        return _section_color_block(s, d, L)
    # Divider chrome follows the look like the cover does: spec `bg` → look
    # `cover_bg` (dark chapter break over a light body) → `canvas` (derived).
    # A light look gets a light divider, not the old hard-coded navy.
    bg = d.get("bg") or _LOOK_COVER_BG or "canvas"
    slide_bg(s, bg)
    _orb_decor(s)   # 소프트 3D 구 데코(grammar.orb_decor); 무선언=무동작.
    dark = is_dark(_resolve(bg))
    # dark-tech divider watermark: huge dim mono section number behind the title.
    # Behind everything (drawn before index/title), no-op without the grammar key.
    _section_watermark(s, d.get("index"))
    if d.get("background"):
        media(s, [0, 0, SLIDE["width_px"], SLIDE["height_px"]], d["background"],
              fit="cover")
    ink = "white" if dark else "navy"
    # Hyundai chapter slate: RIGHT-edge 4-stack chevron (+ left starburst for multi-
    # word titles) + GIANT CYAN title (the pack: 챕터 슬레이트는 #66FFFF 사이안 — 표지/종료만
    # 흰색). No-op without grammar.chevron / on a bg image.
    chev = _chevron_cfg() is not None and not d.get("background")
    title_str = d.get("title", d.get("section_title", ""))
    if chev:
        _chevron_section(s, multiword=(len(str(title_str).split()) > 1))
    if d.get("index"):
        text(s, L["index"], str(d["index"]), "h1",
             color="blue-light" if dark else "blue")
    text(s, L["section_title"], title_str, "h1",
         color=("blue" if chev else ink), anchor="top")
    if d.get("caption"):
        text(s, L["caption"], d["caption"], "body-lg",
             color=("blue-light" if chev else ("slate-300" if dark else "slate-500")))
    return s


def build_agenda(prs, d):
    s = new_slide(prs)
    L = LAYOUTS["archetypes"]["agenda"]["slots"]
    slide_bg(s, "canvas")
    eb = _mono_kicker_color("blue" if not DARK else "blue-light")
    # Hyundai permanent header mark replaces the page kicker on agenda (Mode B).
    if _chevron_cfg() is not None and _chevron_cfg().get("header_mark", True):
        _chevron_header_mark(s)
    elif d.get("kicker") or d.get("eyebrow"):
        text(s, L["kicker"], _mono_kicker_prefix(d.get("kicker") or d.get("eyebrow", "AGENDA")),
             "eyebrow", color=eb, upper=True)
    text(s, L["title"], d.get("title", "목차"), "h1", color="ink")
    items = d.get("items", [])
    n = max(1, len(items))
    rows = _split_rows(L["items"], n, gap=8)
    BADGE = 56          # 번호 배지 폭
    TEXT_GAP = 24       # 배지와 텍스트 사이 간격
    COL_GAP = 40        # 타이틀/노트 열 사이 간격
    for i, it in enumerate(items):
        rx, ry, rw, rh = rows[i]
        tx = rx + BADGE + TEXT_GAP          # 텍스트 영역 좌측 시작
        avail = rw - (BADGE + TEXT_GAP)     # 배지 뒤 가용 폭
        square = GRAMMAR.get("badge") == "square"
        if _mono_cfg() is not None:
            # 모노메타: 채움 배지 대신 터미널 코멘트 번호(`// 01`). 채움 칩이 없어
            # 가용폭이 늘지만, 기존 tx/avail(배지 폭만큼 들여쓴)을 유지해 정렬 일관.
            text(s, [rx, ry, BADGE + TEXT_GAP, rh], _mono_index(i + 1), "label",
                 color=_mono_kicker_color("blue"), align="left", anchor="middle",
                 upper=True)
        else:
            numbox = rect(s, [rx, ry + (rh - BADGE) / 2, BADGE, BADGE],
                          fill="navy" if square else "blue",
                          radius=0 if square else RAD["md"])
            shape_text(numbox, f"{i+1:02d}", "h4",
                       color="white" if square else "on-accent", wrap=False)
        title = it if isinstance(it, str) else it.get("title", "")
        note = it.get("note") if isinstance(it, dict) else None
        if note:
            # 가용폭을 타이틀/노트로 균등 분할 + 갭. 노트가 ~48%(~800px) 확보 → 1줄.
            col_w = (avail - COL_GAP) / 2
            text(s, [tx, ry, col_w, rh], title, "h3", color="ink",
                 anchor="middle")
            text(s, [tx + col_w + COL_GAP, ry, col_w, rh], note, "body",
                 color="muted", anchor="middle", align="right")
        else:
            # 노트 없으면 타이틀이 가용폭 전부 사용
            text(s, [tx, ry, avail, rh], title, "h3", color="ink",
                 anchor="middle")
        if i < n - 1:
            line(s, rx, ry + rh + 4, rx + rw, ry + rh + 4, color="divider", w=1)
    return s


def build_closing(prs, d):
    s = new_slide(prs)
    L = LAYOUTS["archetypes"]["closing"]["slots"]
    # Closing mirrors the cover (bookend): same bg precedence so a light look
    # ends as light as it opened, and a dark-cover look stays dark on both ends.
    bg = d.get("bg") or _LOOK_COVER_BG or "canvas"
    slide_bg(s, bg)
    _orb_decor(s)   # 소프트 3D 구 데코(grammar.orb_decor); 무선언=무동작.
    dark = is_dark(_resolve(bg))
    ink = "white" if dark else "navy"
    # Hyundai closing = cover mirror: LEFT 5-stack chevron + eye, RIGHT-aligned WHITE
    # message (bookend symmetry). No page number / header mark (Mode A). No-op else.
    chev = _chevron_cfg() is not None
    if chev:
        _chevron_cover(s, "left")
        W = SLIDE["width_px"]
        rx = W * 0.50
        rw = W * 0.93 - rx
        if d.get("kicker") or d.get("eyebrow"):
            text(s, [rx, L["kicker"][1], rw, L["kicker"][3]],
                 d.get("kicker") or d["eyebrow"], "eyebrow",
                 color="blue-light", align="right", upper=True)
        text(s, [rx, L["title"][1], rw, L["title"][3]], d.get("title", ""),
             "h1", color="white", align="right")
        if d.get("subtitle"):
            text(s, [rx, L["subtitle"][1], rw, L["subtitle"][3]], d["subtitle"],
                 "body-lg", color="slate-300", align="right")
        return s
    if d.get("kicker") or d.get("eyebrow"):
        text(s, L["kicker"], _mono_kicker_prefix(d.get("kicker") or d["eyebrow"]),
             "eyebrow", color=_mono_kicker_color("blue-light" if dark else "blue"),
             upper=True)
    text(s, L["title"], d.get("title", ""), "h1", color=ink)
    if d.get("subtitle"):
        text(s, L["subtitle"], d["subtitle"], "body-lg",
             color="slate-300" if dark else "slate-500")
    if d.get("cta"):
        btn = card(s, L["cta"], "card", fill="blue")
        shape_text(btn, d["cta"], "body-md", color="on-accent")
    return s


# ===========================================================================
# Family 2 — Focus
# ===========================================================================
def build_statement(prs, d):
    s = new_slide(prs)
    if GRAMMAR.get("cover") == "color_block" and not d.get("image"):
        return _statement_color_block(s, d)
    header(s, d)
    L = LAYOUTS["archetypes"]["statement"]["slots"]
    if d.get("lead"):
        # 리드는 헤더 제목(h2/48px)과 같은 크기로 찍으면 '헤더가 둘'처럼 보인다.
        # 풀쿼트로 처리: 좌측 액센트 바 + 더 큰 h1(64px) + 세로 중앙 정렬 →
        # 슬라이드의 주인공(핵심 문장/인용)으로 명확히 위계가 잡힌다.
        lx, ly, lw, lh = L["lead"]
        bar_w = 6
        rect(s, [lx, ly + 6, bar_w, lh - 12], fill="blue", radius=RAD["sm"])
        text(s, [lx + bar_w + 24, ly, lw - bar_w - 24, lh], d["lead"], "h1",
             color="ink", anchor="middle")
    if d.get("image"):
        media(s, L["image"], d["image"], fit=d["image"].get("fit", "cover"))
        # 이미지 있으면 근거는 좌측만(미디어가 우측 차지)
        sx, sy, sw, sh = L["support"]
        sup_box = [sx, sy, L["image"][0] - sx - GAP, sh]
    else:
        sup_box = L["support"]
    if d.get("support"):
        if isinstance(d["support"], list):
            bullets(s, sup_box, d["support"], "body-lg")
        else:
            text(s, sup_box, d["support"], "body-lg", color="ink-2")
    return s


def build_feature(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["feature"]["slots"]
    if d.get("bullets"):
        bullets(s, L["body"], d["bullets"], "body-lg")
    elif d.get("body"):
        text(s, L["body"], d["body"], "body-lg", color="ink-2")
    media(s, L["media"], d.get("media", {}), fit=(d.get("media") or {}).get("fit", "cover"))
    return s


def _showcase_img(slide, x, y, w, h, src):
    """좌측·상단 정렬 이미지 — 비율 보존 rect(무크롭·무왜곡) + **활성 룩의 카드 스킨 상속**.
    카드에 border/shadow가 정의된 룩(brutalism·dark-luxury·glass…)이면 이미지에도 같은
    테두리·그림자를 입혀 덱과 결을 맞추고, 플랫 룩(swiss 등 border=None)이면 맨몸으로 둔다."""
    pic = slide.shapes.add_picture(src, IN(x), IN(y), IN(w), IN(h))
    pic.shadow.inherit = False
    c = COMPONENTS.get("card", {})
    b = c.get("border") or {}
    if b.get("color"):
        pic.line.color.rgb = C(b["color"])
        pic.line.width = Pt(b.get("width_pt", 1))
    else:
        pic.line.fill.background()
    sh = c.get("shadow")
    if sh:
        try:
            shadow(pic, sh.get("blur_pt", 12), sh.get("dist_pt", 4),
                   sh.get("dir_deg", 90), sh.get("color", "#000000"), sh.get("alpha", 78))
        except Exception:
            pass
    return pic


def build_showcase(prs, d):
    """좌측정렬 미디어 전시 + 불릿. 이미지는 헤더와 같은 좌측 마진에 정렬(무크롭),
    불릿은 **이미지 종횡비로 자동 배치** — 영역(넓고 낮음)보다 세로로 길면(대부분
    스크린샷) 높이에 꽉 차고 우측 여백에 불릿, 영역보다 가로로 길면(울트라와이드)
    폭에 꽉 차고 하단에 불릿. 분기 신호 = 이미지 비율(사용자 결정 불필요)."""
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["showcase"]["slots"]
    x0, y0, W, H = L["media"]
    m = d.get("media") or d.get("image") or {}
    bul = d.get("bullets") or []
    cap = m.get("caption")
    cap_h = 30 if cap else 0
    GAP = 40
    src = _resolve_src(m.get("src"))

    ar = None
    if src:
        try:
            from PIL import Image
            with Image.open(src) as im:
                ar = im.size[0] / im.size[1] if im.size[1] else None
        except Exception:
            ar = None

    def caption_under(bx, by, bw):
        if cap:
            text(s, [bx, by + 4, bw, cap_h], cap, "caption", color="muted")

    if not src or not ar:                                  # 플레이스홀더
        if bul:
            iw = W * 0.56
            _media_placeholder(s, [x0, y0, iw, H], m.get("label", "이미지 자리"))
            bullets(s, [x0 + iw + GAP, y0, W - iw - GAP, H], bul, "body-lg")
        else:
            _media_placeholder(s, [x0, y0, W, H], m.get("label", "이미지 자리"))
        return s

    region_ar = W / H
    avail_h = H - cap_h
    if ar < region_ar:                                     # 높이 바운드(대부분 스크린샷)
        iw, ih = avail_h * ar, avail_h
        leftover = W - iw - GAP
        if bul and leftover >= 480:                        # 좌 이미지 + 우 불릿
            _showcase_img(s, x0, y0, iw, ih, src)
            caption_under(x0, y0 + ih, iw)
            bullets(s, [x0 + iw + GAP, y0, leftover, H], bul, "body-lg")
        elif bul:                                          # 우측 폭 부족 → 이미지 축소 + 하단 불릿
            ih, iw = avail_h * 0.58, avail_h * 0.58 * ar
            _showcase_img(s, x0, y0, iw, ih, src)
            caption_under(x0, y0 + ih, iw)
            bullets(s, [x0, y0 + ih + cap_h + GAP, W, H - ih - cap_h - GAP], bul, "body")
        else:                                              # 불릿 없음 → 좌측정렬 큰 이미지
            _showcase_img(s, x0, y0, iw, ih, src)
            caption_under(x0, y0 + ih, iw)
    else:                                                  # 폭 바운드(울트라와이드) → 상단 이미지 + 하단 불릿
        if bul:                                            # 하단 불릿 띠 확보 위해 이미지 높이 캡(좌측정렬 유지)
            ih = min(W / ar, avail_h * 0.62)
            iw = ih * ar
            _showcase_img(s, x0, y0, iw, ih, src)
            caption_under(x0, y0 + ih, iw)
            bullets(s, [x0, y0 + ih + cap_h + GAP, W, H - ih - cap_h - GAP], bul, "body")
        else:                                              # 불릿 없으면 전폭
            iw, ih = W, W / ar
            _showcase_img(s, x0, y0, iw, ih, src)
            caption_under(x0, y0 + ih, iw)
    return s


# ===========================================================================
# Family 3 — Set
# ===========================================================================
def build_duo(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["duo"]["slots"]
    items = d.get("items") or [d.get("left", {}), d.get("right", {})]
    accents = accent_cycle(2)
    _content_card(s, list(L["item_l"]), items[0] if items else {}, accents[0])
    if len(items) > 1:
        _content_card(s, list(L["item_r"]), items[1], accents[1])
    return s


def build_trio(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["trio"]["slots"]
    items = d.get("items") or d.get("columns", [])
    accents = accent_cycle(3)
    for i, key in enumerate(["col1", "col2", "col3"]):
        if i < len(items):
            _content_card(s, list(L[key]), items[i], accents[i % 3])
    return s


def build_grid(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["grid"]["slots"]
    items = d.get("items") or d.get("cells", [])
    n = len(items)
    ncol = 3 if n > 4 else 2
    nrow = 2 if n > 2 else 1
    boxes = _split_grid(L["cells"], ncol, nrow)
    accents = accent_cycle(max(1, n))
    for i, it in enumerate(items):
        if i < len(boxes):
            _content_card(s, boxes[i], it, accents[i % len(accents)])
    return s


def build_list(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["list"]["slots"]
    items = d.get("items", [])
    n = max(1, len(items))
    rows = _split_rows(L["items"], n, gap=10)
    # mono-meta list indices read as `// 01` (bare terminal numbers), which need a
    # wider gutter than the 44px accent chip. The title column shifts right to match
    # so the index + title never collide. No grammar → legacy 44px chip + 64 gutter.
    bw, gut = (84, 104) if _mono_cfg() is not None else (44, 64)
    for i, it in enumerate(items):
        rx, ry, rw, rh = rows[i]
        _num_badge(s, [rx, ry + (rh - 44) / 2, bw, 44], str(i + 1))
        title = it if isinstance(it, str) else it.get("title", "")
        text(s, [rx + gut, ry, rw - gut, rh], title, "h4", color="ink",
             anchor="middle")
        sub = it.get("body") if isinstance(it, dict) else None
        if sub:
            text(s, [rx + gut, ry + rh * 0.55, rw - gut, rh * 0.4], sub, "body",
                 color="muted")
    return s


# ===========================================================================
# Family 4 — Contrast
# ===========================================================================
def build_versus(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["versus"]["slots"]
    left = d.get("left", {})
    right = d.get("right", {})
    # 좌(어둡게) vs 우(밝게)
    lx, ly, lw, lh = L["left"]
    # 프리즘/글래스: 두 패널 모두 글래스("2~3 글래스 패널, 추천 패널만 보더 글로우") —
    # 프리즘=투명+그라디언트 헤어라인 보더, 비비드=frosted 흰빛 글래스. 좌 패널은 비추천
    # → 무강조 글래스, 텍스트 ink. 그 외 룩은 기존 slab 강조 타일(잉크면+반전 텍스트) 유지
    # (무회귀). vivid 의 versus 는 검정 slab 대신 글래스 패널이라야 채도 메시 언어와 일관.
    _glass_v = (_prism_cfg() is not None) or bool(GRAMMAR.get("glass"))
    # quiet_compare (engineered-dark): hairline panels, NO fill, single-accent markers —
    # the pack forbids 2+ accents (no pro/con red/green) and a filled winner tile. The
    # loser panel is a bare hairline card with muted markers; the winner is a hairline
    # card with an accent border (+ glow seam later). No grammar → unchanged (other looks
    # keep their slab/glass versus). DRY: one flag flips the whole archetype's palette.
    _quiet_v = bool(GRAMMAR.get("quiet_compare"))
    if _quiet_v:
        card(s, L["left"], "card", fill=None)
        ltxt, lmk = "ink", "muted"
    elif _glass_v:
        card(s, L["left"], "card", fill="surface-2")
        ltxt, lmk = "ink", "muted"
    else:
        rect(s, L["left"], fill="slab", radius=RAD["lg"])
        ltxt, lmk = "on-slab", "red"
    text(s, [lx + PAD, ly + 28, lw - 2 * PAD, 30],
         left.get("label", "BEFORE"), "label", color=ltxt, upper=True)
    if left.get("heading"):
        text(s, [lx + PAD, ly + 64, lw - 2 * PAD, 56], left["heading"], "h3",
             color=ltxt)
    if left.get("points"):
        bullets(s, [lx + PAD, ly + 140, lw - 2 * PAD, lh - 170], left["points"],
                "body", color=ltxt, marker_color=lmk)
    rx, ry, rw, rh = L["right"]
    _rfill = None if _quiet_v else "surface-3"
    rcard = card(s, L["right"], "card", fill=_rfill, line_color="blue", line_w=1.5)
    if _quiet_v and DARK:
        # 추천 패널만 보더 글로우(스펙: "추천 열 1개만 보더 #8B7BF0 1.5px + 코너 라디얼 글로우").
        # cards 글로우 패밀리는 정밀 톤상 꺼져 있으므로(carte blanche halo 방지) 여기서만
        # 절제된 직접 halo 를 건다 — 이 한 요소가 '추천' 시그널의 핵심이라 의도적 예외.
        glow(rcard, rad_px=9, color=_resolve("blue"), alpha=30)
    text(s, [rx + PAD, ry + 28, rw - 2 * PAD, 30],
         right.get("label", "AFTER"), "label", color="blue", upper=True)
    if right.get("heading"):
        text(s, [rx + PAD, ry + 64, rw - 2 * PAD, 56], right["heading"], "h3",
             color="ink")
    if right.get("points"):
        bullets(s, [rx + PAD, ry + 140, rw - 2 * PAD, rh - 170], right["points"],
                "body", color="ink-2", marker_color=("blue" if _quiet_v else "green"))
    # 중앙 VS 마크
    vs = oval(s, [(lx + lw + rx) / 2 - 36, ly + lh / 2 - 36, 72, 72], fill="surface",
              line_color="hairline", line_w=1.5)
    shape_text(vs, "VS", "h4", color="ink")
    return s


def build_matrix(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["matrix"]["slots"]
    x, y, w, h = L["table"]
    cols = d.get("headers", [])
    rows = d.get("rows", [])
    ncol = max(1, len(cols))
    cw = w / ncol
    head_h = 60
    # Row height adapts to the active body size: small text keeps compact 72px
    # rows; large text (e.g. lecture looks) gets a row tall enough to hold a
    # WRAPPED 2-line cell cleanly. Capped by the slot so the table never spills.
    _bpx = _size_px("body")
    _blh = (_type_override().get("line") or {}).get("body", TYPO["body"].get("line", 1.2))
    two_line = _bpx * _blh * 2 + 24
    rh = min((h - head_h) / max(1, len(rows)), max(72, two_line))
    rect(s, [x, y, w, head_h], fill="slab")
    for i, c in enumerate(cols):
        text(s, [x + i * cw + 24, y + 16, cw - 32, 32], str(c), "body-md",
             color="on-slab", align="left" if i == 0 else "center")
    for r, row in enumerate(rows):
        ry = y + head_h + r * rh
        if r % 2 == 1:
            rect(s, [x, ry, w, rh], fill="surface-2")
        cells = row if isinstance(row, list) else row.get("cells", [])
        for i, cell in enumerate(cells):
            # Full-height box + middle anchor so long cells WRAP to 2 lines and
            # stay vertically centered, instead of clipping at a fixed 28px line
            # (matters at larger type scales, e.g. lecture looks). Short cells are
            # unchanged — one centered line.
            text(s, [x + i * cw + 24, ry + 6, cw - 32, rh - 12], str(cell),
                 "body", color="ink" if i == 0 else "body",
                 align="left" if i == 0 else "center", anchor="middle")
    return s


def build_rank(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["rank"]["slots"]
    x, y, w, h = L["rows"]
    rows = d.get("rows", [])
    n = max(1, len(rows))
    rh = min(86, h / n)
    for r, row in enumerate(rows):
        ry = y + r * rh
        hl = row.get("highlight")
        if hl:
            card(s, [x, ry, w, rh - 8], "card", fill="surface-3")
        _num_badge(s, [x + 12, ry + (rh - 8) / 2 - 24, 48, 48],
                   str(row.get("rank", r + 1)),
                   oval_fill="blue" if hl else "divider",
                   oval_txt="on-accent" if hl else "muted")
        text(s, [x + 80, ry, w - 300, rh - 8], row.get("label", ""), "body-md",
             color="ink", anchor="middle")
        delta = str(row.get("delta", ""))
        dc = "green" if delta.startswith("▲") else (
            "red" if delta.startswith("▼") else "subtle")
        text(s, [x + w - 200, ry, 180, rh - 8], delta, "body-md", color=dc,
             align="right", anchor="middle")
    return s


# ===========================================================================
# Family 5 — Field
# ===========================================================================
def build_quadrant(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["quadrant"]["slots"]
    cells = d.get("cells", [])
    boxes = _split_grid(L["cells"], 2, 2)
    # 4셀 채움은 균일한 라이트 서피스(slate-50)로 통일 — 셀별 강조는
    # kicker/heading 색(accents)에서만 표현한다(랜덤하게 한 셀만 짙어지지 않게).
    # 단, block_fill 문법(브루탈리즘 "사분면마다 다른 평면 원색")이면 각 셀을 팔레트
    # 원색으로 채우고 텍스트 색을 셀 휘도 기준으로 플립(옐로 셀=어두운 글자).
    # accent_cycle 자체가 팔레트 부족 시 톤 변주로 채우므로(글로벌), 4셀이 팔레트 색 수를
    # 넘어도 같은 솔리드 색이 반복되지 않는다(인-팔레트 톤). 셀별 채움/라벨에 그대로 사용.
    accents = accent_cycle(4)
    filled = bool(GRAMMAR.get("block_fill"))
    # rule_field 문법(스위스): 셀 박스 대신 중앙 십자 룰만 — "박스 채움 없음, 십자 룰만".
    rule = bool(GRAMMAR.get("rule_field"))
    spot = ACCENTS[0] if ACCENTS else "blue"
    if rule:
        gx, gy, gw, gh = L["cells"]
        line(s, gx + gw / 2, gy, gx + gw / 2, gy + gh, color="navy", w=2)
        line(s, gx, gy + gh / 2, gx + gw, gy + gh / 2, color="navy", w=2)
    for i in range(min(4, len(cells))):
        bx, by, bw, bh = boxes[i]
        cell = cells[i]
        on = _ink_on(accents[i]) if filled else None
        if not rule:
            card(s, boxes[i], "card", fill=accents[i] if filled else "surface-2")
        lbl_col = on if filled else (spot if rule else accents[i])
        mk_col = on if filled else (spot if rule else accents[i])
        text(s, [bx + PAD, by + 24, bw - 2 * PAD, 30], cell.get("label", ""),
             "label", color=lbl_col, upper=True)
        if cell.get("heading"):
            text(s, [bx + PAD, by + 62, bw - 2 * PAD, 44], cell["heading"], "h3",
                 color=on if filled else "ink")
        if cell.get("bullets"):
            bullets(s, [bx + PAD, by + 124, bw - 2 * PAD, bh - 150],
                    cell["bullets"], "body", color=on if filled else "ink-2",
                    marker_color=mk_col)
        elif cell.get("body"):
            text(s, [bx + PAD, by + 124, bw - 2 * PAD, bh - 150], cell["body"],
                 "body", color=on if filled else "body")
    # 축 라벨(가장자리)
    cx, cy, cw, ch = L["cells"]
    if d.get("x_axis"):
        text(s, [cx, cy + ch + 6, cw, 24], d["x_axis"], "caption",
             color="muted", align="center")
    if d.get("y_axis"):
        text(s, [cx - 60, cy, 60, ch], d["y_axis"], "caption", color="muted",
             align="center", anchor="middle")
    return s


def build_map(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["map"]["slots"]
    x, y, w, h = L["plot"]
    # 데이터 기하(잠금): 좌표 축은 안 기울인다. 문법별 시각언어:
    #  brutal = 굵은 검정 십자축 + 검정보더 사각 마커 / swiss(rule_field) = 박스 없음 +
    #  얇은 잉크 십자축 + 무채색 사각 마커(강조 1개만 스팟) / 기본 = surface 박스+얇은 회색축+원형점.
    brutal = GRAMMAR.get("chart_style") == "brutal"
    rule = bool(GRAMMAR.get("rule_field"))
    if not rule:
        card(s, L["plot"], "card", fill="surface-2")
    if brutal:
        axc, axw = "navy", 4
    elif rule:
        axc, axw = "navy", 2.5
    else:
        axc, axw = "hairline", 1.5
    line(s, x + w / 2, y + 16, x + w / 2, y + h - 16, color=axc, w=axw)
    line(s, x + 16, y + h / 2, x + w - 16, y + h / 2, color=axc, w=axw)
    if d.get("x_axis"):
        text(s, [x, y + h + 6, w, 24], d["x_axis"], "caption", color="muted",
             align="center")
    if d.get("y_axis"):
        text(s, [x - 60, y, 60, h], d["y_axis"], "caption", color="muted",
             align="center", anchor="middle")
    cc = COMPONENTS.get("card", {})
    bcol = (cc.get("border") or {}).get("color", "navy")
    bwid = (cc.get("border") or {}).get("width_pt", 4)
    for p in d.get("points", []):
        px = x + 20 + p.get("x", 0.5) * (w - 40)
        py = y + h - 20 - p.get("y", 0.5) * (h - 40)
        hl = p.get("highlight")
        mk = (ACCENTS[0] if ACCENTS else "blue") if hl else (
            ACCENTS[1] if (brutal and len(ACCENTS) > 1) else "blue")
        if brutal:
            rect(s, [px - 16, py - 16, 32, 32], fill=mk, line_color=bcol,
                 line_w=bwid, radius=0)
        elif rule:
            # 무채색 사각 마커, 강조 1개만 스팟(스위스)
            smk = (ACCENTS[0] if ACCENTS else "blue") if hl else "navy"
            rect(s, [px - 11, py - 11, 22, 22], fill=smk, radius=0)
        else:
            oval(s, [px - 14, py - 14, 28, 28], fill="orange" if hl else "blue")
        text(s, [px + 22, py - 14, 280, 28], p.get("label", ""), "body-md",
             color="ink" if hl else "body")
    return s


# ===========================================================================
# Family 6 — Structure
# ===========================================================================
def build_flow(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["flow"]["slots"]
    steps = d.get("steps", [])
    n = max(1, len(steps))
    boxes = _split_cols(L["steps"], n, gap=56)
    square = GRAMMAR.get("badge") == "square"
    # block_fill 문법: 프로세스 플로우는 "전부 한 색, 핵심만 다른 원색"(quadrant·KPI의
    # 색-셀별-다름과 달리 *시퀀스*라 균일 바탕+강조가 의미적으로 맞음). 바탕 = 가장 밝은
    # (차분한) 액센트, 강조 = 브랜드 1차 액센트(ACCENTS[0]). highlight 미지정 시 마지막 단계.
    filled = bool(GRAMMAR.get("block_fill"))
    base = hot = neutrals = None
    hl = d.get("highlight")
    if filled and ACCENTS:
        # base = the calm field color, hot = the highlighted step. Default: brightest
        # accent is the calm field + the primary pops. A look can flip this via
        # grammar.flow_field_primary so the PRIMARY accent is the neutral field and the
        # LOUDEST (brightest) accent is the highlight — so the field stays charcoal/teal
        # and only the active step flares (bold-block-infographic 앰버 강조). Opt-in only
        # → every other block_fill look (brutalism/memphis/monochrome) is byte-identical.
        if GRAMMAR.get("flow_field_primary"):
            base = ACCENTS[0]
            hot = max(ACCENTS, key=lambda a: _rel_lum(_resolve(a)))
            # Non-highlight steps ALTERNATE the remaining (neutral) accents so a 3-step
            # flow reads charcoal→teal→amber, not charcoal→charcoal→amber (the pack:
            # 차콜→틸 교대 + 활성만 앰버). One neutral left → uniform field as before.
            neutrals = [a for a in ACCENTS if _resolve(a) != _resolve(hot)] or [ACCENTS[0]]
        else:
            base = max(ACCENTS, key=lambda a: _rel_lum(_resolve(a)))
            hot = ACCENTS[0]
        if hl is None:
            hl = n - 1
    for i, st in enumerate(steps):
        bx, by, bw, bh = boxes[i]
        # 노드(자유 블록)는 틸트 대상 — 카드+배지+텍스트를 한 묶음으로 회전.
        # 커넥터는 묶음 밖(틸트 후)에 그려 직선 유지(노드 사이를 잇는 기준선).
        n0 = len(s.shapes)
        if not filled:
            nodefill = "surface-2"
        elif i == hl:
            nodefill = hot
        elif neutrals is not None:
            nodefill = neutrals[i % len(neutrals)]
        else:
            nodefill = base
        on = _ink_on(nodefill) if filled else None
        card(s, [bx, by + 40, bw, bh - 40], "card", fill=nodefill)
        if square:
            bd = rect(s, [bx + bw / 2 - 32, by, 64, 64], fill="navy", radius=0)
            shape_text(bd, str(i + 1), "h3", color="white", wrap=False)
        else:
            circ = oval(s, [bx + bw / 2 - 32, by + 8, 64, 64], fill="blue")
            shape_text(circ, str(i + 1), "h3", color="on-accent")
            _glow_shape(circ, "badges", color=_resolve("blue"))
        text(s, [bx + 20, by + 96, bw - 40, 56], st.get("heading", ""), "h4",
             color=on if filled else "ink", align="center")
        if st.get("body"):
            text(s, [bx + 20, by + 160, bw - 40, bh - 200], st["body"], "body",
                 color=on if filled else "body", align="center")
        _tilt_group(s, n0, bx + bw / 2.0, by + bh / 2.0)
        if i < n - 1:
            ny = by + 40 + (bh - 40) / 2
            if square:
                # 굵은 검정 직선 + 각진 화살촉
                line(s, bx + bw + 8, ny, bx + bw + 44, ny, color="navy", w=4)
                _arrowhead(s, bx + bw + 50, ny, 18, "navy")
            else:
                # 발광 룩이면 커넥터를 액센트 네온 라인으로(무선언=hairline 그대로).
                cc = "blue" if _glow_cfg() else "hairline"
                line(s, bx + bw + 10, ny, bx + bw + 46, ny, color=cc, w=2,
                     glow_line=True)
    return s


def build_system(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["system"]["slots"]
    nodes = d.get("nodes", [])
    n = len(nodes)
    # 코어는 nodes 영역(그리드 파생) 중앙에 두고, 노드는 그 둘레에 배치해
    # 코어가 어떤 노드도 가리지 않게 한다(오빗수학 폐기·그리드 셀 파생).
    nx, ny, nw, nh = L["nodes"]
    # 코어 박스: nodes 영역의 가운데 컬럼 블록(전체 폭의 중앙 ~30%).
    core_w = nw * 0.30
    core_h = nh * 0.40
    core_box = [nx + (nw - core_w) / 2, ny + (nh - core_h) / 2, core_w, core_h]
    ccx = core_box[0] + core_box[2] / 2
    ccy = core_box[1] + core_box[3] / 2
    # 노드 둘레 배치: 좌/우 측면 컬럼 블록에 노드를 위→아래로 분배.
    # 좌측에 ceil(n/2), 우측에 나머지. 코어가 차지하는 중앙 컬럼은 비운다.
    side_w = (nw - core_w) / 2 - GAP
    left_n = (n + 1) // 2
    right_n = n - left_n
    boxes = []
    if left_n:
        boxes += _split_rows([nx, ny, side_w, nh], left_n, gap=GAP)
    if right_n:
        rx = nx + nw - side_w
        boxes += _split_rows([rx, ny, side_w, nh], right_n, gap=GAP)
    # 먼저 커넥터(코어→노드). 발광 룩이면 액센트 네온 라인(무선언=hairline 그대로).
    _ccol = "blue" if _glow_cfg() else "hairline"
    _cw = 2 if _glow_cfg() else 1.5
    for i in range(min(len(boxes), n)):
        bx, by, bw, bh = boxes[i]
        line(s, ccx, ccy, bx + bw / 2, by + bh / 2, color=_ccol, w=_cw,
             glow_line=True)
    # 노드 박스
    for i in range(min(len(boxes), n)):
        bx, by, bw, bh = boxes[i]
        nd = nodes[i]
        nb = card(s, [bx, by, bw, bh], "card", fill="surface-3",
                  line_color="blue", line_w=1.5)
        shape_text(nb, nd if isinstance(nd, str) else nd.get("label", ""),
                   "body-md", color="ink")
    # 코어(맨 위에 — 둘레가 비어 있으므로 어떤 노드도 가리지 않음). 강조 잉크 타일.
    cb = rect(s, core_box, fill="blue", radius=RAD["lg"])
    shape_text(cb, d.get("core", "Core"), "h3", color="on-accent")
    _glow_shape(cb, "badges", color=_resolve("blue"))   # 발광 룩이면 코어가 빛나는 허브
    return s


# ===========================================================================
# Family 7 — Data
# ===========================================================================
def build_share(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["share"]["slots"]
    spec = dict(d)
    spec.setdefault("type", "doughnut")
    spec["legend"] = False  # share draws its own color-chip legend below (bug B)
    chart(s, L["chart"], spec)
    # 범례/주석 영역: 항목별 색칩 + 라벨 + 값
    items = []
    cats = d.get("categories", [])
    vals = d.get("values") or (d.get("series", [{}])[0].get("values", []) if d.get("series") else [])
    for i, c in enumerate(cats):
        items.append((c, vals[i] if i < len(vals) else ""))
    lx, ly, lw, lh = L["legend"]
    rowh = min(48, lh / max(1, len(items)))
    # 차트 문법별 칩: 브루탈=원색+검정보더 사각(도넛 조각색 일치), 스위스=무채색 잉크+
    # 강조 1개만 스팟(프로포션 바 세그먼트색과 일치), 기본=accent 라운드 칩.
    _cs = GRAMMAR.get("chart_style")
    brutal = _cs == "brutal"
    swiss = _cs == "swiss"
    prism = _cs == "prism"
    # 프리즘 칩 = 프리즘 스톱색(세그먼트 링 색과 일치) — 무지개로 늘리지 않고 3스톱 순환.
    prism_stops = ["#" + s for s in _prism_stops()] if prism else None
    chip_cols = (_brutal_palette(len(items)) if brutal else
                 _swiss_share_colors(vals, d) if swiss else None)
    bcol, bwid, _sh = _brutal_skin() if brutal else (None, 0, None)
    for i, (c, v) in enumerate(items):
        ry = ly + i * rowh
        if brutal:
            rect(s, [lx, ry + 4, 28, 28], fill=chip_cols[i], line_color=bcol,
                 line_w=bwid, radius=0)
        elif swiss:
            rect(s, [lx, ry + 5, 24, 24], fill=chip_cols[i], radius=0)
        elif prism:
            ch = rect(s, [lx, ry + 6, 20, 20],
                      fill=prism_stops[i % len(prism_stops)], radius=RAD["sm"])
            _glow_shape(ch, "chart", color=prism_stops[i % len(prism_stops)])
        else:
            rect(s, [lx, ry + 6, 20, 20], fill=chart_color(i), radius=RAD["sm"])
        text(s, [lx + 40, ry, lw - 130, rowh], str(c), "body-md", color="ink",
             anchor="middle")
        if v != "":
            text(s, [lx + lw - 90, ry, 90, rowh], str(v), "body-md",
                 color="body", align="right", anchor="middle")
    return s


def build_bars(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "bar_clustered")
    chart(s, LAYOUTS["archetypes"]["bars"]["slots"]["chart"], spec)
    return s


def build_trend(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "line_markers" if len(d.get("series", [])) > 1 else "line")
    chart(s, LAYOUTS["archetypes"]["trend"]["slots"]["chart"], spec)
    return s


def build_spread(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "column_clustered")  # 히스토그램은 컬럼으로 표현
    spec.setdefault("legend", False)
    # 가독성 레이어: 히스토그램은 ctype=column 이지만 *조밀한 분포*라 busy 배경 위에서
    # 가독성이 약하다 → 막대(bars=ROBUST)와 달리 백드롭 대상(fragile)임을 명시 힌트.
    spec.setdefault("_fragile", True)
    chart(s, LAYOUTS["archetypes"]["spread"]["slots"]["chart"], spec)
    return s


def build_correlate(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "bubble" if any("size" in p for p in d.get("points", [])) else "xy_scatter")
    chart(s, LAYOUTS["archetypes"]["correlate"]["slots"]["chart"], spec)
    return s


def build_kpi(prs, d):
    s = new_slide(prs)
    header(s, d)
    # KPI 카드 배열은 items(권장) / kpis / stats 중 하나로 받는다(Set·List 계열과 키 일관).
    spec = {"type": "kpi_cards", "kpis": d.get("kpis") or d.get("stats") or d.get("items", [])}
    chart(s, LAYOUTS["archetypes"]["kpi"]["slots"]["cells"], spec)
    return s


# ---------------------------------------------------------------------------
# Routing table — 새 슬러그 → 빌드 함수.
# ---------------------------------------------------------------------------
ARCHETYPE_TO_BUILDER = {
    # Frame
    "cover": build_cover, "section": build_section, "agenda": build_agenda,
    "closing": build_closing,
    # Focus
    "statement": build_statement, "feature": build_feature, "showcase": build_showcase,
    # Set
    "duo": build_duo, "trio": build_trio, "grid": build_grid, "list": build_list,
    # Contrast
    "versus": build_versus, "matrix": build_matrix, "rank": build_rank,
    # Field
    "quadrant": build_quadrant, "map": build_map,
    # Structure
    "flow": build_flow, "system": build_system,
    # Data
    "share": build_share, "bars": build_bars, "trend": build_trend,
    "spread": build_spread, "correlate": build_correlate, "kpi": build_kpi,
}


# ---------------------------------------------------------------------------
# Demo spec (--spec 없을 때) — 모든 아키타입 1장씩, 중립 플레이스홀더 콘텐츠.
# ---------------------------------------------------------------------------
def _demo_spec():
    return {
        "title": "ppt-lab Rebuild — 전 아키타입 데모",
        "slides": [
            {"archetype": "cover", "data": {"kicker": "REBUILD",
                "title": "클린룸 PPT 엔진\n그리드 SSOT 데모",
                "subtitle": "grid.json → layouts.json → 23 아키타입",
                "meta": "2026 · ppt-grid-deck"}},
            {"archetype": "agenda", "data": {"kicker": "AGENDA", "title": "목차",
                "items": ["프레임", "포커스", "세트", "대비", "필드", "구조", "데이터"]}},
            {"archetype": "section", "data": {"index": "01", "title": "데이터 패밀리",
                "caption": "Zelazny 5비교유형 + KPI"}},
            {"archetype": "statement", "data": {"kicker": "주장", "title": "한 슬라이드, 한 메시지",
                "lead": "핵심 주장을 먼저 단언한다", "support": ["근거 1", "근거 2", "근거 3"]}},
            {"archetype": "feature", "data": {"kicker": "FEATURE", "title": "히어로 + 인사이트",
                "bullets": ["좌측 본문", "우측 미디어"], "media": {"caption": "제품 스크린샷"}}},
            {"archetype": "duo", "data": {"kicker": "비교 아님", "title": "두 동등 항목",
                "items": [{"label": "A", "heading": "항목 A", "bullets": ["요점 1", "요점 2"]},
                          {"label": "B", "heading": "항목 B", "image": {}, "body": "이미지 카드 변종"}]}},
            {"archetype": "trio", "data": {"kicker": "세 가지", "title": "세 동등 항목",
                "items": [{"label": "S", "heading": "속도", "body": "빠름"},
                          {"label": "C", "heading": "일관성", "body": "한 톤"},
                          {"label": "Q", "heading": "품질", "body": "정제"}]}},
            {"archetype": "grid", "data": {"kicker": "카드", "title": "4업 그리드",
                "items": [{"heading": f"카드 {i+1}", "body": "설명"} for i in range(4)]}},
            {"archetype": "list", "data": {"kicker": "리스트", "title": "순차 항목",
                "items": [{"title": f"단계 {i+1}", "body": "설명"} for i in range(4)]}},
            {"archetype": "versus", "data": {"kicker": "대비", "title": "Before vs After",
                "left": {"label": "BEFORE", "heading": "수작업", "points": ["느림", "불일치"]},
                "right": {"label": "AFTER", "heading": "자동화", "points": ["빠름", "일관"]}}},
            {"archetype": "matrix", "data": {"kicker": "표", "title": "기준 × 옵션",
                "headers": ["기준", "옵션 A", "옵션 B"],
                "rows": [["속도", "높음", "중간"], ["비용", "낮음", "높음"]]}},
            {"archetype": "rank", "data": {"kicker": "서열", "title": "순위 변동",
                "rows": [{"rank": 1, "label": "1위", "delta": "▲2", "highlight": True},
                         {"rank": 2, "label": "2위", "delta": "▼1"},
                         {"rank": 3, "label": "3위", "delta": "—"}]}},
            {"archetype": "quadrant", "data": {"kicker": "2×2", "title": "사분면",
                "x_axis": "x축", "y_axis": "y축",
                "cells": [{"label": L, "heading": L, "bullets": ["내용"]}
                          for L in ["강점", "약점", "기회", "위협"]]}},
            {"archetype": "map", "data": {"kicker": "포지셔닝", "title": "연속 2축 맵",
                "x_axis": "가격", "y_axis": "품질",
                "points": [{"x": 0.7, "y": 0.8, "label": "우리", "highlight": True},
                           {"x": 0.3, "y": 0.4, "label": "경쟁사"}]}},
            {"archetype": "flow", "data": {"kicker": "과정", "title": "3단계 흐름",
                "steps": [{"heading": f"단계 {i+1}", "body": "설명"} for i in range(3)]}},
            {"archetype": "system", "data": {"kicker": "구조", "title": "아키텍처",
                "core": "코어", "nodes": ["노드 A", "노드 B", "노드 C", "노드 D"]}},
            {"archetype": "share", "data": {"kicker": "구성", "title": "점유율",
                "categories": ["A", "B", "C", "D"], "values": [40, 30, 20, 10]}},
            {"archetype": "bars", "data": {"kicker": "항목", "title": "항목 비교",
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "series": [{"name": "건수", "values": [120, 240, 410, 690]}]}},
            {"archetype": "trend", "data": {"kicker": "시계열", "title": "추이",
                "categories": ["1월", "2월", "3월", "4월", "5월"],
                "series": [{"name": "Organic", "values": [10, 18, 27, 35, 52]},
                           {"name": "Referral", "values": [5, 9, 14, 20, 28]}]}},
            {"archetype": "spread", "data": {"kicker": "빈도", "title": "분포",
                "categories": ["0-10", "10-20", "20-30", "30-40", "40+"],
                "values": [3, 12, 25, 18, 7]}},
            {"archetype": "correlate", "data": {"kicker": "상관", "title": "산점도",
                "points": [{"x": i, "y": i * 1.4 + (i % 3)} for i in range(8)]}},
            {"archetype": "kpi", "data": {"kicker": "지표", "title": "핵심 지표",
                "kpis": [{"label": "제작시간", "value": "-78%", "note": "건당"},
                         {"label": "MAU", "value": "12K", "note": "+24%"},
                         {"label": "재사용", "value": "23", "note": "아키타입"},
                         {"label": "NPS", "value": "62", "note": "베타"}]}},
            {"archetype": "closing", "data": {"kicker": "NEXT", "bg": "navy",
                "title": "지금, 한 번의 요청으로", "subtitle": "자료 첨부 → 빌드",
                "cta": "데모 빌드 →"}},
        ],
    }


DEMO_SPEC = _demo_spec()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build(spec, out_path, palette_override=None, style_override=None, look_override=None,
          size_override=None, chart_backdrop_opacity=None):
    warnings = []
    look_id = look_override or spec.get("look")
    lk = apply_look(look_id) if look_id else None      # 룩: fonts+components+palette 통째
    if look_id and lk is None:
        warnings.append(f"unknown look '{look_id}' -> ignored (house/base fallback)")
    # 명시 --style 이 있거나 룩이 없을 때만 스타일 적용(룩의 컴포넌트 보존)
    style_id = style_override or spec.get("style")
    if style_id or lk is None:
        st = apply_style(style_id)
        # 명시 style_id 가 주어졌는데 매칭 실패 → house 로 폴백됨을 경고
        if style_id and st is not None and STYLES.get(str(style_id)) is None:
            valid_styles = ", ".join(k for k in STYLES if k != "_meta")
            warnings.append(
                f"style '{style_id}' not found (valid: {valid_styles}) -> using house")
    # 명시 --palette 가 있거나 룩이 없을 때만 팔레트 적용(룩의 색 보존)
    pal = palette_override or spec.get("palette")
    if pal is not None or lk is None:
        applied = apply_palette(pal)
        # 명시 pal 이 주어졌는데 매칭 실패 → base 로 폴백됨을 경고
        if pal is not None and applied is None:
            valid_pals = ", ".join(str(p["id"]) for p in PALETTES)
            warnings.append(
                f"palette {pal} not found (valid: {valid_pals}) -> using base")
    else:
        applied = {"name": f"look:{look_id} ({lk['palette'].get('blue','')})"}
    # Size axis — applied AFTER look/style settle GRAMMAR (style override may reset
    # it). '소'/None leaves the look's natural type scale untouched.
    size_tier = size_override or spec.get("size")
    if size_tier:
        canon = apply_size(size_tier)
        if canon is None:
            warnings.append(
                f"size '{size_tier}' not recognized (use 소/중/대) -> look default")
    # HITL 노브 주입(차트 백드롭 불투명도 = 일관성↔가독성 균형). CLI > spec > 룩 토큰 > 엔진 바닥선.
    # apply_look/style 가 GRAMMAR 를 settle한 *뒤*에 주입해야 clear 에 안 지워진다.
    cbo = chart_backdrop_opacity if chart_backdrop_opacity is not None \
        else spec.get("chart_backdrop_opacity")
    if cbo is not None:
        GRAMMAR["chart_backdrop_opacity"] = cbo
    # Deck-level dark control: any deck (even on a light look) can declare its own
    # `canvas` and the role palette auto-derives to match that background. This is
    # the generalization — "the deck picks a palette that fits its background".
    # Re-settle DARK/roles after style+palette overrides (ACCENTS may have changed).
    # NOTE: this re-settle historically re-derives ALL dark roles from the canvas,
    # clobbering a look's hand-pinned surfaces — harmless for opaque-card looks but
    # FATAL for glass/prism looks, whose card fill is translucent so the exact surface
    # hue is the visible tint: prism = near-transparent OUTLINE nodes (surface = faint
    # glass tint), vivid-gradient-future = frosted WHITE-tinted glass panels (surfaces
    # pinned light so a 16% fill reads as bright frosted glass over the mesh, not a dark
    # slab). So when the active look declares prism OR glass, preserve its explicit roles
    # (scoped to those grammars to keep every other look — incl. dark-tech — byte-
    # identical). A deck-level canvas override still re-derives everything to match.
    if spec.get("canvas"):
        COLORS["canvas"] = spec["canvas"]
        resolve_dark(force=spec.get("dark"))
    elif GRAMMAR.get("prism") or GRAMMAR.get("glass"):
        resolve_dark(_LOOK_EXPLICIT_ROLES, force=spec.get("dark"))
    else:
        resolve_dark(force=spec.get("dark"))
    prs = Presentation()
    prs.slide_width = Inches(SLIDE["width_in"])
    prs.slide_height = Inches(SLIDE["height_in"])
    for i, sl in enumerate(spec.get("slides", [])):
        # 슬롯 라우팅 키는 'archetype'(신규 슬러그). 'variant' 도 별칭 허용.
        arch = sl.get("archetype") or sl.get("variant", "statement")
        data = sl.get("data", {})
        fn = ARCHETYPE_TO_BUILDER.get(arch)
        if fn is None:
            warnings.append(f"slide {i+1}: unknown archetype '{arch}' -> fallback statement")
            fn = build_statement
        try:
            fn(prs, data)
        except Exception as e:
            warnings.append(f"slide {i+1} ({arch}): build error {e} -> blank")
            new_slide(prs)
        # 장식 레이어 포스트패스(grammar.decor 선언 룩만 — 무선언 룩은 무회귀).
        # 방금 만든 슬라이드(prs.slides[-1]) 위 가장자리에 도형을 흩뿌린다.
        if GRAMMAR.get("decor") == "memphis" and len(prs.slides._sldIdLst) > 0:
            try:
                _memphis_decor(prs.slides[len(prs.slides._sldIdLst) - 1], i)
            except Exception as e:
                warnings.append(f"slide {i+1}: decor error {e}")
    try:
        prs.save(out_path)
    except PermissionError:
        # 산출 .pptx 가 PowerPoint(또는 잔존 COM 렌더)에서 열려 있으면 쓰기가 거부된다.
        # 조용히 실패해 옛 파일이 그대로 남는 함정을 막기 위해 명확히 알린다.
        sys.stderr.write(
            f"[ppt-lab] ERROR: '{out_path}' 쓰기 거부(PermissionError) — "
            f"PowerPoint 등에서 파일이 열려 있으면 닫고 다시 빌드하세요.\n")
        raise SystemExit(2)
    return applied, warnings


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("out", help="output .pptx path")
    ap.add_argument("--palette", type=int, default=None,
                    help="palette id (valid: 1,2,3,4,5,6,8 — there is no 7)")
    ap.add_argument("--style", default=None, help="style preset id (default: house)")
    ap.add_argument("--look", default=None, help="design-pick look slug (self-contained style+palette; layout stays orthogonal)")
    ap.add_argument("--size", default=None,
                    help="type size tier: 소/중/대 (s/m/l). orthogonal to look. omit=look default(소)")
    ap.add_argument("--spec", default=None, help="spec.json path (omit for demo)")
    ap.add_argument("--chart-backdrop-opacity", type=float, default=None,
                    help="차트 가독성 백드롭 불투명도 0~100 (일관성↔가독성 HITL 노브; 미설정=가독 바닥선)")
    args = ap.parse_args()

    global SPEC_DIR
    if args.spec:
        spec = json.load(open(args.spec, encoding="utf-8"))
        SPEC_DIR = os.path.dirname(os.path.abspath(args.spec))  # 미디어 상대경로 기준
    else:
        spec = DEMO_SPEC

    applied, warnings = build(spec, args.out, args.palette, args.style, args.look, args.size,
                              args.chart_backdrop_opacity)
    pal_name = applied["name"] if applied else "base (#2563EB)"
    print(f"[ppt-lab] built {len(spec.get('slides', []))} slides -> {args.out}")
    if args.look or spec.get("look"):
        print(f"[ppt-lab] look:    {args.look or spec.get('look')}")
    print(f"[ppt-lab] palette: {pal_name}")
    print(f"[ppt-lab] size:    {args.size or spec.get('size') or '소 (look default)'}")
    print(f"[ppt-lab] style:   {args.style or ('(from look)' if (args.look or spec.get('look')) else 'house')}  "
          f"(latin={ACTIVE_FONT['latin']} / ea={ACTIVE_FONT['ea']})")
    if warnings:
        print("[ppt-lab] WARNINGS:", file=sys.stderr)
        for w in warnings:
            print(f"[ppt-lab] WARNING: {w}", file=sys.stderr)


if __name__ == "__main__":
    main()
