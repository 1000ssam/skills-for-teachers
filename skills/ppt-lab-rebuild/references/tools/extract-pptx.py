#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
extract-pptx.py — PPTX 결정론 측정기 (ppt-lab 흡수 1단계)

맘에 드는 .pptx 한 장(또는 전체)을 읽어 도형의 좌표·크기·색·폰트·텍스트를
1920x1080 px 좌표로 정규화해 뽑아낸다. "측정"은 코드가, "추상화(역할 해석)"는 사람/AI가.

출력:
  - 화면: 슬라이드별 도형 표 + 색 히스토그램
  - <pptx>_extract.json: 기계가 읽을 측정 리포트 + suggested_layout / suggested_tokens

Usage:
  python3 extract-pptx.py deck.pptx            # 전체
  python3 extract-pptx.py deck.pptx --slide 3  # 특정 슬라이드 1장
"""
import sys, os, json, argparse, collections
from pptx import Presentation
from pptx.util import Emu

EMU_PER_PX = 9525  # 914400 EMU/inch ÷ 96 px/inch


def hexof(color):
    try:
        return "#" + str(color.rgb)
    except Exception:
        return None


def shape_fill(sp):
    try:
        f = sp.fill
        if f.type is not None and f.fore_color and f.fore_color.type is not None:
            return hexof(f.fore_color)
    except Exception:
        pass
    return None


def first_run_font(sp):
    try:
        for p in sp.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    sz = None
                    try: sz = round(r.font.size.pt / 0.75) if r.font.size else None  # pt->px
                    except Exception: pass
                    col = None
                    try: col = hexof(r.font.color) if r.font.color and r.font.color.type is not None else None
                    except Exception: pass
                    return {"font": r.font.name, "size_px": sz, "bold": r.font.bold, "color": col}
    except Exception:
        pass
    return {}


def kind(sp):
    try:
        if sp.has_table: return "table"
        if sp.has_chart: return "chart"
        if sp.shape_type is not None and "PICTURE" in str(sp.shape_type): return "image"
        if sp.has_text_frame and sp.text_frame.text.strip(): return "text"
        return str(sp.shape_type).split()[0].lower() if sp.shape_type else "shape"
    except Exception:
        return "shape"


def role_guess(s):
    """좌표·크기·폰트로 역할 1차 추정 (사람이 검수)."""
    y, h, fs = s["y"], s["h"], (s.get("size_px") or 0)
    if s["kind"] == "table": return "table"
    if s["kind"] == "chart": return "chart"
    if s["kind"] == "image": return "image"
    if fs >= 60: return "display/title"
    if fs >= 40: return "title"
    if fs >= 28: return "heading"
    if 18 <= fs < 28 and y < 200: return "eyebrow?"
    if fs and fs < 16: return "caption/label"
    return "body"


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--slide", type=int, default=None)
    a = ap.parse_args()
    prs = Presentation(a.pptx)

    sw_px = prs.slide_width / EMU_PER_PX
    sh_px = prs.slide_height / EMU_PER_PX
    sx = 1920.0 / sw_px      # 1920 정규화 스케일
    sy = 1080.0 / sh_px

    report = {"source": os.path.basename(a.pptx),
              "native_px": [round(sw_px), round(sh_px)],
              "normalized": [1920, 1080], "slides": []}
    colors = collections.Counter()

    for idx, slide in enumerate(prs.slides, 1):
        if a.slide and idx != a.slide:
            continue
        shapes = []
        for sp in slide.shapes:
            try:
                x = round(sp.left / EMU_PER_PX * sx); y = round(sp.top / EMU_PER_PX * sy)
                w = round(sp.width / EMU_PER_PX * sx); h = round(sp.height / EMU_PER_PX * sy)
            except Exception:
                continue
            info = {"kind": kind(sp), "x": x, "y": y, "w": w, "h": h,
                    "fill": shape_fill(sp)}
            info.update(first_run_font(sp))
            try:
                t = sp.text_frame.text.strip().replace("\n", " ") if sp.has_text_frame else ""
                info["text"] = (t[:40] + "…") if len(t) > 40 else t
            except Exception:
                info["text"] = ""
            info["role"] = role_guess(info)
            for c in (info.get("fill"), info.get("color")):
                if c: colors[c] += 1
            shapes.append(info)
        shapes.sort(key=lambda s: (s["y"], s["x"]))
        # 제안 좌표 블록 (이미지/배경 제외한 주요 도형)
        sug = {f"slot{i+1}": [s["x"], s["y"], s["w"], s["h"]]
               for i, s in enumerate(shapes) if s["w"] > 60 and s["h"] > 20}
        report["slides"].append({"index": idx, "shapes": shapes, "suggested_layout": sug})

    report["suggested_tokens"] = [{"hex": c, "count": n} for c, n in colors.most_common()]

    out = os.path.splitext(a.pptx)[0] + "_extract.json"
    json.dump(report, open(out, "w", encoding="utf-8"), ensure_ascii=False, indent=2)

    # 화면 요약
    for sl in report["slides"]:
        print(f"\n── 슬라이드 {sl['index']} ({len(sl['shapes'])} 도형) ──")
        for s in sl["shapes"]:
            fs = s.get("size_px") or "-"
            print(f"  [{s['role']:>14}] {s['kind']:>6} ({s['x']},{s['y']},{s['w']},{s['h']}) "
                  f"fill={s['fill'] or '-':>8} {fs}px  {s['text']}")
    print("\n── 색 히스토그램(빈도순) ──")
    for c, n in colors.most_common(12):
        print(f"  {c}  ×{n}")
    print(f"\n✅ 리포트 저장: {out}")
    print("→ 다음: suggested_layout 을 layouts.json 에 코드명으로 등록 + build_<코드>() 작성")


if __name__ == "__main__":
    main()
