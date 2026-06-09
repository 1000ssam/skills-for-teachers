#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build-template.py — PPTX builder (ppt-lab clean-room rebuild).

Generates 1920x1080 (20 x 11.25 inch) slides from a spec.json using python-pptx.
- Loads design-tokens.json + layouts.json (no hardcoded hex / fonts).
- 좌표는 전부 layouts.json (grid.json SSOT 파생). 매직넘버 금지.
- ARCHETYPE_TO_BUILDER dict routes an archetype slug -> builder function.
- 23 아키타입(7 family) + media() 미디어 슬롯 + chart() ~20 차트형.
- --palette N (valid ids: 1,2,3,4,5,6,8 — no 7) 팔레트 스왑 · --style / --look 직교(레이아웃 불변).
- Korean font fallback set on latin + east-asian + complex-script runs.

Usage:
    python3 build-template.py OUT.pptx [--palette N] [--style S] [--look L] [--spec spec.json]
    python3 build-template.py /tmp/demo.pptx              # 내장 데모(전 아키타입)

spec.json format:
    {
      "title": "...",
      "palette": 2,
      "slides": [
        { "archetype": "cover", "data": {...} },
        { "archetype": "bars",  "data": {...} },
        ...
      ]
    }
"""
import json
import os
import sys
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

REF_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # tools/ -> references/

# ---------------------------------------------------------------------------
# Token loading + palette
# ---------------------------------------------------------------------------
TOKENS = json.load(open(os.path.join(REF_DIR, "design-tokens.json"), encoding="utf-8"))
LAYOUTS = json.load(open(os.path.join(REF_DIR, "layouts.json"), encoding="utf-8"))

COLORS = dict(TOKENS["colors"])          # mutated by palette swap
PALETTES = TOKENS["palettes"]
TYPO = TOKENS["typography"]
SLIDE = TOKENS["slide"]
SP = TOKENS["spacing"]
RAD = TOKENS["radius"]
FONT_PRIMARY = TOKENS["fonts"]["primary"]
FONTS = TOKENS["fonts"]
STYLES = TOKENS.get("styles", {})
LOOKS = TOKENS.get("looks", {})          # design-pick 보존 룩(자기완결 style+palette)

# spec.json 의 디렉토리 — 미디어 상대경로 해소 기준(media() 헬퍼가 참조).
# main()/build()에서 spec 로드 시 설정한다.
SPEC_DIR = None

# Active style state — mutated by apply_style() like COLORS is by apply_palette().
# A "style" = the absorbed identity of a reference deck: its fonts + component look.
# Default split: Latin/numerals -> english face (Inter), Hangul -> primary (Pretendard).
# A second optional tier — *_display — is used for headline tiers (DISPLAY_TIERS) so a
# look can pair a display face (e.g. high-contrast 송명) for covers/headings with a
# readable body face. Falls back to the base latin/ea when a look omits `display`.
ACTIVE_FONT = {"latin": FONTS.get("english", FONT_PRIMARY), "ea": FONT_PRIMARY,
               "latin_display": FONTS.get("english", FONT_PRIMARY),
               "ea_display": FONT_PRIMARY}
# Tiers that render in the display face. Mirrors the headline pre-wrap set exactly:
# covers/section/slide-titles + big stat numbers. Sub-headings (h3/h4) stay on body.
DISPLAY_TIERS = ("display", "h1", "h2")
COMPONENTS = {}                          # active component presets (card/callout/...)
ACCENTS = ["blue", "orange"]             # active palette's accent ROLES (primary, secondary, ...); set by apply_palette
DARK = False                             # True when the active look declares a dark canvas; set by apply_look

PX_TO_PT = SLIDE["px_to_pt"]             # 0.75


def apply_palette(pal_id):
    """Overlay palette overrides onto base COLORS."""
    if not pal_id:
        return None
    for p in PALETTES:
        if p["id"] == int(pal_id):
            COLORS.update(p["overrides"])
            ACCENTS[:] = p.get("accents", ["blue", "orange"])
            return p
    return None


def accent_cycle(n):
    """n colors from the active palette's accent ROLES (primary→secondary→...),
    padding with primary shades when the palette defines fewer accents than n.

    Mono palette → mono; multi palette → multi — but ALWAYS inside the chosen
    deck's palette. No out-of-palette color can leak (the cause of the old
    stats-green bug). 'One deck, one DEFINED palette' (mono or multi)."""
    pad = ["blue-2", "blue-light", "blue-pale"]
    return [ACCENTS[i] if i < len(ACCENTS) else pad[(i - len(ACCENTS)) % len(pad)]
            for i in range(n)]


def apply_style(style_id):
    """Resolve a style preset (fonts + components). Deterministic per chosen deck.

    Falls back to the 'house' style when style_id is None or unknown, so every
    build has a coherent default. Latin/EA fonts are resolved separately: Latin
    follows the deck's typeface, Hangul (ea/cs) falls back to a Korean-capable
    font so absorbing a Western reference (e.g. Calibri) never tofu-boxes Korean.
    """
    st = STYLES.get(str(style_id)) if style_id else None
    if st is None:
        st = STYLES.get("house")
    if st is None:
        return None
    _apply_fonts(st.get("fonts", {}))
    COMPONENTS.clear()
    COMPONENTS.update(st.get("components", {}))
    return st


def _apply_fonts(f):
    """Set ACTIVE_FONT from a look/style `fonts` block. The optional nested
    `display: {latin, ea}` overrides the headline-tier faces; absent → same as
    body faces (so existing single-font looks are unchanged)."""
    base_latin = f.get("latin", FONTS.get("english", FONT_PRIMARY))
    base_ea = f.get("ea", FONT_PRIMARY)
    disp = f.get("display", {})
    ACTIVE_FONT["latin"] = base_latin
    ACTIVE_FONT["ea"] = base_ea
    ACTIVE_FONT["latin_display"] = disp.get("latin", base_latin)
    ACTIVE_FONT["ea_display"] = disp.get("ea", base_ea)


def apply_look(look_id):
    """Apply a self-contained design-pick look in one shot: fonts + components +
    palette (colors + accent roles). A look bundles a style and a palette so that
    `--look <slug>` reproduces that look's identity, while the LAYOUT axis (variant)
    stays orthogonal. Explicit --style / --palette can still override afterwards.
    Returns the look dict, or None if unknown."""
    lk = LOOKS.get(str(look_id))
    if lk is None:
        return None
    _apply_fonts(lk.get("fonts", {}))
    COMPONENTS.clear()
    COMPONENTS.update(lk.get("components", {}))
    pal = dict(lk.get("palette", {}))
    accents = pal.pop("accents", ["blue", "orange"])
    explicit_roles = set(pal.keys())   # roles the look set by hand (respected)
    COLORS.update(pal)
    ACCENTS[:] = accents
    resolve_dark(explicit_roles, force=lk.get("dark"))
    return lk


def resolve_dark(explicit_roles=frozenset(), force=None):
    """Decide DARK and fill the role palette to match the active canvas.
    A look is dark when its `canvas` token resolves to a dark color (or force=True).
    For a dark canvas, every role NOT explicitly set by the look is auto-derived
    from the canvas hue, so any look becomes readable dark by declaring one color.
    Also fixes on-accent in BOTH modes from accent luminance (light accents get
    dark text). Idempotent — safe to call again after a palette/style override."""
    global DARK
    canvas_hex = _resolve("canvas")
    accent_hex = _resolve(ACCENTS[0] if ACCENTS else "blue")
    DARK = bool(force) if force is not None else is_dark(canvas_hex)
    if DARK:
        for k, v in derive_dark_roles(canvas_hex, accent_hex).items():
            if k not in explicit_roles:
                COLORS[k] = v
    else:
        # Light mode: roles keep their base aliases; only correct on-accent so a
        # light accent (e.g. mint/lime) doesn't get invisible white text.
        if "on-accent" not in explicit_roles:
            COLORS["on-accent"] = "#0F172A" if _rel_lum(accent_hex) > 0.6 else "white"


# ---------------------------------------------------------------------------
# Unit + color helpers
# ---------------------------------------------------------------------------
def IN(px):
    return Inches(px / 96.0)


def PT_PX(px):
    """px (Figma) -> pt for python-pptx."""
    return Pt(px * PX_TO_PT)


def _resolve(token, _depth=0):
    """Resolve a token through one or more alias hops to a raw #hex.
    Role tokens (ink/muted/surface/canvas…) alias other tokens, so a dark look
    can repoint a role in its palette and every builder that reads the role flips
    at once. Stops at a value that is a #hex literal (or any non-key string)."""
    v = COLORS.get(token, token)
    if isinstance(v, str) and not v.startswith("#") and v in COLORS and _depth < 8:
        return _resolve(v, _depth + 1)
    return v


def C(token):
    """token name, role alias, or raw #hex -> RGBColor."""
    if token is None:
        token = "ink"
    hexv = _resolve(token)
    return RGBColor.from_string(hexv.lstrip("#").upper())


# ---------------------------------------------------------------------------
# Dark-canvas auto-derivation — generalizes dark mode to ANY look.
# A look declares ONE color: `canvas` (its content background). Everything else
# (ink ladder, surfaces, dividers, on-accent text) is derived deterministically
# from that canvas hue so the surface palette always MATCHES the deck's
# background — no per-look hand tuning. Light looks keep the base aliases.
# ---------------------------------------------------------------------------
def _hex2rgb(h):
    h = str(h).lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _rgb2hex(r, g, b):
    return "#%02X%02X%02X" % tuple(
        max(0, min(255, int(round(c)))) for c in (r, g, b))


def _rel_lum(hexv):
    """WCAG relative luminance 0..1 from a #hex."""
    def lin(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = _hex2rgb(hexv)
    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def _mix(a, b, t):
    """Linear blend a->b by t in [0,1], in sRGB space (good enough for UI tints)."""
    ar, ag, ab = _hex2rgb(a)
    br, bg, bb = _hex2rgb(b)
    return _rgb2hex(ar + (br - ar) * t, ag + (bg - ag) * t, ab + (bb - ab) * t)


def is_dark(hexv):
    """A canvas is dark when its luminance sits below mid — text must invert."""
    try:
        return _rel_lum(hexv) < 0.30
    except Exception:
        return False


def derive_dark_roles(canvas_hex, accent_hex):
    """Deterministically build the full role palette for a dark canvas.
    Ink lifts toward white at decreasing strength (primary→tertiary); surfaces
    lift the canvas itself by small amounts (so cards read as elevated panels of
    the SAME hue); dividers/hairlines lift a touch more. on-* text inverts to the
    canvas so chips/accents stay legible. Accent-on-light is handled by luminance:
    when the accent is bright, text on it goes dark."""
    W = "#FFFFFF"
    roles = {
        "ink": W,
        "ink-2": _mix(canvas_hex, W, 0.85),
        "body": _mix(canvas_hex, W, 0.70),
        "muted": _mix(canvas_hex, W, 0.56),
        "subtle": _mix(canvas_hex, W, 0.44),
        "surface": _mix(canvas_hex, W, 0.07),
        "surface-2": _mix(canvas_hex, W, 0.11),
        "surface-3": _mix(canvas_hex, W, 0.17),
        "band": _mix(canvas_hex, W, 0.11),
        "divider": _mix(canvas_hex, W, 0.17),
        "hairline": _mix(canvas_hex, W, 0.30),
        # Emphasis slab (versus card / table header): in light mode a near-black
        # navy block; in dark mode a *subtly elevated dark* panel (NOT inverted to
        # white, which would scream on a dark deck). Text on it stays light.
        "slab": _mix(canvas_hex, W, 0.05),
        "on-slab": W,
        "on-ink": canvas_hex,
    }
    roles["on-accent"] = canvas_hex if _rel_lum(accent_hex) > 0.55 else W
    return roles


def style(key):
    """typography key -> (pt, bold)."""
    s = TYPO[key]
    bold = bool(s.get("bold") or s.get("medium"))
    return PT_PX(s["size_px"]), bold, s.get("line", 1.2), bool(s.get("upper"))


# ---------------------------------------------------------------------------
# Korean-safe pre-wrap (headline tiers only)
# ---------------------------------------------------------------------------
# PowerPoint wraps CJK text at ANY character boundary, so a Korean 어절 like
# "한다" can split as "한"/"다" at a line end. The OOXML eaLnBrk attribute that
# should control this is ignored by PowerPoint's render engine (confirmed
# empirically) and U+2060 word-joiners are ignored too — so the ONLY reliable
# fix is to measure the rendered width with the real font and insert explicit
# line breaks at whitespace (어절) boundaries BEFORE the text reaches PowerPoint.
# Applied to headline tiers (display/h1/h2) only, where mid-word breaks are most
# visible; body/caption keep PowerPoint's native wrap. Pure measurement — the
# rendered runs are unchanged. Units are Figma px: PX_TO_PT=0.75 makes the
# 96-DPI export pixel size equal the token size_px, and layout boxes are Figma
# px too, so box width and font px share one space (no DPI juggling).
import os as _os
# Optional font warehouse for Korean headline pre-wrap measurement. Point the
# PPT_LAB_FONT_DIR env var at a folder of per-family subdirs (e.g. Inter/,
# Song_Myung/) each holding .ttf/.otf files. When unset or missing, the pre-wrap
# degrades gracefully to PowerPoint's native CJK wrap — builds are unaffected.
_FONT_DIR = _os.environ.get("PPT_LAB_FONT_DIR", "")
_PIL_FONT_CACHE = {}
_FONT_PATH_CACHE = {}


def _font_path(face, bold):
    key = (face, bold)
    if key in _FONT_PATH_CACHE:
        return _FONT_PATH_CACHE[key]
    path = None
    for folder in (face.replace(" ", "_"), face.split(" ")[0]):
        d = _os.path.join(_FONT_DIR, folder)
        if _os.path.isdir(d):
            ttfs = [f for f in _os.listdir(d) if f.lower().endswith((".ttf", ".otf"))]
            pool = [f for f in ttfs if "italic" not in f.lower()] or ttfs
            want = "bold" if bold else "regular"
            pick = next((f for f in pool if want in f.lower()
                         and "extrabold" not in f.lower()
                         and "semibold" not in f.lower()), None)
            pick = pick or next((f for f in pool if want in f.lower()), None)
            pick = pick or (pool[0] if pool else None)
            if pick:
                path = _os.path.join(d, pick)
                break
    if path is None and face != "Inter":  # proprietary/missing → proportional fallback
        path = _font_path("Inter", bold)
    _FONT_PATH_CACHE[key] = path
    return path


def _pil_font(face, bold, px):
    key = (face, bold, px)
    if key not in _PIL_FONT_CACHE:
        from PIL import ImageFont
        p = _font_path(face, bold)
        _PIL_FONT_CACHE[key] = ImageFont.truetype(p, px) if p else ImageFont.load_default()
    return _PIL_FONT_CACHE[key]


def _is_cjk(ch):
    o = ord(ch)
    return (0xAC00 <= o <= 0xD7A3 or 0x3130 <= o <= 0x318F
            or 0x3000 <= o <= 0x303F or 0x4E00 <= o <= 0x9FFF
            or 0xFF00 <= o <= 0xFFEF)


def _measure(s, latin_face, ea_face, bold, px):
    """Width of s in px, routing CJK glyphs to the ea font, others to latin."""
    total, seg, seg_cjk = 0.0, "", None
    for ch in s:
        c = _is_cjk(ch)
        if seg and c != seg_cjk:
            total += _pil_font(ea_face if seg_cjk else latin_face, bold, px).getlength(seg)
            seg = ""
        seg += ch
        seg_cjk = c
    if seg:
        total += _pil_font(ea_face if seg_cjk else latin_face, bold, px).getlength(seg)
    return total


def _prewrap(content, box_w, size_px, bold, latin_face, ea_face, upper=False):
    """Greedy word-pack at whitespace so no 어절 splits mid-line. Returns a list
    of lines. On any failure (missing font, PIL error) returns [content] so the
    build never breaks; the worst case is PowerPoint's original native wrap."""
    try:
        px = max(1, int(round(size_px)))
        limit = box_w - 2  # tiny safety margin vs render rounding
        out = []
        for src in str(content).split("\n"):
            cur = ""
            for w in src.split(" "):
                trial = w if not cur else cur + " " + w
                m = _measure(trial.upper() if upper else trial,
                             latin_face, ea_face, bold, px)
                if not cur or m <= limit:
                    cur = trial
                else:
                    out.append(cur)
                    cur = w
            out.append(cur)
        return out or [content]
    except Exception:
        return [content]


# ---------------------------------------------------------------------------
# Low-level drawing primitives
# ---------------------------------------------------------------------------
def _set_run_font(run, latin=None, ea=None):
    """Resolve typeface from the active style: Latin glyphs use the deck's latin
    face, Hangul (east-asian + complex-script) uses the ea face (Korean fallback).
    Both default to the active style set by apply_style()."""
    latin = latin or ACTIVE_FONT["latin"]
    ea = ea or ACTIVE_FONT["ea"]
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", ea)


def text(slide, box, content, style_key, color="navy", align="left",
         anchor="top", wrap=True, line=None, upper=None, shrink=False):
    """Add a textbox. content may be a string or list of (string,) paragraphs."""
    x, y, w, h = box
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if shrink:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.NONE
    pt, bold, lh, up = style(style_key)
    if upper is not None:
        up = upper
    # Headline tiers render in the display face (falls back to body face when the
    # look declares no `display` font), so covers/headings can use a display serif.
    is_disp = style_key in DISPLAY_TIERS
    lat = ACTIVE_FONT["latin_display"] if is_disp else ACTIVE_FONT["latin"]
    ea = ACTIVE_FONT["ea_display"] if is_disp else ACTIVE_FONT["ea"]
    # Korean-safe pre-wrap for headline tiers: break at 어절 boundaries so a word
    # never splits mid-line (PowerPoint's CJK char-level wrap, eaLnBrk ignored).
    if wrap and is_disp and isinstance(content, str):
        content = _prewrap(content, w, TYPO[style_key]["size_px"], bold,
                           lat, ea, upper=up)
    al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT}[align]
    lines = content if isinstance(content, list) else [content]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        p.line_spacing = line or lh
        s = (str(ln).upper() if up else str(ln))
        r = p.add_run()
        r.text = s
        r.font.size = pt
        r.font.bold = bold
        r.font.color.rgb = C(color)
        _set_run_font(r, lat, ea)
    return tb


def bullets(slide, box, items, style_key="body", color="ink-2",
            marker=True, marker_color="blue", gap=1.5):
    x, y, w, h = box
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    pt, bold, lh, _ = style(style_key)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap
        p.space_after = Pt(6)
        if marker:
            rb = p.add_run()
            rb.text = "•  "
            rb.font.size = pt
            rb.font.bold = True
            rb.font.color.rgb = C(marker_color)
            _set_run_font(rb)
        r = p.add_run()
        r.text = str(it)
        r.font.size = pt
        r.font.bold = bold
        r.font.color.rgb = C(color)
        _set_run_font(r)
    return tb


def rect(slide, box, fill=None, line_color=None, line_w=0, radius=0):
    x, y, w, h = box
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp_type, IN(x), IN(y), IN(w), IN(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line_color:
        sp.line.color.rgb = C(line_color)
        sp.line.width = Pt(line_w or 1)
    else:
        sp.line.fill.background()
    if radius:
        try:
            sp.adjustments[0] = min(0.5, (radius / 96.0) / (min(w, h) / 96.0))
        except Exception:
            pass
    return sp


def shadow(sp, blur_pt=12, dist_pt=4, dir_deg=90, color="#000000", alpha=78):
    """Inject an outer drop shadow via OOXML — python-pptx's high-level API does
    not expose outerShdw (the same reason extract-style.py reads it via lxml).
    blur/dist in pt, dir in degrees, alpha 0..100 (opacity). effectLst is appended
    after <a:ln> per the CT_ShapeProperties child order."""
    spPr = sp._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur_pt * 12700)),
        "dist": str(int(dist_pt * 12700)),
        "dir": str(int(dir_deg * 60000)),
        "rotWithShape": "0",
    })
    clr = sh.makeelement(qn("a:srgbClr"), {"val": color.lstrip("#").upper()})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - alpha) * 1000))})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)
    return sp


def _accent_edge(slide, box, color, edge):
    """Draw a flat accent bar on one side of a card. The bar's SIDE+width come from
    the active style's `accent_edge` preset, so the *placement* of emphasis is a
    style property (top bar vs left rule) while the *color* stays the palette's
    accent role. Default = top/8px (reproduces the legacy hardcoded column bar)."""
    x, y, w, h = box
    t = edge.get("width_px", 8)
    side = edge.get("side", "top")
    bar = {"top": [x, y, w, t],
           "bottom": [x, y + h - t, w, t],
           "left": [x, y, t, h],
           "right": [x + w - t, y, t, h]}.get(side, [x, y, w, t])
    rect(slide, bar, fill=color, radius=0)


def card(slide, box, preset="card", fill=None, line_color=None, line_w=None,
         accent=None):
    """Draw a styled container from the active component preset (radius + border +
    shadow as one bundle). Falls back to a plain rect when the preset is absent,
    so this is safe even with a flat 'house' style. Per-call fill/line override
    the preset. This is the seam where the absorbed STYLE axis reaches the canvas.

    When `accent` (a color token) is given, an accent edge is drawn whose side+width
    come from the style preset's `accent_edge` (default top/8px). This lets an
    absorbed look move emphasis from a top bar to a left rule without touching
    layout code — the consulting (design-pick) look uses left/4px."""
    c = COMPONENTS.get(preset, {})
    rad = c.get("radius", 0)
    rad = RAD.get(rad, rad) if isinstance(rad, str) else (rad or 0)
    border = c.get("border") or {}
    lc = line_color if line_color is not None else border.get("color")
    lw = line_w if line_w is not None else border.get("width_pt", 0)
    fillv = fill if fill is not None else c.get("fill")
    sp = rect(slide, box, fill=fillv, line_color=lc, line_w=lw or 0, radius=rad)
    sh = c.get("shadow")
    if sh:
        shadow(sp, sh.get("blur_pt", 12), sh.get("dist_pt", 4),
               sh.get("dir_deg", 90), sh.get("color", "#000000"),
               sh.get("alpha", 78))
    if accent is not None:
        _accent_edge(slide, box, accent, c.get("accent_edge") or {"side": "top", "width_px": 8})
    return sp


def line(slide, x1, y1, x2, y2, color="slate-300", w=1):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    IN(x1), IN(y1), IN(x2), IN(y2))
    cn.line.color.rgb = C(color)
    cn.line.width = Pt(w)
    return cn


def oval(slide, box, fill=None, line_color=None, line_w=0):
    x, y, w, h = box
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(x), IN(y), IN(w), IN(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line_color:
        sp.line.color.rgb = C(line_color)
        sp.line.width = Pt(line_w or 1)
    else:
        sp.line.fill.background()
    return sp


def shape_text(sp, content, style_key, color="white", align="center",
               anchor="middle", upper=None, wrap=True):
    tf = sp.text_frame
    tf.word_wrap = wrap
    # No-wrap chips (e.g. number badges) must keep the glyphs on one line: the
    # default ~0.1" L/R text margins shrink a small box's usable width enough to
    # wrap a 2-digit number like "02" onto two lines. Zero them so the run stays
    # single-line and optically centered. (bug E)
    if not wrap:
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    pt, bold, lh, up = style(style_key)
    if upper is not None:
        up = upper
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = str(content).upper() if up else str(content)
    r.font.size = pt
    r.font.bold = bold
    r.font.color.rgb = C(color)
    _set_run_font(r)


def slide_bg(slide, color):
    # Skip when the (resolved) color is white — a white canvas needs no fill.
    # Resolve first so the role token "canvas" (=white in light mode) is a no-op.
    if _resolve(color).upper() in ("#FFFFFF", "WHITE"):
        return
    rect(slide, [0, 0, SLIDE["width_px"], SLIDE["height_px"]], fill=color)


def header(slide, data, dark=None):
    """헤더밴드: kicker(eyebrow) + title + divider. 좌표는 grid.json 파생
    (_meta.header_band). 슬롯명은 archetypes.md/grid.json 의 kicker/title/divider.
    data 는 'kicker' 또는 'eyebrow'(별칭) 둘 다 허용한다.
    dark=None 이면 활성 룩의 DARK 플래그를 따른다 → 다크 캔버스에서 제목·eyebrow·
    디바이더 색을 한 곳에서 자동 플립(17 헤더 아키타입 공용 시맨)."""
    if dark is None:
        dark = DARK
    hb = LAYOUTS["_meta"]["header_band"]
    ti = "ink"
    eb = "blue" if not dark else "blue-light"
    kick = data.get("kicker") or data.get("eyebrow")
    if kick:
        text(slide, hb["kicker"], kick, "eyebrow", color=eb, upper=True)
    if data.get("title"):
        text(slide, hb["title"], data["title"], "h2", color=ti)
    dv = hb["divider"]
    line(slide, dv[0], dv[1], dv[0] + dv[2], dv[1], color="divider", w=1.5)


def footer(slide, data):
    """선택적 푸터밴드: caption(좌) + pageno(우). 좌표는 grid.json 파생."""
    fb = LAYOUTS["_meta"]["footer_band"]
    if data.get("caption"):
        text(slide, fb["caption"], data["caption"], "caption", color="subtle")
    if data.get("pageno") is not None:
        text(slide, fb["pageno"], str(data["pageno"]), "caption",
             color="subtle", align="right")


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Dark canvas seam: paint the content background ONCE here so every archetype
    # (statement·duo·trio·grid·bars·share·…) gets a dark canvas without each
    # builder calling slide_bg. In light mode DARK=False → no-op (white canvas).
    # Frame archetypes (cover/section/closing) draw their own bg over this.
    if DARK:
        slide_bg(s, "canvas")
    return s


# ---------------------------------------------------------------------------
# Archetype builders (24 슬러그) — 좌표는 전부 layouts.json(grid.json 파생).
# 매직넘버 금지: 슬롯 박스에서 패딩만 상수로 빼서 파생한다.
# ---------------------------------------------------------------------------
PAD = 32          # 카드 내부 패딩(px) — 박스 안에서만 쓰는 로컬 상수
GAP = 24          # 카드/셀 사이 간격(px)


def _slots(slug):
    return LAYOUTS["archetypes"][slug]["slots"]


def _split_cols(box, n, gap=GAP):
    """box[x,y,w,h] 를 n 개 가로 컬럼으로 등분. 각 컬럼 박스 리스트 반환."""
    x, y, w, h = box
    if n <= 0:
        return []
    cw = (w - gap * (n - 1)) / n
    return [[x + i * (cw + gap), y, cw, h] for i in range(n)]


def _split_rows(box, n, gap=12):
    """box 를 n 개 가로 행으로 등분."""
    x, y, w, h = box
    if n <= 0:
        return []
    rh = (h - gap * (n - 1)) / n
    return [[x, y + i * (rh + gap), w, rh] for i in range(n)]


def _split_grid(box, ncol, nrow, gap=GAP):
    """box 를 ncol×nrow 셀로 분할. 행우선(row-major) 셀 박스 리스트."""
    x, y, w, h = box
    cw = (w - gap * (ncol - 1)) / ncol
    ch = (h - gap * (nrow - 1)) / nrow
    out = []
    for r in range(nrow):
        for c in range(ncol):
            out.append([x + c * (cw + gap), y + r * (ch + gap), cw, ch])
    return out


# ---------------------------------------------------------------------------
# 미디어 슬롯 헬퍼 — 이미지 또는 플레이스홀더.
# ---------------------------------------------------------------------------
def _resolve_src(src):
    """이미지 경로 해소. 절대경로면 그대로, 아니면 spec 디렉토리 기준.
    (WSL→Windows 공유 시 /mnt/c/... 마운트 경로.) 존재하지 않으면 None."""
    if not src:
        return None
    if not os.path.isabs(src):
        src = os.path.join(SPEC_DIR or os.getcwd(), src)
    return src if os.path.exists(src) else None


def media(slide, box, spec, fit=None):
    """미디어 슬롯 채우기. spec={src,fit,caption,focal}.
    - src 있고 파일 존재 → add_picture, box 에 맞춤(cover=크롭 / contain=레터박스).
    - 없으면 → 플레이스홀더(slate-50 + 1px 테두리 + 중앙 라벨 + 비율힌트, 이모지 금지).
    - caption 있으면 박스 하단에 작은 캡션.
    좌표는 호출자가 layouts.json 슬롯에서 넘긴 box 그대로 사용(매직넘버 없음)."""
    x, y, w, h = box
    spec = spec or {}
    fit = fit or spec.get("fit", "cover")
    src = _resolve_src(spec.get("src"))
    cap = spec.get("caption")
    cap_h = 28 if cap else 0
    img_h = h - cap_h
    if src:
        try:
            _place_picture(slide, [x, y, w, img_h], src, fit, spec.get("focal", "center"))
        except Exception as e:
            _media_placeholder(slide, [x, y, w, img_h],
                               f"이미지 로드 실패: {os.path.basename(src)}")
    else:
        _media_placeholder(slide, [x, y, w, img_h], spec.get("label", "이미지 자리"),
                           ratio=f"{w:.0f}×{img_h:.0f}")
    if cap:
        text(slide, [x, y + h - cap_h, w, cap_h], cap, "caption",
             color="muted")


def _place_picture(slide, box, src, fit, focal):
    """실제 이미지 삽입. cover=박스 꽉 채우고 크롭, contain=비율유지 레터박스.
    PIL 로 원본 비율을 읽어 결정론적으로 크롭/레터박스 계산한다."""
    x, y, w, h = box
    try:
        from PIL import Image
        with Image.open(src) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = w, h
    box_ar = w / h if h else 1
    img_ar = iw / ih if ih else 1
    if fit == "contain":
        # 비율 유지 레터박스: 박스 안에 들어가도록 축소, 가운데 정렬
        if img_ar > box_ar:
            nw = w; nh = w / img_ar
        else:
            nh = h; nw = h * img_ar
        px = x + (w - nw) / 2
        py = y + (h - nh) / 2
        slide.shapes.add_picture(src, IN(px), IN(py), IN(nw), IN(nh))
    else:
        # cover: 박스를 꽉 채우고 넘치는 부분 크롭. add_picture 후 crop_* 로 잘라낸다.
        pic = slide.shapes.add_picture(src, IN(x), IN(y), IN(w), IN(h))
        if img_ar > box_ar:
            # 이미지가 더 넓음 → 좌우 크롭
            crop = (1 - box_ar / img_ar) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_ar < box_ar:
            # 이미지가 더 높음 → 상하 크롭(focal 로 미세조정)
            crop = (1 - img_ar / box_ar) / 2
            if focal == "top":
                pic.crop_top = 0; pic.crop_bottom = crop * 2
            elif focal == "bottom":
                pic.crop_top = crop * 2; pic.crop_bottom = 0
            else:
                pic.crop_top = crop; pic.crop_bottom = crop
        return pic


def _media_placeholder(slide, box, label="이미지 자리", ratio=None):
    """이미지 없을 때 자리표시. slate-50 채움 + 1px 테두리 + 중앙 라벨 + 단순 도형
    아이콘(이모지 금지: 산 모양 삼각형 + 해 원 으로 '사진' 픽토그램)."""
    x, y, w, h = box
    rect(slide, box, fill="surface-2", line_color="divider", line_w=1, radius=0)
    # 픽토그램(도형만): 작은 원(해) + 삼각형(산) — 박스 중앙 위쪽
    icon = min(w, h) * 0.18
    cx, cy = x + w / 2, y + h / 2 - icon * 0.4
    oval(slide, [cx - icon * 0.55, cy - icon * 0.7, icon * 0.5, icon * 0.5],
         fill="hairline")
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 IN(cx - icon * 0.6), IN(cy - icon * 0.2),
                                 IN(icon * 1.4), IN(icon * 0.9))
    tri.shadow.inherit = False
    tri.fill.solid(); tri.fill.fore_color.rgb = C("hairline")
    tri.line.fill.background()
    # 라벨 + 비율 힌트
    lbl = label + (f"\n{ratio}" if ratio else "")
    text(slide, [x, y + h / 2 + icon * 0.6, w, 50], lbl, "caption",
         color="subtle", align="center", anchor="top")


# ---------------------------------------------------------------------------
# 차트 헬퍼 — 약 20종 디스패치(native python-pptx + 도형 합성).
# ---------------------------------------------------------------------------
# Categorical series color = the deck's DECLARED accents first (ACCENTS roles),
# then the primary accent's own tint ladder as padding. So:
#   - mono palette (1 accent)  -> [blue, blue-light, blue-2, blue-pale]  (= old
#     behavior exactly: one coherent tint ladder, strong 2-series contrast)
#   - dual+ palette/look       -> [blue, orange, blue-light, …]  (the look's 2nd
#     accent differentiates series 2 — e.g. ppt-dark-tech mint vs violet)
# This honors a look's deliberate multi-accent identity (esp. line/slope charts
# where color is the ONLY differentiator) while never leaking OFF-palette tokens:
# only ACCENTS (palette-defined) + primary-accent tints are ever used — the old
# green/slate/navy bleed bug came from raw semantic tokens, not from this.
# C() resolves each token against the look-mutated COLORS. >ramp series cycle.
def chart_color(i):
    """i-th chart series color token (accent-first, tint-padded). See note above."""
    pad = ["blue-light", "blue-2", "blue-pale"]
    return ACCENTS[i] if i < len(ACCENTS) else pad[(i - len(ACCENTS)) % len(pad)]

_NATIVE = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column_stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "radar": XL_CHART_TYPE.RADAR,
}


def _theme_axes(ch):
    """Flip native-chart axis ink for canvas contrast. Tick labels -> ink,
    axis lines + gridlines -> divider. Resolves the active role tokens, so in
    light mode this reproduces the dark-on-white default and in dark mode it
    lifts labels/lines to readable values. python-pptx's high-level API exposes
    these, so no OOXML surgery needed. Best-effort per axis (some chart types
    lack a value/category axis)."""
    ink = C("ink")
    div = C("divider")
    for getter in ("category_axis", "value_axis"):
        try:
            ax = getattr(ch, getter)
        except Exception:
            continue
        try:
            ax.tick_labels.font.color.rgb = ink
        except Exception:
            pass
        try:
            ax.format.line.color.rgb = div
        except Exception:
            pass
        try:
            if ax.has_major_gridlines:
                ax.major_gridlines.format.line.color.rgb = div
        except Exception:
            pass


def _theme_legend(ch):
    """Legend text -> ink so series names stay readable on a dark canvas."""
    try:
        ch.legend.font.color.rgb = C("ink")
    except Exception:
        pass


def chart(slide, box, spec):
    """차트 디스패처. spec={type, categories, series:[{name,values}], values, ...}.
    native(python-pptx) 12종 + 도형 합성 8종. 팔레트 램프로 시리즈 색 지정.
    빈 데이터면 안내 플레이스홀더. 좌표는 호출자가 넘긴 box."""
    spec = spec or {}
    ctype = spec.get("type", "column_clustered")
    cats = spec.get("categories", [])
    series = spec.get("series")
    if not series and "values" in spec:
        series = [{"name": spec.get("series_name", ""), "values": spec["values"]}]
    series = series or []
    # Normalized series must be written back: _chart_native reads spec["series"]
    # directly, so a chart given a bare `values` list (e.g. spread/histogram)
    # would otherwise render with ZERO series — an empty plot area. (bug C)
    spec["series"] = series
    has_data = bool(series) and any(s.get("values") for s in series)

    if ctype in ("xy_scatter", "bubble"):
        return _chart_xy(slide, box, spec, ctype)
    composed = {
        "waterfall": _chart_waterfall, "funnel": _chart_funnel,
        "gauge": _chart_gauge, "progress": _chart_gauge,
        "gantt": _chart_gantt, "timeline": _chart_gantt,
        "tam_sam_som": _chart_tam_sam_som, "bullet": _chart_bullet,
        "slope": _chart_slope, "kpi_cards": _chart_kpi_cards,
    }
    if ctype in composed:
        return composed[ctype](slide, box, spec)
    if ctype in ("pie", "doughnut"):
        return _chart_pie(slide, box, spec, ctype)
    if not has_data:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    return _chart_native(slide, box, spec, ctype)


def _chart_native(slide, box, spec, ctype):
    """CategoryChartData 기반 native 차트. 시리즈색=팔레트 램프."""
    x, y, w, h = box
    xl = _NATIVE.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
    cd = CategoryChartData()
    cd.categories = spec.get("categories", [])
    for i, s in enumerate(spec.get("series", [])):
        cd.add_series(s.get("name", f"S{i+1}"), s.get("values", []),
                      number_format=spec.get("number_format"))
    gf = slide.shapes.add_chart(xl, IN(x), IN(y), IN(w), IN(h), cd)
    ch = gf.chart
    multi = len(spec.get("series", [])) > 1
    ch.has_legend = spec.get("legend", multi)
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = PT_PX(TYPO["caption"]["size_px"])
        _theme_legend(ch)
    # Series color = brand ramp. Bars/areas read .format.fill; lines read
    # .format.line; markers need their own fill/line. Set ALL so EVERY chart type
    # follows the active palette (the off-brand steel-blue/red line defaults were
    # Office's, not ours). Each in its own try so one unsupported facet never
    # aborts the rest.
    for i, ps in enumerate(ch.series):
        col = C(chart_color(i))
        try:
            ps.format.fill.solid()
            ps.format.fill.fore_color.rgb = col
        except Exception:
            pass
        try:
            ps.format.line.color.rgb = col
        except Exception:
            pass
        try:
            ps.marker.format.fill.solid()
            ps.marker.format.fill.fore_color.rgb = col
            ps.marker.format.line.color.rgb = col
        except Exception:
            pass
    try:
        ch.category_axis.tick_labels.font.size = PT_PX(13)
        va = ch.value_axis
        va.tick_labels.font.size = PT_PX(13)
        va.has_major_gridlines = True
    except Exception:
        pass
    _theme_axes(ch)
    return ch


def _chart_pie(slide, box, spec, ctype):
    x, y, w, h = box
    cats = spec.get("categories", [])
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    cd = CategoryChartData()
    cd.categories = cats or [str(i + 1) for i in range(len(vals))]
    cd.add_series("", vals)
    xl = XL_CHART_TYPE.DOUGHNUT if ctype == "doughnut" else XL_CHART_TYPE.PIE
    gf = slide.shapes.add_chart(xl, IN(x), IN(y), IN(w), IN(h), cd)
    ch = gf.chart
    # Respect spec legend: callers that draw their OWN legend (e.g. build_share's
    # color-chip list) pass legend=False to avoid a duplicate native legend. (bug B)
    ch.has_legend = spec.get("legend", True)
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.RIGHT
        ch.legend.include_in_layout = False
        ch.legend.font.size = PT_PX(14)
        _theme_legend(ch)
    try:
        pts = ch.plots[0].series[0].points
        for i, p in enumerate(pts):
            p.format.fill.solid()
            p.format.fill.fore_color.rgb = C(chart_color(i))
    except Exception:
        pass
    return ch


def _chart_xy(slide, box, spec, ctype):
    """스캐터/버블. XyChartData / BubbleChartData 사용.
    spec.points=[{x,y,(size),label}] 또는 series 형식."""
    x, y, w, h = box
    pts = spec.get("points", [])
    if not pts:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    if ctype == "bubble":
        from pptx.chart.data import BubbleChartData
        cd = BubbleChartData()
        sr = cd.add_series("")
        for p in pts:
            sr.add_data_point(p.get("x", 0), p.get("y", 0), p.get("size", 1))
        xl = XL_CHART_TYPE.BUBBLE
    else:
        from pptx.chart.data import XyChartData
        cd = XyChartData()
        sr = cd.add_series("")
        for p in pts:
            sr.add_data_point(p.get("x", 0), p.get("y", 0))
        xl = XL_CHART_TYPE.XY_SCATTER
    gf = slide.shapes.add_chart(xl, IN(x), IN(y), IN(w), IN(h), cd)
    ch = gf.chart
    ch.has_legend = False
    try:
        ch.series[0].format.fill.solid()
        ch.series[0].format.fill.fore_color.rgb = C("blue")
    except Exception:
        pass
    return ch


# ----- 도형 합성 차트 (native 로 안 되는 8종, box 에서 결정론적으로 파생) -----
def _series_values(spec):
    """첫 시리즈 (cats, vals) 추출."""
    cats = spec.get("categories", [])
    vals = spec.get("values")
    if vals is None and spec.get("series"):
        vals = spec["series"][0].get("values", [])
    return cats, (vals or [])


def _chart_waterfall(slide, box, spec):
    """워터폴: 시작값 + 증감 막대 누적. spec.values 의 각 값이 증감,
    첫/마지막은 total 로 본다(spec.totals=[idx] 로 표시 가능)."""
    x, y, w, h = box
    cats, vals = _series_values(spec)
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    totals = set(spec.get("totals", [0, len(vals) - 1]))
    n = len(vals)
    bw = (w - GAP * (n - 1)) / n
    # 누적 범위 계산
    cum = 0
    lows, run = [], []
    for i, v in enumerate(vals):
        if i in totals:
            run.append((0, v)); cum = v
        else:
            start = cum; cum += v
            run.append((min(start, cum), max(start, cum)))
    vmax = max(hi for _, hi in run) or 1
    base_y = y + h - 28
    plot_h = h - 40
    for i, (lo, hi) in enumerate(run):
        bx = x + i * (bw + GAP)
        bh = (hi - lo) / vmax * plot_h
        by = base_y - (hi / vmax * plot_h)
        col = "blue" if i in totals else ("green" if vals[i] >= 0 else "orange")
        rect(slide, [bx, by, bw, max(2, bh)], fill=col, radius=RAD["sm"])
        if i < len(cats):
            text(slide, [bx - 4, base_y + 4, bw + 8, 24], str(cats[i]), "caption",
                 color="muted", align="center")


def _chart_funnel(slide, box, spec):
    """퍼널: 위→아래로 좁아지는 사다리꼴 막대(값 비례 폭)."""
    x, y, w, h = box
    cats, vals = _series_values(spec)
    if not vals:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    vmax = max(vals) or 1
    n = len(vals)
    rh = (h - 12 * (n - 1)) / n
    for i, v in enumerate(vals):
        bw = v / vmax * w
        bx = x + (w - bw) / 2
        by = y + i * (rh + 12)
        rect(slide, [bx, by, bw, rh], fill=chart_color(i),
             radius=RAD["sm"])
        lbl = (str(cats[i]) + "  " if i < len(cats) else "") + str(v)
        text(slide, [x, by + (rh - 20) / 2, w, 24], lbl, "body-md",
             color="on-accent", align="center")


def _chart_gauge(slide, box, spec):
    """게이지/진행률: 트랙 + 채움 막대 + 중앙 %텍스트. spec.percent(0~100)."""
    x, y, w, h = box
    pct = float(spec.get("percent", spec.get("value", 0)) or 0)
    pct = max(0, min(100, pct))
    # 도넛 게이지 대신 명시적 수평 진행바(읽기 쉬움)
    bar_h = min(48, h * 0.18)
    by = y + h / 2 + 10
    rect(slide, [x, by, w, bar_h], fill="band", radius=RAD["pill"])
    if pct > 0:
        rect(slide, [x, by, w * pct / 100, bar_h], fill="blue", radius=RAD["pill"])
    text(slide, [x, y + h / 2 - 90, w, 100], f"{pct:.0f}%", "display",
         color="blue", align="center")
    if spec.get("caption"):
        text(slide, [x, by + bar_h + 10, w, 40], spec["caption"], "body-lg",
             color="muted", align="center")


def _chart_gantt(slide, box, spec):
    """간트/타임라인: 행별 막대(start,end 0~1 또는 일수). spec.tasks=[{label,start,end}]."""
    x, y, w, h = box
    tasks = spec.get("tasks", [])
    if not tasks:
        cats, vals = _series_values(spec)
        tasks = [{"label": c, "start": 0, "end": v} for c, v in zip(cats, vals)]
    if not tasks:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    mx = max((t.get("end", 1) for t in tasks), default=1) or 1
    label_w = w * 0.25
    track_x = x + label_w
    track_w = w - label_w
    n = len(tasks)
    rh = (h - 10 * (n - 1)) / n
    for i, t in enumerate(tasks):
        ry = y + i * (rh + 10)
        text(slide, [x, ry, label_w - 12, rh], t.get("label", ""), "body-md",
             color="ink", anchor="middle")
        s = t.get("start", 0) / mx
        e = t.get("end", 1) / mx
        rect(slide, [track_x + s * track_w, ry + rh * 0.2,
                     max(4, (e - s) * track_w), rh * 0.6],
             fill=chart_color(i), radius=RAD["sm"])


def _chart_tam_sam_som(slide, box, spec):
    """TAM-SAM-SOM: 동심원 3겹 + 라벨. spec.bubbles=[{label,value}] (큰→작은)."""
    x, y, w, h = box
    items = spec.get("bubbles") or spec.get("layers") or []
    if not items:
        cats, vals = _series_values(spec)
        items = [{"label": c, "value": v} for c, v in zip(cats, vals)]
    items = items[:3]
    if not items:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    d = min(w * 0.55, h)
    cx = x + d / 2
    cy = y + h / 2
    fills = ["blue-pale", "blue-light", "blue"]
    sizes = [d, d * 0.66, d * 0.36]
    for i in range(len(items)):
        sz = sizes[i]
        oval(slide, [cx - sz / 2, cy - sz / 2, sz, sz], fill=fills[i])
    lx = x + d + 40
    lh = h / max(1, len(items))
    for i, it in enumerate(items):
        ly = y + i * lh
        text(slide, [lx, ly + 8, w - d - 60, 36],
             it.get("label", ""), "h4",
             color="blue" if fills[i] != "blue-pale" else "ink")
        if it.get("value") is not None:
            text(slide, [lx, ly + 50, w - d - 60, 44], str(it["value"]), "h2",
                 color="ink")


def _chart_bullet(slide, box, spec):
    """불릿 차트: 행별 (실측 막대 vs 목표 마커 + 정성 밴드). spec.rows=[{label,value,target,max}]."""
    x, y, w, h = box
    rows = spec.get("rows", [])
    if not rows:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    label_w = w * 0.22
    track_x = x + label_w
    track_w = w - label_w
    n = len(rows)
    rh = (h - 16 * (n - 1)) / n
    for i, r in enumerate(rows):
        ry = y + i * (rh + 16)
        mx = r.get("max", max(r.get("value", 0), r.get("target", 0)) * 1.25) or 1
        text(slide, [x, ry, label_w - 12, rh], r.get("label", ""), "body-md",
             color="ink", anchor="middle")
        # 정성 밴드(연한 배경)
        rect(slide, [track_x, ry + rh * 0.25, track_w, rh * 0.5],
             fill="band", radius=0)
        val = r.get("value", 0) / mx
        rect(slide, [track_x, ry + rh * 0.35, val * track_w, rh * 0.3],
             fill="blue", radius=0)
        if r.get("target") is not None:
            tx = track_x + r["target"] / mx * track_w
            rect(slide, [tx - 2, ry + rh * 0.2, 4, rh * 0.6], fill="ink")


def _chart_slope(slide, box, spec):
    """슬로프(기울기) 차트: 좌→우 두 시점 값 연결선. spec.series=[{name,values:[a,b]}]."""
    x, y, w, h = box
    series = spec.get("series", [])
    if not series:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    allv = [v for s in series for v in s.get("values", [])[:2]]
    if not allv:
        return _media_placeholder(slide, box, "차트 데이터 없음")
    vmin, vmax = min(allv), max(allv)
    span = (vmax - vmin) or 1
    lx, rx = x + w * 0.18, x + w * 0.82
    top, bot = y + 20, y + h - 30

    def py(v):
        return bot - (v - vmin) / span * (bot - top)
    cats = spec.get("categories", ["전", "후"])
    text(slide, [x, y + h - 26, w * 0.4, 24], str(cats[0] if cats else ""),
         "caption", color="muted", align="left")
    text(slide, [x + w * 0.6, y + h - 26, w * 0.4, 24],
         str(cats[1] if len(cats) > 1 else ""), "caption",
         color="muted", align="right")
    for i, s in enumerate(series):
        vals = s.get("values", [])
        if len(vals) < 2:
            continue
        col = chart_color(i)
        line(slide, lx, py(vals[0]), rx, py(vals[1]), color=col, w=2.5)
        oval(slide, [lx - 6, py(vals[0]) - 6, 12, 12], fill=col)
        oval(slide, [rx - 6, py(vals[1]) - 6, 12, 12], fill=col)
        text(slide, [x, py(vals[0]) - 14, w * 0.16, 28],
             f"{s.get('name','')} {vals[0]}", "caption", color="ink",
             align="right")
        text(slide, [rx + 8, py(vals[1]) - 14, w * 0.16, 28], str(vals[1]),
             "caption", color="ink")


def _chart_kpi_cards(slide, box, spec):
    """빅넘버 KPI 카드 N개(가로 분할). spec.kpis=[{label,value,note}]."""
    kpis = spec.get("kpis") or spec.get("stats") or []
    if not kpis:
        return _media_placeholder(slide, box, "데이터 없음")
    boxes = _split_cols(box, len(kpis))
    accents = accent_cycle(len(kpis))
    for i, k in enumerate(kpis):
        bx, by, bw, bh = boxes[i]
        card(slide, boxes[i], "card", accent=accents[i % len(accents)])
        text(slide, [bx + PAD, by + bh * 0.18, bw - 2 * PAD, 36],
             k.get("label", ""), "label", color="muted",
             align="center", upper=True)
        text(slide, [bx + PAD, by + bh * 0.38, bw - 2 * PAD, bh * 0.34],
             k.get("value", ""), "h1", color="ink", align="center")
        if k.get("note"):
            text(slide, [bx + PAD, by + bh * 0.74, bw - 2 * PAD, bh * 0.2],
                 k["note"], "body", color="muted", align="center")


# ---------------------------------------------------------------------------
# 카드 콘텐츠 헬퍼 (Set/Contrast 공용) — 이미지 카드 변종 지원.
# ---------------------------------------------------------------------------
def _content_card(slide, box, item, accent="blue"):
    """카드: (선택적 이미지 밴드 상단 ~55%) + label/heading/body 하단.
    item={image:{src,...}, label, heading, body|bullets}. 이미지 있으면 이미지 카드."""
    x, y, w, h = box
    # 'image' 키가 있으면(빈 dict {} 라도) 이미지 카드 — src 없으면 media()가
    # 플레이스홀더를 그린다. ({} 는 falsy 이므로 get() 대신 키 존재로 판정.)
    # fill/border 는 활성 style/look 의 components.card 프리셋이 결정한다
    # (하드코딩 금지). accent 엣지는 이미지 카드엔 그리지 않는다(이미지가 상단을
    # 덮으므로) — 텍스트 카드일 때만 프리셋의 accent_edge 로 그린다.
    if "image" in item:
        card(slide, box, "card")
        img = item.get("image") or {}
        band_h = h * 0.55
        media(slide, [x, y, w, band_h], img, fit=img.get("fit", "cover"))
        ty = y + band_h + 16
        th = h - band_h - 24
    else:
        card(slide, box, "card", accent=accent)
        ty = y + 24
        th = h - 40
    inx = x + PAD
    inw = w - 2 * PAD
    cy = ty
    if item.get("label"):
        text(slide, [inx, cy, inw, 26], item["label"], "label",
             color=accent, upper=True)
        cy += 32
    if item.get("heading"):
        text(slide, [inx, cy, inw, 44], item["heading"], "h3", color="ink")
        cy += 56
    rest = max(40, (ty + th) - cy)
    if item.get("bullets"):
        bullets(slide, [inx, cy, inw, rest], item["bullets"], "body",
                color="body", marker_color=accent)
    elif item.get("body"):
        text(slide, [inx, cy, inw, rest], item["body"], "body", color="body")


# ===========================================================================
# Family 1 — Frame
# ===========================================================================
def build_cover(prs, d):
    s = new_slide(prs)
    A = LAYOUTS["archetypes"]["cover"]
    L = A["slots"]
    bg = d.get("bg", "canvas" if DARK else "white")
    slide_bg(s, bg)
    dark = DARK or bg not in ("white", "slate-50", "blue-faint")
    if d.get("background"):
        media(s, [0, 0, SLIDE["width_px"], SLIDE["height_px"]], d["background"],
              fit="cover")
    ink = "white" if dark else "navy"
    sub = "slate-300" if dark else "slate-500"
    if d.get("kicker") or d.get("eyebrow"):
        text(s, L["kicker"], d.get("kicker") or d["eyebrow"], "eyebrow",
             color="blue-light" if dark else "blue", upper=True)
    text(s, L["title"], d.get("title", ""), "display", color=ink, anchor="top")
    if d.get("subtitle"):
        text(s, L["subtitle"], d["subtitle"], "h4", color=sub)
    if d.get("meta"):
        text(s, L["meta"], d["meta"], "caption", color="slate-400")
    return s


def build_section(prs, d):
    s = new_slide(prs)
    L = LAYOUTS["archetypes"]["section"]["slots"]
    bg = d.get("bg", "navy")
    slide_bg(s, bg)
    dark = bg not in ("white", "slate-50", "blue-faint")
    if d.get("background"):
        media(s, [0, 0, SLIDE["width_px"], SLIDE["height_px"]], d["background"],
              fit="cover")
    ink = "white" if dark else "navy"
    if d.get("index"):
        text(s, L["index"], str(d["index"]), "h1",
             color="blue-light" if dark else "blue")
    text(s, L["section_title"], d.get("title", d.get("section_title", "")),
         "h1", color=ink, anchor="top")
    if d.get("caption"):
        text(s, L["caption"], d["caption"], "body-lg",
             color="slate-300" if dark else "slate-500")
    return s


def build_agenda(prs, d):
    s = new_slide(prs)
    L = LAYOUTS["archetypes"]["agenda"]["slots"]
    slide_bg(s, "canvas")
    eb = "blue" if not DARK else "blue-light"
    if d.get("kicker") or d.get("eyebrow"):
        text(s, L["kicker"], d.get("kicker") or d.get("eyebrow", "AGENDA"),
             "eyebrow", color=eb, upper=True)
    text(s, L["title"], d.get("title", "목차"), "h1", color="ink")
    items = d.get("items", [])
    n = max(1, len(items))
    rows = _split_rows(L["items"], n, gap=8)
    BADGE = 56          # 번호 배지 폭
    TEXT_GAP = 24       # 배지와 텍스트 사이 간격
    COL_GAP = 40        # 타이틀/노트 열 사이 간격
    for i, it in enumerate(items):
        rx, ry, rw, rh = rows[i]
        tx = rx + BADGE + TEXT_GAP          # 텍스트 영역 좌측 시작
        avail = rw - (BADGE + TEXT_GAP)     # 배지 뒤 가용 폭
        numbox = rect(s, [rx, ry + (rh - BADGE) / 2, BADGE, BADGE], fill="blue",
                      radius=RAD["md"])
        shape_text(numbox, f"{i+1:02d}", "h4", color="on-accent", wrap=False)
        title = it if isinstance(it, str) else it.get("title", "")
        note = it.get("note") if isinstance(it, dict) else None
        if note:
            # 가용폭을 타이틀/노트로 균등 분할 + 갭. 노트가 ~48%(~800px) 확보 → 1줄.
            col_w = (avail - COL_GAP) / 2
            text(s, [tx, ry, col_w, rh], title, "h3", color="ink",
                 anchor="middle")
            text(s, [tx + col_w + COL_GAP, ry, col_w, rh], note, "body",
                 color="muted", anchor="middle", align="right")
        else:
            # 노트 없으면 타이틀이 가용폭 전부 사용
            text(s, [tx, ry, avail, rh], title, "h3", color="ink",
                 anchor="middle")
        if i < n - 1:
            line(s, rx, ry + rh + 4, rx + rw, ry + rh + 4, color="divider", w=1)
    return s


def build_closing(prs, d):
    s = new_slide(prs)
    L = LAYOUTS["archetypes"]["closing"]["slots"]
    bg = d.get("bg", "navy")
    slide_bg(s, bg)
    dark = bg not in ("white", "slate-50", "blue-faint")
    ink = "white" if dark else "navy"
    if d.get("kicker") or d.get("eyebrow"):
        text(s, L["kicker"], d.get("kicker") or d["eyebrow"], "eyebrow",
             color="blue-light" if dark else "blue", upper=True)
    text(s, L["title"], d.get("title", ""), "h1", color=ink)
    if d.get("subtitle"):
        text(s, L["subtitle"], d["subtitle"], "body-lg",
             color="slate-300" if dark else "slate-500")
    if d.get("cta"):
        btn = card(s, L["cta"], "card", fill="blue")
        shape_text(btn, d["cta"], "body-md", color="on-accent")
    return s


# ===========================================================================
# Family 2 — Focus
# ===========================================================================
def build_statement(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["statement"]["slots"]
    if d.get("lead"):
        # 리드는 헤더 제목(h2/48px)과 같은 크기로 찍으면 '헤더가 둘'처럼 보인다.
        # 풀쿼트로 처리: 좌측 액센트 바 + 더 큰 h1(64px) + 세로 중앙 정렬 →
        # 슬라이드의 주인공(핵심 문장/인용)으로 명확히 위계가 잡힌다.
        lx, ly, lw, lh = L["lead"]
        bar_w = 6
        rect(s, [lx, ly + 6, bar_w, lh - 12], fill="blue", radius=RAD["sm"])
        text(s, [lx + bar_w + 24, ly, lw - bar_w - 24, lh], d["lead"], "h1",
             color="ink", anchor="middle")
    if d.get("image"):
        media(s, L["image"], d["image"], fit=d["image"].get("fit", "cover"))
        # 이미지 있으면 근거는 좌측만(미디어가 우측 차지)
        sx, sy, sw, sh = L["support"]
        sup_box = [sx, sy, L["image"][0] - sx - GAP, sh]
    else:
        sup_box = L["support"]
    if d.get("support"):
        if isinstance(d["support"], list):
            bullets(s, sup_box, d["support"], "body-lg")
        else:
            text(s, sup_box, d["support"], "body-lg", color="ink-2")
    return s


def build_feature(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["feature"]["slots"]
    if d.get("bullets"):
        bullets(s, L["body"], d["bullets"], "body-lg")
    elif d.get("body"):
        text(s, L["body"], d["body"], "body-lg", color="ink-2")
    media(s, L["media"], d.get("media", {}), fit=(d.get("media") or {}).get("fit", "cover"))
    return s


# ===========================================================================
# Family 3 — Set
# ===========================================================================
def build_duo(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["duo"]["slots"]
    items = d.get("items") or [d.get("left", {}), d.get("right", {})]
    accents = accent_cycle(2)
    _content_card(s, list(L["item_l"]), items[0] if items else {}, accents[0])
    if len(items) > 1:
        _content_card(s, list(L["item_r"]), items[1], accents[1])
    return s


def build_trio(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["trio"]["slots"]
    items = d.get("items") or d.get("columns", [])
    accents = accent_cycle(3)
    for i, key in enumerate(["col1", "col2", "col3"]):
        if i < len(items):
            _content_card(s, list(L[key]), items[i], accents[i % 3])
    return s


def build_grid(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["grid"]["slots"]
    items = d.get("items") or d.get("cells", [])
    n = len(items)
    ncol = 3 if n > 4 else 2
    nrow = 2 if n > 2 else 1
    boxes = _split_grid(L["cells"], ncol, nrow)
    accents = accent_cycle(max(1, n))
    for i, it in enumerate(items):
        if i < len(boxes):
            _content_card(s, boxes[i], it, accents[i % len(accents)])
    return s


def build_list(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["list"]["slots"]
    items = d.get("items", [])
    n = max(1, len(items))
    rows = _split_rows(L["items"], n, gap=10)
    for i, it in enumerate(items):
        rx, ry, rw, rh = rows[i]
        num = oval(s, [rx, ry + (rh - 44) / 2, 44, 44], fill="blue")
        shape_text(num, str(i + 1), "h4", color="on-accent")
        title = it if isinstance(it, str) else it.get("title", "")
        text(s, [rx + 64, ry, rw - 64, rh], title, "h4", color="ink",
             anchor="middle")
        sub = it.get("body") if isinstance(it, dict) else None
        if sub:
            text(s, [rx + 64, ry + rh * 0.55, rw - 64, rh * 0.4], sub, "body",
                 color="muted")
    return s


# ===========================================================================
# Family 4 — Contrast
# ===========================================================================
def build_versus(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["versus"]["slots"]
    left = d.get("left", {})
    right = d.get("right", {})
    # 좌(어둡게) vs 우(밝게)
    lx, ly, lw, lh = L["left"]
    # 좌 카드는 "강조 타일"(잉크면 + 반전 텍스트) — 라이트=네이비/흰글씨,
    # 다크=흰면/어두운글씨. ink/on-ink 역할로 두 모드 모두 대비 보장.
    rect(s, L["left"], fill="slab", radius=RAD["lg"])
    text(s, [lx + PAD, ly + 28, lw - 2 * PAD, 30],
         left.get("label", "BEFORE"), "label", color="on-slab", upper=True)
    if left.get("heading"):
        text(s, [lx + PAD, ly + 64, lw - 2 * PAD, 56], left["heading"], "h3",
             color="on-slab")
    if left.get("points"):
        bullets(s, [lx + PAD, ly + 140, lw - 2 * PAD, lh - 170], left["points"],
                "body", color="on-slab", marker_color="red")
    rx, ry, rw, rh = L["right"]
    card(s, L["right"], "card", fill="surface-3", line_color="blue", line_w=1.5)
    text(s, [rx + PAD, ry + 28, rw - 2 * PAD, 30],
         right.get("label", "AFTER"), "label", color="blue", upper=True)
    if right.get("heading"):
        text(s, [rx + PAD, ry + 64, rw - 2 * PAD, 56], right["heading"], "h3",
             color="ink")
    if right.get("points"):
        bullets(s, [rx + PAD, ry + 140, rw - 2 * PAD, rh - 170], right["points"],
                "body", color="ink-2", marker_color="green")
    # 중앙 VS 마크
    vs = oval(s, [(lx + lw + rx) / 2 - 36, ly + lh / 2 - 36, 72, 72], fill="surface",
              line_color="hairline", line_w=1.5)
    shape_text(vs, "VS", "h4", color="ink")
    return s


def build_matrix(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["matrix"]["slots"]
    x, y, w, h = L["table"]
    cols = d.get("headers", [])
    rows = d.get("rows", [])
    ncol = max(1, len(cols))
    cw = w / ncol
    head_h = 60
    rh = min(72, (h - head_h) / max(1, len(rows)))
    rect(s, [x, y, w, head_h], fill="slab")
    for i, c in enumerate(cols):
        text(s, [x + i * cw + 24, y + 16, cw - 32, 32], str(c), "body-md",
             color="on-slab", align="left" if i == 0 else "center")
    for r, row in enumerate(rows):
        ry = y + head_h + r * rh
        if r % 2 == 1:
            rect(s, [x, ry, w, rh], fill="surface-2")
        cells = row if isinstance(row, list) else row.get("cells", [])
        for i, cell in enumerate(cells):
            text(s, [x + i * cw + 24, ry + (rh - 24) / 2, cw - 32, 28], str(cell),
                 "body", color="ink" if i == 0 else "body",
                 align="left" if i == 0 else "center")
    return s


def build_rank(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["rank"]["slots"]
    x, y, w, h = L["rows"]
    rows = d.get("rows", [])
    n = max(1, len(rows))
    rh = min(86, h / n)
    for r, row in enumerate(rows):
        ry = y + r * rh
        hl = row.get("highlight")
        if hl:
            card(s, [x, ry, w, rh - 8], "card", fill="surface-3")
        rk = oval(s, [x + 12, ry + (rh - 8) / 2 - 24, 48, 48],
                  fill="blue" if hl else "divider")
        shape_text(rk, str(row.get("rank", r + 1)), "h4",
                   color="on-accent" if hl else "muted")
        text(s, [x + 80, ry, w - 300, rh - 8], row.get("label", ""), "body-md",
             color="ink", anchor="middle")
        delta = str(row.get("delta", ""))
        dc = "green" if delta.startswith("▲") else (
            "red" if delta.startswith("▼") else "subtle")
        text(s, [x + w - 200, ry, 180, rh - 8], delta, "body-md", color=dc,
             align="right", anchor="middle")
    return s


# ===========================================================================
# Family 5 — Field
# ===========================================================================
def build_quadrant(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["quadrant"]["slots"]
    cells = d.get("cells", [])
    boxes = _split_grid(L["cells"], 2, 2)
    # 4셀 채움은 균일한 라이트 서피스(slate-50)로 통일 — 셀별 강조는
    # kicker/heading 색(accents)에서만 표현한다(랜덤하게 한 셀만 짙어지지 않게).
    accents = accent_cycle(4)
    for i in range(min(4, len(cells))):
        bx, by, bw, bh = boxes[i]
        cell = cells[i]
        card(s, boxes[i], "card", fill="surface-2")
        text(s, [bx + PAD, by + 24, bw - 2 * PAD, 30], cell.get("label", ""),
             "label", color=accents[i], upper=True)
        if cell.get("heading"):
            text(s, [bx + PAD, by + 62, bw - 2 * PAD, 44], cell["heading"], "h3",
                 color="ink")
        if cell.get("bullets"):
            bullets(s, [bx + PAD, by + 124, bw - 2 * PAD, bh - 150],
                    cell["bullets"], "body", marker_color=accents[i])
        elif cell.get("body"):
            text(s, [bx + PAD, by + 124, bw - 2 * PAD, bh - 150], cell["body"],
                 "body", color="body")
    # 축 라벨(가장자리)
    cx, cy, cw, ch = L["cells"]
    if d.get("x_axis"):
        text(s, [cx, cy + ch + 6, cw, 24], d["x_axis"], "caption",
             color="muted", align="center")
    if d.get("y_axis"):
        text(s, [cx - 60, cy, 60, ch], d["y_axis"], "caption", color="muted",
             align="center", anchor="middle")
    return s


def build_map(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["map"]["slots"]
    x, y, w, h = L["plot"]
    card(s, L["plot"], "card", fill="surface-2")
    line(s, x + w / 2, y + 16, x + w / 2, y + h - 16, color="hairline", w=1.5)
    line(s, x + 16, y + h / 2, x + w - 16, y + h / 2, color="hairline", w=1.5)
    if d.get("x_axis"):
        text(s, [x, y + h + 6, w, 24], d["x_axis"], "caption", color="muted",
             align="center")
    if d.get("y_axis"):
        text(s, [x - 60, y, 60, h], d["y_axis"], "caption", color="muted",
             align="center", anchor="middle")
    for p in d.get("points", []):
        px = x + 20 + p.get("x", 0.5) * (w - 40)
        py = y + h - 20 - p.get("y", 0.5) * (h - 40)
        hl = p.get("highlight")
        dot = oval(s, [px - 14, py - 14, 28, 28],
                   fill="orange" if hl else "blue")
        text(s, [px + 20, py - 14, 280, 28], p.get("label", ""), "body-md",
             color="ink" if hl else "body")
    return s


# ===========================================================================
# Family 6 — Structure
# ===========================================================================
def build_flow(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["flow"]["slots"]
    steps = d.get("steps", [])
    n = max(1, len(steps))
    boxes = _split_cols(L["steps"], n, gap=56)
    for i, st in enumerate(steps):
        bx, by, bw, bh = boxes[i]
        card(s, [bx, by + 40, bw, bh - 40], "card", fill="surface-2")
        circ = oval(s, [bx + bw / 2 - 32, by + 8, 64, 64], fill="blue")
        shape_text(circ, str(i + 1), "h3", color="on-accent")
        text(s, [bx + 20, by + 96, bw - 40, 56], st.get("heading", ""), "h4",
             color="ink", align="center")
        if st.get("body"):
            text(s, [bx + 20, by + 160, bw - 40, bh - 200], st["body"], "body",
                 color="body", align="center")
        if i < n - 1:
            ny = by + 40 + (bh - 40) / 2
            line(s, bx + bw + 10, ny, bx + bw + 46, ny, color="hairline", w=2)
    return s


def build_system(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["system"]["slots"]
    nodes = d.get("nodes", [])
    n = len(nodes)
    # 코어는 nodes 영역(그리드 파생) 중앙에 두고, 노드는 그 둘레에 배치해
    # 코어가 어떤 노드도 가리지 않게 한다(오빗수학 폐기·그리드 셀 파생).
    nx, ny, nw, nh = L["nodes"]
    # 코어 박스: nodes 영역의 가운데 컬럼 블록(전체 폭의 중앙 ~30%).
    core_w = nw * 0.30
    core_h = nh * 0.40
    core_box = [nx + (nw - core_w) / 2, ny + (nh - core_h) / 2, core_w, core_h]
    ccx = core_box[0] + core_box[2] / 2
    ccy = core_box[1] + core_box[3] / 2
    # 노드 둘레 배치: 좌/우 측면 컬럼 블록에 노드를 위→아래로 분배.
    # 좌측에 ceil(n/2), 우측에 나머지. 코어가 차지하는 중앙 컬럼은 비운다.
    side_w = (nw - core_w) / 2 - GAP
    left_n = (n + 1) // 2
    right_n = n - left_n
    boxes = []
    if left_n:
        boxes += _split_rows([nx, ny, side_w, nh], left_n, gap=GAP)
    if right_n:
        rx = nx + nw - side_w
        boxes += _split_rows([rx, ny, side_w, nh], right_n, gap=GAP)
    # 먼저 커넥터(코어→노드)
    for i in range(min(len(boxes), n)):
        bx, by, bw, bh = boxes[i]
        line(s, ccx, ccy, bx + bw / 2, by + bh / 2, color="hairline", w=1.5)
    # 노드 박스
    for i in range(min(len(boxes), n)):
        bx, by, bw, bh = boxes[i]
        nd = nodes[i]
        nb = card(s, [bx, by, bw, bh], "card", fill="surface-3",
                  line_color="blue", line_w=1.5)
        shape_text(nb, nd if isinstance(nd, str) else nd.get("label", ""),
                   "body-md", color="ink")
    # 코어(맨 위에 — 둘레가 비어 있으므로 어떤 노드도 가리지 않음). 강조 잉크 타일.
    cb = rect(s, core_box, fill="blue", radius=RAD["lg"])
    shape_text(cb, d.get("core", "Core"), "h3", color="on-accent")
    return s


# ===========================================================================
# Family 7 — Data
# ===========================================================================
def build_share(prs, d):
    s = new_slide(prs)
    header(s, d)
    L = LAYOUTS["archetypes"]["share"]["slots"]
    spec = dict(d)
    spec.setdefault("type", "doughnut")
    spec["legend"] = False  # share draws its own color-chip legend below (bug B)
    chart(s, L["chart"], spec)
    # 범례/주석 영역: 항목별 색칩 + 라벨 + 값
    items = []
    cats = d.get("categories", [])
    vals = d.get("values") or (d.get("series", [{}])[0].get("values", []) if d.get("series") else [])
    for i, c in enumerate(cats):
        items.append((c, vals[i] if i < len(vals) else ""))
    lx, ly, lw, lh = L["legend"]
    rowh = min(48, lh / max(1, len(items)))
    for i, (c, v) in enumerate(items):
        ry = ly + i * rowh
        rect(s, [lx, ry + 6, 20, 20], fill=chart_color(i),
             radius=RAD["sm"])
        text(s, [lx + 32, ry, lw - 120, rowh], str(c), "body-md", color="ink",
             anchor="middle")
        if v != "":
            text(s, [lx + lw - 90, ry, 90, rowh], str(v), "body-md",
                 color="body", align="right", anchor="middle")
    return s


def build_bars(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "bar_clustered")
    chart(s, LAYOUTS["archetypes"]["bars"]["slots"]["chart"], spec)
    return s


def build_trend(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "line_markers" if len(d.get("series", [])) > 1 else "line")
    chart(s, LAYOUTS["archetypes"]["trend"]["slots"]["chart"], spec)
    return s


def build_spread(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "column_clustered")  # 히스토그램은 컬럼으로 표현
    spec.setdefault("legend", False)
    chart(s, LAYOUTS["archetypes"]["spread"]["slots"]["chart"], spec)
    return s


def build_correlate(prs, d):
    s = new_slide(prs)
    header(s, d)
    spec = dict(d)
    spec.setdefault("type", "bubble" if any("size" in p for p in d.get("points", [])) else "xy_scatter")
    chart(s, LAYOUTS["archetypes"]["correlate"]["slots"]["chart"], spec)
    return s


def build_kpi(prs, d):
    s = new_slide(prs)
    header(s, d)
    # KPI 카드 배열은 items(권장) / kpis / stats 중 하나로 받는다(Set·List 계열과 키 일관).
    spec = {"type": "kpi_cards", "kpis": d.get("kpis") or d.get("stats") or d.get("items", [])}
    chart(s, LAYOUTS["archetypes"]["kpi"]["slots"]["cells"], spec)
    return s


# ---------------------------------------------------------------------------
# Routing table — 새 슬러그 → 빌드 함수.
# ---------------------------------------------------------------------------
ARCHETYPE_TO_BUILDER = {
    # Frame
    "cover": build_cover, "section": build_section, "agenda": build_agenda,
    "closing": build_closing,
    # Focus
    "statement": build_statement, "feature": build_feature,
    # Set
    "duo": build_duo, "trio": build_trio, "grid": build_grid, "list": build_list,
    # Contrast
    "versus": build_versus, "matrix": build_matrix, "rank": build_rank,
    # Field
    "quadrant": build_quadrant, "map": build_map,
    # Structure
    "flow": build_flow, "system": build_system,
    # Data
    "share": build_share, "bars": build_bars, "trend": build_trend,
    "spread": build_spread, "correlate": build_correlate, "kpi": build_kpi,
}


# ---------------------------------------------------------------------------
# Demo spec (--spec 없을 때) — 모든 아키타입 1장씩, 중립 플레이스홀더 콘텐츠.
# ---------------------------------------------------------------------------
def _demo_spec():
    return {
        "title": "ppt-lab Rebuild — 전 아키타입 데모",
        "slides": [
            {"archetype": "cover", "data": {"kicker": "REBUILD", "bg": "navy",
                "title": "클린룸 PPT 엔진\n그리드 SSOT 데모",
                "subtitle": "grid.json → layouts.json → 23 아키타입",
                "meta": "2026 · ppt-lab-rebuild"}},
            {"archetype": "agenda", "data": {"kicker": "AGENDA", "title": "목차",
                "items": ["프레임", "포커스", "세트", "대비", "필드", "구조", "데이터"]}},
            {"archetype": "section", "data": {"index": "01", "title": "데이터 패밀리",
                "caption": "Zelazny 5비교유형 + KPI"}},
            {"archetype": "statement", "data": {"kicker": "주장", "title": "한 슬라이드, 한 메시지",
                "lead": "핵심 주장을 먼저 단언한다", "support": ["근거 1", "근거 2", "근거 3"]}},
            {"archetype": "feature", "data": {"kicker": "FEATURE", "title": "히어로 + 인사이트",
                "bullets": ["좌측 본문", "우측 미디어"], "media": {"caption": "제품 스크린샷"}}},
            {"archetype": "duo", "data": {"kicker": "비교 아님", "title": "두 동등 항목",
                "items": [{"label": "A", "heading": "항목 A", "bullets": ["요점 1", "요점 2"]},
                          {"label": "B", "heading": "항목 B", "image": {}, "body": "이미지 카드 변종"}]}},
            {"archetype": "trio", "data": {"kicker": "세 가지", "title": "세 동등 항목",
                "items": [{"label": "S", "heading": "속도", "body": "빠름"},
                          {"label": "C", "heading": "일관성", "body": "한 톤"},
                          {"label": "Q", "heading": "품질", "body": "정제"}]}},
            {"archetype": "grid", "data": {"kicker": "카드", "title": "4업 그리드",
                "items": [{"heading": f"카드 {i+1}", "body": "설명"} for i in range(4)]}},
            {"archetype": "list", "data": {"kicker": "리스트", "title": "순차 항목",
                "items": [{"title": f"단계 {i+1}", "body": "설명"} for i in range(4)]}},
            {"archetype": "versus", "data": {"kicker": "대비", "title": "Before vs After",
                "left": {"label": "BEFORE", "heading": "수작업", "points": ["느림", "불일치"]},
                "right": {"label": "AFTER", "heading": "자동화", "points": ["빠름", "일관"]}}},
            {"archetype": "matrix", "data": {"kicker": "표", "title": "기준 × 옵션",
                "headers": ["기준", "옵션 A", "옵션 B"],
                "rows": [["속도", "높음", "중간"], ["비용", "낮음", "높음"]]}},
            {"archetype": "rank", "data": {"kicker": "서열", "title": "순위 변동",
                "rows": [{"rank": 1, "label": "1위", "delta": "▲2", "highlight": True},
                         {"rank": 2, "label": "2위", "delta": "▼1"},
                         {"rank": 3, "label": "3위", "delta": "—"}]}},
            {"archetype": "quadrant", "data": {"kicker": "2×2", "title": "사분면",
                "x_axis": "x축", "y_axis": "y축",
                "cells": [{"label": L, "heading": L, "bullets": ["내용"]}
                          for L in ["강점", "약점", "기회", "위협"]]}},
            {"archetype": "map", "data": {"kicker": "포지셔닝", "title": "연속 2축 맵",
                "x_axis": "가격", "y_axis": "품질",
                "points": [{"x": 0.7, "y": 0.8, "label": "우리", "highlight": True},
                           {"x": 0.3, "y": 0.4, "label": "경쟁사"}]}},
            {"archetype": "flow", "data": {"kicker": "과정", "title": "3단계 흐름",
                "steps": [{"heading": f"단계 {i+1}", "body": "설명"} for i in range(3)]}},
            {"archetype": "system", "data": {"kicker": "구조", "title": "아키텍처",
                "core": "코어", "nodes": ["노드 A", "노드 B", "노드 C", "노드 D"]}},
            {"archetype": "share", "data": {"kicker": "구성", "title": "점유율",
                "categories": ["A", "B", "C", "D"], "values": [40, 30, 20, 10]}},
            {"archetype": "bars", "data": {"kicker": "항목", "title": "항목 비교",
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "series": [{"name": "건수", "values": [120, 240, 410, 690]}]}},
            {"archetype": "trend", "data": {"kicker": "시계열", "title": "추이",
                "categories": ["1월", "2월", "3월", "4월", "5월"],
                "series": [{"name": "Organic", "values": [10, 18, 27, 35, 52]},
                           {"name": "Referral", "values": [5, 9, 14, 20, 28]}]}},
            {"archetype": "spread", "data": {"kicker": "빈도", "title": "분포",
                "categories": ["0-10", "10-20", "20-30", "30-40", "40+"],
                "values": [3, 12, 25, 18, 7]}},
            {"archetype": "correlate", "data": {"kicker": "상관", "title": "산점도",
                "points": [{"x": i, "y": i * 1.4 + (i % 3)} for i in range(8)]}},
            {"archetype": "kpi", "data": {"kicker": "지표", "title": "핵심 지표",
                "kpis": [{"label": "제작시간", "value": "-78%", "note": "건당"},
                         {"label": "MAU", "value": "12K", "note": "+24%"},
                         {"label": "재사용", "value": "23", "note": "아키타입"},
                         {"label": "NPS", "value": "62", "note": "베타"}]}},
            {"archetype": "closing", "data": {"kicker": "NEXT", "bg": "navy",
                "title": "지금, 한 번의 요청으로", "subtitle": "자료 첨부 → 빌드",
                "cta": "데모 빌드 →"}},
        ],
    }


DEMO_SPEC = _demo_spec()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build(spec, out_path, palette_override=None, style_override=None, look_override=None):
    warnings = []
    look_id = look_override or spec.get("look")
    lk = apply_look(look_id) if look_id else None      # 룩: fonts+components+palette 통째
    if look_id and lk is None:
        warnings.append(f"unknown look '{look_id}' -> ignored (house/base fallback)")
    # 명시 --style 이 있거나 룩이 없을 때만 스타일 적용(룩의 컴포넌트 보존)
    style_id = style_override or spec.get("style")
    if style_id or lk is None:
        st = apply_style(style_id)
        # 명시 style_id 가 주어졌는데 매칭 실패 → house 로 폴백됨을 경고
        if style_id and st is not None and STYLES.get(str(style_id)) is None:
            valid_styles = ", ".join(k for k in STYLES if k != "_meta")
            warnings.append(
                f"style '{style_id}' not found (valid: {valid_styles}) -> using house")
    # 명시 --palette 가 있거나 룩이 없을 때만 팔레트 적용(룩의 색 보존)
    pal = palette_override or spec.get("palette")
    if pal is not None or lk is None:
        applied = apply_palette(pal)
        # 명시 pal 이 주어졌는데 매칭 실패 → base 로 폴백됨을 경고
        if pal is not None and applied is None:
            valid_pals = ", ".join(str(p["id"]) for p in PALETTES)
            warnings.append(
                f"palette {pal} not found (valid: {valid_pals}) -> using base")
    else:
        applied = {"name": f"look:{look_id} ({lk['palette'].get('blue','')})"}
    # Deck-level dark control: any deck (even on a light look) can declare its own
    # `canvas` and the role palette auto-derives to match that background. This is
    # the generalization — "the deck picks a palette that fits its background".
    # Re-settle DARK/roles after style+palette overrides (ACCENTS may have changed).
    if spec.get("canvas"):
        COLORS["canvas"] = spec["canvas"]
    resolve_dark(force=spec.get("dark"))
    prs = Presentation()
    prs.slide_width = Inches(SLIDE["width_in"])
    prs.slide_height = Inches(SLIDE["height_in"])
    for i, sl in enumerate(spec.get("slides", [])):
        # 슬롯 라우팅 키는 'archetype'(신규 슬러그). 'variant' 도 별칭 허용.
        arch = sl.get("archetype") or sl.get("variant", "statement")
        data = sl.get("data", {})
        fn = ARCHETYPE_TO_BUILDER.get(arch)
        if fn is None:
            warnings.append(f"slide {i+1}: unknown archetype '{arch}' -> fallback statement")
            fn = build_statement
        try:
            fn(prs, data)
        except Exception as e:
            warnings.append(f"slide {i+1} ({arch}): build error {e} -> blank")
            new_slide(prs)
    try:
        prs.save(out_path)
    except PermissionError:
        # 산출 .pptx 가 PowerPoint(또는 잔존 COM 렌더)에서 열려 있으면 쓰기가 거부된다.
        # 조용히 실패해 옛 파일이 그대로 남는 함정을 막기 위해 명확히 알린다.
        sys.stderr.write(
            f"[ppt-lab] ERROR: '{out_path}' 쓰기 거부(PermissionError) — "
            f"PowerPoint 등에서 파일이 열려 있으면 닫고 다시 빌드하세요.\n")
        raise SystemExit(2)
    return applied, warnings


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("out", help="output .pptx path")
    ap.add_argument("--palette", type=int, default=None,
                    help="palette id (valid: 1,2,3,4,5,6,8 — there is no 7)")
    ap.add_argument("--style", default=None, help="style preset id (default: house)")
    ap.add_argument("--look", default=None, help="design-pick look slug (self-contained style+palette; layout stays orthogonal)")
    ap.add_argument("--spec", default=None, help="spec.json path (omit for demo)")
    args = ap.parse_args()

    global SPEC_DIR
    if args.spec:
        spec = json.load(open(args.spec, encoding="utf-8"))
        SPEC_DIR = os.path.dirname(os.path.abspath(args.spec))  # 미디어 상대경로 기준
    else:
        spec = DEMO_SPEC

    applied, warnings = build(spec, args.out, args.palette, args.style, args.look)
    pal_name = applied["name"] if applied else "base (#2563EB)"
    print(f"[ppt-lab] built {len(spec.get('slides', []))} slides -> {args.out}")
    if args.look or spec.get("look"):
        print(f"[ppt-lab] look:    {args.look or spec.get('look')}")
    print(f"[ppt-lab] palette: {pal_name}")
    print(f"[ppt-lab] style:   {args.style or ('(from look)' if (args.look or spec.get('look')) else 'house')}  "
          f"(latin={ACTIVE_FONT['latin']} / ea={ACTIVE_FONT['ea']})")
    if warnings:
        print("[ppt-lab] WARNINGS:", file=sys.stderr)
        for w in warnings:
            print(f"[ppt-lab] WARNING: {w}", file=sys.stderr)


if __name__ == "__main__":
    main()
